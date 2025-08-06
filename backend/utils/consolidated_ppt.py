import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import uuid, os

def _add_title_slide(prs, main_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    frame = title_box.text_frame
    frame.text = "Asia Crystal Commodity LLP"
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)

    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1))
    frame = subtitle.text_frame
    frame.text = main_title
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    footer = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.33), Inches(1))
    frame = footer.text_frame
    frame.text = "ACCLLP"
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)

def _add_table_slide(prs, title_text, df, percent_cols=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.75))
    frame = title.text_frame
    frame.text = title_text
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)

    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5)).table

    for i in range(cols):
        table.columns[i].width = Inches(3)

    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    for row_idx in range(rows):
        is_total = str(df.iloc[row_idx, 0]).upper() in ['TOTAL', 'GRAND TOTAL']
        for col_idx in range(cols):
            val = df.iloc[row_idx, col_idx]
            cell = table.cell(row_idx + 1, col_idx)

            if percent_cols and col_idx in percent_cols:
                cell.text = f"{val}%" if pd.notnull(val) else ""
            else:
                cell.text = str(val)

            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if is_total:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(211, 211, 211)
                cell.text_frame.paragraphs[0].font.bold = True
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

def generate_consolidated_ppt(payload):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title = payload.get("title", "ACCLLP Consolidated Report")
    dfs_with_titles = payload.get("dfs_with_titles", [])

    _add_title_slide(prs, title)

    for item in dfs_with_titles:
        title_text = item.get("title")
        columns = item.get("columns", [])
        percent_cols = item.get("percent_cols", [])
        data = item.get("data", [])

        if not data or not columns:
            continue

        df = pd.DataFrame(data)[columns]

        for col in df.columns:
            if pd.api.types.is_numeric_dtype(df[col]):
                df[col] = df[col].round(0).astype("Int64")

        _add_table_slide(prs, title_text, df, percent_cols)

    out_path = f"static/Consolidated_Report_{uuid.uuid4().hex[:6]}.pptx"
    prs.save(out_path)
    return out_path
