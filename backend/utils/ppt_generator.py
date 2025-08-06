import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
import os
from io import BytesIO
import uuid
import json
import logging

def _add_title_slide(prs, month_title, logo_file=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Company name
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    frame = title_box.text_frame
    frame.text = "Asia Crystal Commodity LLP"
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)

    if logo_file is not None:
        try:
            from io import BytesIO
            logo_stream = BytesIO(logo_file.read())  # Read file-like object into BytesIO
            slide.shapes.add_picture(logo_stream, Inches(5.665), Inches(1.5), width=Inches(2), height=Inches(2))
            logo_file.seek(0)  # Reset pointer for reuse
        except Exception as e:
            print(f"⚠️ Failed to insert logo: {e}")

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(1))
    frame = title.text_frame
    frame.text = f"Monthly Review Meeting – {month_title}"
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.33), Inches(1))
    frame = subtitle.text_frame
    frame.text = "ACCLLP"
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)

def _add_table_slide(prs, title_text, df, percent_col_index=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.75))
    frame = title.text_frame
    frame.text = title_text
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)

    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5)).table
    for i in range(cols):
        table.columns[i].width = Inches(3)
        
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    for row_idx in range(rows):
        is_total = df.iloc[row_idx, 0].upper() in ['TOTAL', 'GRAND TOTAL'] if isinstance(df.iloc[row_idx, 0], str) else False
        for col_idx in range(cols):
            val = df.iloc[row_idx, col_idx]
            text = f"{val}%" if percent_col_index is not None and col_idx == percent_col_index else str(val)
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if is_total:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(211, 211, 211)
                cell.text_frame.paragraphs[0].font.bold = True
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

def create_consolidated_ppt(all_dfs_with_titles, title="ACCLLP Consolidated Report"):
    """Create a consolidated PPT with all report data. Automatically loads static logo."""
    try:
        from io import BytesIO
        import os

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 🔥 Load logo from static folder (internal handling)
        logo_stream = None
        static_logo_path = os.path.join("static", "logo.jpeg")  # adjust if .jpg/.png

        if os.path.exists(static_logo_path):
            try:
                with open(static_logo_path, "rb") as f:
                    logo_stream = BytesIO(f.read())
                    logo_stream.seek(0)
                logger.info(f"✅ Logo loaded from {static_logo_path}")
            except Exception as e:
                logger.warning(f"⚠️ Could not read logo from static: {e}")
        else:
            logger.warning("⚠️ No static logo found — proceeding without logo.")

        # 🧠 Title slide
        create_title_slide(prs, title, logo_stream)

        # 📊 Table slides
        for df_info in all_dfs_with_titles:
            if df_info and 'df' in df_info and 'title' in df_info:
                add_table_slide(
                    prs,
                    df_info['df'],
                    df_info['title'],
                    df_info.get('percent_cols')
                )

        # 📤 Output PPT
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer

    except Exception as e:
        logger.error(f"❌ Error creating consolidated PPT: {e}", exc_info=True)
        return None

def generate_budget_ppt(payload):
    try:
        ppt = Presentation()
        ppt.slide_width = Inches(13.33)
        ppt.slide_height = Inches(7.5)

        month = payload.get("month", "Month")

        logo_stream = None
        static_logo_path = os.path.join("static", "logo.jpeg")  # adjust if .jpg/.png

        if os.path.exists(static_logo_path):
            try:
                with open(static_logo_path, "rb") as f:
                    logo_stream = BytesIO(f.read())
                    logo_stream.seek(0)
                logger.info(f"✅ Logo loaded from {static_logo_path}")
            except Exception as e:
                logger.warning(f"⚠️ Could not read logo from static: {e}")
        else:
            logger.warning("⚠️ No static logo found — proceeding without logo.")

        _add_title_slide(ppt, month, logo_stream)

        def convert(data):
            df = pd.DataFrame(data) if data else pd.DataFrame() 
            for col in df.columns:
                if pd.api.types.is_numeric_dtype(df[col]):
                    df[col] = df[col].round(0).astype('Int64')
            return df
        
        def safe_percent_index(df):
            try:
                return df.columns.get_loc('%')
            except Exception:
                return None
            
        def convert_ordered(section):
            if isinstance(section, dict) and 'data' in section and 'columns' in section:
                df = pd.DataFrame(section['data'])[section['columns']]  # Use only specified order
            else:
                df = pd.DataFrame(section or [])

            for col in df.columns:
                if pd.api.types.is_numeric_dtype(df[col]):
                    df[col] = df[col].round(0).astype('Int64')
            return df
            
        qty_df = convert_ordered(payload.get("budget_vs_billed_qty"))
        val_df = convert_ordered(payload.get("budget_vs_billed_value"))
        overall_qty_df = convert_ordered(payload.get("overall_sales_qty"))
        overall_val_df = convert_ordered(payload.get("overall_sales_value"))


        if not qty_df.empty:
            _add_table_slide(ppt, f"BUDGET AGAINST BILLED (Qty in Mt) – {month}", qty_df, percent_col_index=safe_percent_index(qty_df))
        if not val_df.empty:
            _add_table_slide(ppt, f"BUDGET AGAINST BILLED (Value in Lakhs) – {month}", val_df, percent_col_index=safe_percent_index(val_df))
        if not overall_qty_df.empty:
            _add_table_slide(ppt, f"OVERALL SALES (Qty in Mt) – {month}", overall_qty_df)
        if not overall_val_df.empty:
            _add_table_slide(ppt, f"OVERALL SALES (Value in Lakhs) – {month}", overall_val_df)

        out_path = f"static/Budget_vs_Billed_{month.replace(' ', '_')}_{uuid.uuid4().hex[:6]}.pptx"
        ppt.save(out_path)
        return out_path
    except Exception as e:
        print("❌ Error in generate_budget_ppt:", str(e))
        raise Exception(f"Error generating PPT: {e}")


logger = logging.getLogger(__name__)

def create_title_slide(prs, title, logo_file=None):
    """Create a standard title slide for PPT."""
    try:
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        title_slide = prs.slides.add_slide(blank_slide_layout)

        # Company Name
        company_name = title_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
        company_frame = company_name.text_frame
        company_frame.text = "Asia Crystal Commodity LLP"
        p = company_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)

        # Logo (optional)
        if logo_file is not None:
            try:
                logo_file.seek(0)  # Ensure we're at the start
                title_slide.shapes.add_picture(logo_file, Inches(5.665), Inches(1.5), width=Inches(2), height=Inches(2))
                logger.info("✅ Logo added to title slide.")
            except Exception as e:
                logger.error(f"⚠️ Failed to add logo: {e}")

        # Title
        title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)

        # Subtitle
        subtitle = title_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.33), Inches(1))
        subtitle_frame = subtitle.text_frame
        subtitle_frame.text = "ACCLLP"
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)

        return title_slide

    except Exception as e:
        logger.error(f"❌ Error generating title slide: {e}")
        return None


def add_table_slide(prs, df, title, percent_cols=None):
    """Add a slide with a table to a PPT presentation."""
    if df is None or df.empty:
        return None
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.75))
    title_frame = title_shape.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5)).table
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    for row_idx in range(df.shape[0]):
        is_total_row = row_idx == rows - 1 and str(df.iloc[row_idx, 0]).upper() in ['TOTAL', 'GRAND TOTAL']
        for col_idx in range(cols):
            cell = table.cell(row_idx + 1, col_idx)
            value = df.iloc[row_idx, col_idx]
            if percent_cols and col_idx in percent_cols:
                cell.text = f"{value}%"
            else:
                cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            if row_idx % 2 == 0 and not is_total_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            if is_total_row:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(211, 211, 211)    
    return slide


def create_od_ppt_updated(df, regional_df, title, logo_file=None):
    """Create a PPT for OD Target vs Collection Report with regional summary (Flask-compatible)."""
    try:
        from io import BytesIO
        import os

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 🔁 Auto-load static logo if not passed
        if logo_file is None:
            static_logo_path = os.path.join("static", "logo.jpeg")
            if os.path.exists(static_logo_path):
                try:
                    with open(static_logo_path, "rb") as f:
                        logo_file = BytesIO(f.read())
                        logo_file.seek(0)
                    logger.info(f"✅ Static logo loaded from {static_logo_path}")
                except Exception as e:
                    logger.warning(f"⚠️ Could not read static logo: {e}")
                    logo_file = None
            else:
                logger.warning("⚠️ Static logo not found — proceeding without logo.")

        # Add Title Slide
        create_title_slide(prs, title, logo_file)

        # Branch-wise Table Slide
        add_table_slide(prs, df, f"Branch-wise Performance - {title}", percent_cols=[3, 6])

        # Regional Summary Table Slide
        if regional_df is not None and not regional_df.empty:
            add_table_slide(prs, regional_df, f"Regional Summary - {title}", percent_cols=[3, 6])

        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)

        return ppt_buffer

    except Exception as e:
        logger.error(f"❌ Error creating OD PPT: {e}", exc_info=True)
        return None



def create_product_growth_ppt(group_results, month_title, logo_file=None):
    try:
        from io import BytesIO
        import os

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 🔁 Load static logo if not provided
        if logo_file is None:
            static_logo_path = os.path.join("static", "logo.jpeg")
            if os.path.exists(static_logo_path):
                try:
                    with open(static_logo_path, "rb") as f:
                        logo_file = BytesIO(f.read())
                        logo_file.seek(0)
                    logger.info(f"✅ Static logo loaded from {static_logo_path}")
                except Exception as e:
                    logger.warning(f"⚠️ Failed to read static logo: {e}")
                    logo_file = None
            else:
                logger.warning("⚠️ Static logo not found, proceeding without logo.")

        # Add title slide
        create_title_slide(prs, f"Product Growth by Company Group – {month_title}", logo_file)

        # Column order templates
        qty_order = ["PRODUCT NAME", "LY_QTY", "BUDGET_QTY", "CY_QTY", "ACHIEVEMENT %"]
        value_order = ["PRODUCT NAME", "LY_VALUE", "BUDGET_VALUE", "CY_VALUE", "ACHIEVEMENT %"]

        # Helper: reorder columns
        def reorder_df(df, order):
            df = df.copy()
            return df[[col for col in order if col in df.columns]]

        # Loop through each company group
        for group, data in group_results.items():
            qty_df = reorder_df(data['qty_df'], qty_order)
            value_df = reorder_df(data['value_df'], value_order)

            # Round numeric data
            for col in qty_df.select_dtypes(include='number').columns:
                qty_df[col] = qty_df[col].round(0)
            for col in value_df.select_dtypes(include='number').columns:
                value_df[col] = value_df[col].round(0)

            # Add slides
            add_table_slide(prs, qty_df, f"{group} - Quantity Growth (Qty in Mt)", percent_cols=[4])
            add_table_slide(prs, value_df, f"{group} - Value Growth (Value in Lakhs)", percent_cols=[4])

        # Export PPT
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer

    except Exception as e:
        logger.error(f"❌ Error creating Product Growth PPT: {e}", exc_info=True)
        return None
def create_customer_ppt_slide(slide, df, title, sorted_months, financial_year):
    """Add a slide with the billed customers table to the presentation."""
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = f"{title} - FY {financial_year}"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER
    data_columns = ['S.No', 'Mapped_Branch'] + sorted_months
    display_columns = ['S.No', 'Branch Name'] + sorted_months
    rows = len(df) + 1
    ncols = len(display_columns)
    table_width = Inches(12.0)
    table_height = Inches(0.3 * len(df) + 0.4)
    left = Inches(0.65)
    top = Inches(1.2)
    table = slide.shapes.add_table(rows, ncols, left, top, table_width, table_height).table
    col_widths = [Inches(0.5), Inches(2.0)] + [Inches(0.75)] * len(sorted_months)
    for col_idx in range(ncols):
        table.columns[col_idx].width = col_widths[col_idx]
    for col_idx, col_name in enumerate(display_columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    total_row_idx = len(df) - 1
    for row_idx in range(len(df)):
        is_total_row = row_idx == total_row_idx
        for col_idx, col_key in enumerate(data_columns):
            cell = table.cell(row_idx + 1, col_idx)
            value = df.iloc[row_idx][col_key] if col_key in df.columns else ""
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            if is_total_row:
                cell.fill.fore_color.rgb = RGBColor(211, 211, 211)
                cell.text_frame.paragraphs[0].font.bold = True
            else:
                cell.fill.fore_color.rgb = RGBColor(221, 235, 247) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
                
def create_nbc_individual_ppt(customer_df, customer_title, sorted_months, financial_year, logo_file=None):
    """Create individual PPT for NBC (Number of Billed Customers) report."""
    try:
        from io import BytesIO
        import os

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 🔁 Auto-load static logo if not provided
        if logo_file is None:
            static_logo_path = os.path.join("static", "logo.jpeg")
            if os.path.exists(static_logo_path):
                try:
                    with open(static_logo_path, "rb") as f:
                        logo_file = BytesIO(f.read())
                        logo_file.seek(0)
                    logger.info(f"✅ Logo loaded from {static_logo_path}")
                except Exception as e:
                    logger.warning(f"⚠️ Failed to load static logo: {e}")
                    logo_file = None
            else:
                logger.warning("⚠️ Static logo not found — proceeding without logo.")

        # Add Title Slide
        create_title_slide(prs, "Number of Billed Customers Report", logo_file)

        # Add Customer Analysis Slide
        slide_layout = prs.slide_layouts[6]  # Blank slide
        slide = prs.slides.add_slide(slide_layout)
        create_customer_ppt_slide(slide, customer_df, customer_title, sorted_months, financial_year)

        # Export as PPT
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer

    except Exception as e:
        logger.error(f"❌ Error creating NBC PPT: {e}", exc_info=True)
        return None

    
def create_od_ppt_slide(slide, df, title):
    """Add a slide with the OD Target table to the presentation."""
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Match internal column names
    columns_to_show = ['Area', 'TARGET']
    display_names = ['Area', 'TARGET (Lakhs)']

    rows, cols = df.shape[0] + 1, len(columns_to_show)
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(5)).table

    # Header
    for col_idx, display_name in enumerate(display_names):
        cell = table.cell(0, col_idx)
        cell.text = display_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx in range(len(df)):
        is_total_row = row_idx == len(df) - 1
        for col_idx, col_key in enumerate(columns_to_show):
            cell = table.cell(row_idx + 1, col_idx)
            val = df.iloc[row_idx][col_key]
            cell.text = str(val) if col_key == 'Area' else f"{float(val):.2f}"
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(255, 255, 204) if is_total_row else
                (RGBColor(221, 235, 247) if row_idx % 2 == 0 else RGBColor(255, 255, 255))
            )
            if is_total_row:
                cell.text_frame.paragraphs[0].font.bold = True
    
def create_od_individual_ppt(od_target_df, od_title, logo_file=None):
    """Create individual PPT for OD Target report only."""
    try:
        from io import BytesIO
        import os

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 🔁 Auto-load static logo if not provided
        if logo_file is None:
            static_logo_path = os.path.join("static", "logo.jpeg")
            if os.path.exists(static_logo_path):
                try:
                    with open(static_logo_path, "rb") as f:
                        logo_file = BytesIO(f.read())
                        logo_file.seek(0)
                    logger.info(f"✅ Logo loaded from {static_logo_path}")
                except Exception as e:
                    logger.warning(f"⚠️ Failed to load static logo: {e}")
                    logo_file = None
            else:
                logger.warning("⚠️ Static logo not found — proceeding without logo.")

        # Create title slide
        create_title_slide(prs, "OD Target Report", logo_file)

        # Add OD Target table slide
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        create_od_ppt_slide(slide, od_target_df, od_title)

        # Save PPT to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer

    except Exception as e:
        logger.error(f"❌ Error creating OD Target PPT: {e}", exc_info=True)
        return None


