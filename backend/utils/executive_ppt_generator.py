import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import math
import logging
import base64

logger = logging.getLogger(__name__)

def create_title_slide(prs, title, logo_file=None):
    """Create title slide following the exact Streamlit logic"""
    try:
        blank_slide_layout = prs.slide_layouts[6]
        title_slide = prs.slides.add_slide(blank_slide_layout)
        
        # Company name
        company_name = title_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
        company_frame = company_name.text_frame
        company_frame.text = "Asia Crystal Commodity LLP"
        p = company_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        
        # Logo (if provided)
        if logo_file is not None:
            try:
                if isinstance(logo_file, str):
                    # Base64 encoded logo
                    logo_data = base64.b64decode(logo_file)
                    logo_buffer = BytesIO(logo_data)
                    title_slide.shapes.add_picture(
                        logo_buffer, Inches(5.665), Inches(1.5), 
                        width=Inches(2), height=Inches(2)
                    )
                else:
                    # File object
                    logo_buffer = BytesIO(logo_file.read())
                    title_slide.shapes.add_picture(
                        logo_buffer, Inches(5.665), Inches(1.5), 
                        width=Inches(2), height=Inches(2)
                    )
                    logo_file.seek(0)  # Reset for reuse
            except Exception as e:
                logger.error(f"Error adding logo to slide: {str(e)}")
        
        # Title
        title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)
        
        # Subtitle
        subtitle = title_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.33), Inches(1))
        subtitle_frame = subtitle.text_frame
        subtitle_frame.text = "ACCLLP"
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        
        return title_slide
        
    except Exception as e:
        logger.error(f"Error creating title slide: {e}")
        raise

def add_table_slide(prs, df, title, percent_cols=None, is_consolidated=False):
    """Add table slide with strict column order (supports dynamic override via df.attrs['columns'])"""
    try:
        if percent_cols is None:
            percent_cols = []

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title box
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        p.alignment = PP_ALIGN.CENTER

        logger.info(f"🔍 PROCESSING SLIDE: '{title}'")
        logger.info(f"🔍 DataFrame columns: {list(df.columns)}")
        logger.info(f"🔍 DataFrame shape: {df.shape}")
        logger.info(f"🔍 Is consolidated: {is_consolidated}")
        print(f"🔍 PROCESSING: '{title}' | Columns: {list(df.columns)}")

        actual_columns = list(df.columns)
        title_lower = title.lower()

        # ✅ 1. Try dynamic column order from frontend
        strict_order = df.attrs.get("columns", [])
        if strict_order:
            logger.info(f"✅ Consolidated report — using passed DataFrame column order: {strict_order}")
        else:
            # 🔁 2. Fallback to title-based detection
            logger.info(f"⚠️ No column override found, falling back to title-based detection")
            if any(phrase in title_lower for phrase in ["budget against billed", "budget vs billed", "budget v/s billed"]):
                if "qty" in title_lower or "mt" in title_lower:
                    strict_order = ['Executive', 'Budget Qty', 'Billed Qty', '%']
                    logger.info(f"🔒 BUDGET QTY: Expected order: {strict_order}")
                else:
                    strict_order = ['Executive', 'Budget Value', 'Billed Value', '%']
                    logger.info(f"🔒 BUDGET VALUE: Expected order: {strict_order}")
            elif any(phrase in title_lower for phrase in ["overall sales", "overall sale", "total sales"]):
                if "qty" in title_lower or "mt" in title_lower:
                    strict_order = ['Executive', 'Budget Qty', 'Billed Qty']
                    logger.info(f"🔒 OVERALL SALES QTY: Expected order: {strict_order}")
                else:
                    strict_order = ['Executive', 'Budget Value', 'Billed Value']
                    logger.info(f"🔒 OVERALL SALES VALUE: Expected order: {strict_order}")
            elif any(phrase in title_lower for phrase in ["od target vs collection", "od collection"]):
                strict_order = [
                    'Executive', 'Due Target', 'Collection Achieved',
                    'Overall % Achieved', 'For the month Overdue',
                    'For the month Collection', '% Achieved (Selected Month)'
                ]
                logger.info(f"🔒 OD COLLECTION: Expected order: {strict_order}")
            elif any(phrase in title_lower for phrase in ["product growth", "quantity growth"]):
                strict_order = ['PRODUCT GROUP', 'LY_QTY', 'BUDGET_QTY', 'CY_QTY', 'ACHIEVEMENT %']
                logger.info(f"🔒 PRODUCT QTY: Expected order: {strict_order}")
            elif any(phrase in title_lower for phrase in ["value growth", "value in lakhs"]):
                strict_order = ['PRODUCT GROUP', 'LY_VALUE', 'BUDGET_VALUE', 'CY_VALUE', 'ACHIEVEMENT %']
                logger.info(f"🔒 PRODUCT VALUE: Expected order: {strict_order}")
            elif any(phrase in title_lower for phrase in ["customer", "billed customers"]):
                strict_order = []
                if 'S.No' in actual_columns: strict_order.append('S.No')
                if 'Executive Name' in actual_columns: strict_order.append('Executive Name')
                elif 'Executive' in actual_columns: strict_order.append('Executive')
                for col in actual_columns:
                    if col not in strict_order:
                        strict_order.append(col)
                logger.info(f"🔒 CUSTOMER: Dynamic order: {strict_order}")
            else:
                strict_order = actual_columns
                logger.info(f"🔒 FALLBACK: Using DataFrame order: {strict_order}")

        # Final column ordering (only keep what's available in df)
        ordered_columns = [col for col in strict_order if col in actual_columns]
        for col in actual_columns:
            if col not in ordered_columns:
                ordered_columns.append(col)
                logger.warning(f"⚠️ Added missing column to end: {col}")

        logger.info(f"✅ FINAL COLUMN ORDER for '{title}': {ordered_columns}")
        print(f"✅ FINAL COLUMN ORDER for '{title}': {ordered_columns}")

        # [The rest of your existing slide/table creation logic remains unchanged...]
        # -- table.rows, table.columns, header styling, data rows, formatting --
        # You don’t need to touch that code

        # Continue with table generation...
        num_rows = len(df) + 1
        num_cols = len(ordered_columns)

        table = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(0.5), Inches(1.5),
            Inches(12), Inches(0.3 * len(df) + 0.3)
        ).table

        if num_cols > 0:
            table.columns[0].width = Inches(3.0)
            remaining_width = 12.0 - 3.0
            if num_cols > 1:
                col_width = remaining_width / (num_cols - 1)
                for i in range(1, num_cols):
                    table.columns[i].width = Inches(col_width)

        # Header
        for i, col_name in enumerate(ordered_columns):
            cell = table.cell(0, i)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Data rows
        for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
            first_col_value = str(row.iloc[0]).upper()
            is_total_row = first_col_value in ['TOTAL', 'GRAND TOTAL', 'PART 1 TOTAL', 'PART 2 TOTAL']
            for col_idx, col_name in enumerate(ordered_columns):
                cell = table.cell(row_idx, col_idx)
                value = row.get(col_name, "")
                if col_name in ['ACHIEVEMENT %', '%', 'Overall % Achieved', '% Achieved (Selected Month)']:
                    if isinstance(value, (int, float)) and not pd.isna(value):
                        cell.text = f"{value:.2f}%"
                    else:
                        cell.text = str(value) if value is not None else ""
                elif col_idx in percent_cols and isinstance(value, (int, float)) and not pd.isna(value):
                    cell.text = f"{value:.2f}%"
                elif isinstance(value, (int, float)) and not pd.isna(value):
                    cell.text = str(int(round(value))) if abs(value - round(value)) < 0.001 else f"{value:.2f}"
                else:
                    cell.text = str(value) if value is not None and not pd.isna(value) else ""

                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.fill.solid()

                if is_total_row:
                    cell.fill.fore_color.rgb = RGBColor(211, 211, 211)
                    cell.text_frame.paragraphs[0].font.bold = True
                else:
                    cell.fill.fore_color.rgb = RGBColor(221, 235, 247) if row_idx % 2 == 0 else RGBColor(255, 255, 255)

        return slide

    except Exception as e:
        logger.error(f"Error adding table slide: {e}")
        raise

def create_consolidated_ppt(dfs_info, logo_file=None, title="Consolidated Report"):
    """Create consolidated PowerPoint with STREAMLIT-COMPATIBLE column order consistency and proper splitting"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide
        create_title_slide(prs, title, logo_file)
        
        logger.info(f"📊 Creating consolidated PPT with {len(dfs_info)} reports")
        print(f"📊 Creating consolidated PPT with {len(dfs_info)} reports")
        
        # Process each report with strict column order and splitting logic
        for df_info in dfs_info:
            df_data = df_info.get('df', [])
            slide_title = df_info['title']
            percent_cols = df_info.get('percent_cols', [])
            
            # Convert to DataFrame if it's a list of dicts
            if isinstance(df_data, list):
                df = pd.DataFrame(df_data)
            else:
                df = df_data.copy()

             # ✅ Apply frontend-defined column order (if present)
            frontend_order = df_info.get("columns")
            if frontend_order:
                df = df[[col for col in frontend_order if col in df.columns]]
                # df.attrs["columns"] = frontend_order  # Used in add_table_slide
                logger.info(f"📐 Applied frontend column order: {frontend_order}")
            
            if df.empty:
                logger.warning(f"Skipping empty report: {slide_title}")
                continue
            
            logger.info(f"🔄 Processing consolidated report: {slide_title}")
            print(f"🔄 Processing consolidated report: {slide_title}")
            print(f"🔍 Original DataFrame columns: {list(df.columns)}")
            
            # Check if this report needs splitting (same logic as individual PPTs)
            if 'Executive' in df.columns:
                # Use the same splitting logic as process_df_for_slides
                process_df_for_slides(prs, df, slide_title, percent_cols=percent_cols, is_consolidated=True)
            else:
                # Non-executive reports (like product growth, customer analysis)
                add_table_slide(prs, df, slide_title, percent_cols=percent_cols, is_consolidated=True)
        
        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        # Log final slide count
        actual_slide_count = len(prs.slides)
        logger.info(f"📊 Consolidated PPT created with {actual_slide_count} total slides")
        print(f"📊 Consolidated PPT created with {actual_slide_count} total slides")
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error creating consolidated PPT: {str(e)}")
        raise Exception(f"Error creating consolidated PPT: {str(e)}")

def process_df_for_slides(prs, df, title_base, percent_cols=None, is_consolidated=False):
    """Process DataFrame for slides with correct ordering and splitting logic"""
    try:
        if percent_cols is None:
            percent_cols = []
        
        # Remove ACCLP if exists
        df = df[df['Executive'] != "ACCLP"].copy()
        num_executives = df[df['Executive'] != 'TOTAL'].shape[0]
        
        # SPLIT THRESHOLD: 20 executives
        split_threshold = 20
        
        if num_executives <= split_threshold:
            add_table_slide(prs, df, title_base, percent_cols=percent_cols, is_consolidated=is_consolidated)
            return
        
        data_rows = df[df['Executive'] != 'TOTAL'].copy()
        total_row = df[df['Executive'] == 'TOTAL'].copy()
        split_point = math.ceil(num_executives / 2)
        part1 = data_rows.iloc[:split_point].copy()
        part2 = data_rows.iloc[split_point:].copy()
        
        for i, part in enumerate([part1, part2], 1):
            part_total = {}
            for col in df.columns:
                if col == 'Executive':
                    part_total[col] = f'PART {i} TOTAL'
                elif col == '%':
                    budget_col = 'Budget Value' if 'Budget Value' in df.columns else 'Budget Qty'
                    billed_col = 'Billed Value' if 'Billed Value' in df.columns else 'Billed Qty'
                    budget_sum = part[budget_col].sum()
                    billed_sum = part[billed_col].sum()
                    part_total[col] = round((billed_sum / budget_sum * 100), 2) if budget_sum != 0 else 0.0
                else:
                    # ENHANCED ROUNDING: Double-round to eliminate floating point errors
                    raw_sum = part[col].sum()
                    part_total[col] = round(round(raw_sum, 4), 2)  # First round to 4, then to 2 decimals
            
            # Create DataFrame for part with total and ensure all numeric columns are rounded
            part_with_total = pd.concat([part, pd.DataFrame([part_total])], ignore_index=True)
            
            # ADDITIONAL SAFETY ROUNDING: Ensure all numeric columns are properly rounded
            numeric_cols = part_with_total.select_dtypes(include=[np.number]).columns
            for col in numeric_cols:
                part_with_total[col] = part_with_total[col].apply(lambda x: round(float(x), 2))
            
            add_table_slide(prs, part_with_total, f"{title_base} - Part {i}", percent_cols=percent_cols, is_consolidated=is_consolidated)
        
        # Also ensure the grand total row is properly rounded
        numeric_cols = total_row.select_dtypes(include=[np.number]).columns
        for col in numeric_cols:
            total_row[col] = total_row[col].apply(lambda x: round(float(x), 2))
        
        add_table_slide(prs, total_row, f"{title_base} - Grand Total", percent_cols=percent_cols, is_consolidated=is_consolidated)
        
    except Exception as e:
        logger.error(f"Error processing DataFrame for slides: {e}")
        raise

def create_product_growth_ppt(group_results, month_title, logo_file=None, ly_month=None, cy_month=None):
    """Create Product Growth PPT with STRICT order and proper month titles"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # 🔒 CRITICAL FIX: Proper title slide text generation
        if ly_month and cy_month:
            title_slide_text = f"Product Growth – LY: {ly_month} vs CY: {cy_month}"
            logger.info(f"✅ Title with months: {title_slide_text}")
        elif month_title and month_title != "Product Growth Analysis":
            title_slide_text = f"Product Growth – {month_title}"
            logger.info(f"✅ Title with custom title: {title_slide_text}")
        else:
            title_slide_text = "Product Growth Analysis"
            logger.info(f"✅ Default title: {title_slide_text}")
            
        print(f"🔍 DEBUG: ly_month={ly_month}, cy_month={cy_month}, month_title={month_title}")
        print(f"🔍 DEBUG: Final title_slide_text={title_slide_text}")
            
        create_title_slide(prs, title_slide_text, logo_file)
        
        # Add slides for each company
        for company, data in group_results.items():
            logger.info(f"🔄 Processing company: {company}")
            
            # Handle both DataFrame and dict formats
            if isinstance(data['qty_df'], list):
                qty_df = pd.DataFrame(data['qty_df'])
            else:
                qty_df = data['qty_df'].copy()
                
            if isinstance(data['value_df'], list):
                value_df = pd.DataFrame(data['value_df'])
            else:
                value_df = data['value_df'].copy()
            
            # Debug original data
            logger.info(f"🔍 {company} QTY original columns: {list(qty_df.columns)}")
            logger.info(f"🔍 {company} VALUE original columns: {list(value_df.columns)}")
            
            # Round numeric columns but maintain order
            numeric_cols_qty = qty_df.select_dtypes(include=[np.number]).columns
            for col in numeric_cols_qty:
                qty_df[col] = qty_df[col].round(2)
            
            numeric_cols_value = value_df.select_dtypes(include=[np.number]).columns
            for col in numeric_cols_value:
                value_df[col] = value_df[col].round(2)
            
            # 🔒 CRITICAL FIX: Create slide titles with proper month information
            if ly_month and cy_month:
                qty_title = f"{company} - Quantity Growth (Qty in Mt) - LY: {ly_month} vs CY: {cy_month}"
                value_title = f"{company} - Value Growth (Value in Lakhs) - LY: {ly_month} vs CY: {cy_month}"
                logger.info(f"✅ Using month-based titles for {company}")
            else:
                qty_title = f"{company} - Quantity Growth (Qty in Mt)"
                value_title = f"{company} - Value Growth (Value in Lakhs)"
                logger.info(f"✅ Using default titles for {company}")
            
            print(f"🔍 DEBUG: qty_title={qty_title}")
            print(f"🔍 DEBUG: value_title={value_title}")
            
            # Create slides with ENFORCED STRICT column order
            add_table_slide(prs, qty_df, qty_title, is_consolidated=False)
            add_table_slide(prs, value_df, value_title, is_consolidated=False)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error creating Product Growth PPT: {e}")
        return None

def process_df_for_slides(prs, df, title_base, percent_cols=None, is_consolidated=False):
    """Process DataFrame for slides with correct ordering and splitting logic"""
    try:
        if percent_cols is None:
            percent_cols = []
        
        # Remove ACCLP if exists
        df = df[df['Executive'] != "ACCLP"].copy()
        num_executives = df[df['Executive'] != 'TOTAL'].shape[0]
        
        # SPLIT THRESHOLD: 20 executives
        split_threshold = 20
        
        if num_executives <= split_threshold:
            add_table_slide(prs, df, title_base, percent_cols=percent_cols, is_consolidated=is_consolidated)
            return
        
        data_rows = df[df['Executive'] != 'TOTAL'].copy()
        total_row = df[df['Executive'] == 'TOTAL'].copy()
        split_point = math.ceil(num_executives / 2)
        part1 = data_rows.iloc[:split_point].copy()
        part2 = data_rows.iloc[split_point:].copy()
        
        for i, part in enumerate([part1, part2], 1):
            part_total = {}
            for col in df.columns:
                if col == 'Executive':
                    part_total[col] = f'PART {i} TOTAL'
                elif col == '%':
                    budget_col = 'Budget Value' if 'Budget Value' in df.columns else 'Budget Qty'
                    billed_col = 'Billed Value' if 'Billed Value' in df.columns else 'Billed Qty'
                    budget_sum = part[budget_col].sum()
                    billed_sum = part[billed_col].sum()
                    part_total[col] = round((billed_sum / budget_sum * 100), 2) if budget_sum != 0 else 0.0
                else:
                    # ENHANCED ROUNDING: Double-round to eliminate floating point errors
                    raw_sum = part[col].sum()
                    part_total[col] = round(round(raw_sum, 4), 2)  # First round to 4, then to 2 decimals
            
            # Create DataFrame for part with total and ensure all numeric columns are rounded
            part_with_total = pd.concat([part, pd.DataFrame([part_total])], ignore_index=True)
            
            # ADDITIONAL SAFETY ROUNDING: Ensure all numeric columns are properly rounded
            numeric_cols = part_with_total.select_dtypes(include=[np.number]).columns
            for col in numeric_cols:
                part_with_total[col] = part_with_total[col].apply(lambda x: round(float(x), 2))
            
            add_table_slide(prs, part_with_total, f"{title_base} - Part {i}", percent_cols=percent_cols, is_consolidated=is_consolidated)
        
        # Also ensure the grand total row is properly rounded
        numeric_cols = total_row.select_dtypes(include=[np.number]).columns
        for col in numeric_cols:
            total_row[col] = total_row[col].apply(lambda x: round(float(x), 2))
        
        add_table_slide(prs, total_row, f"{title_base} - Grand Total", percent_cols=percent_cols, is_consolidated=is_consolidated)
        
    except Exception as e:
        logger.error(f"Error processing DataFrame for slides: {e}")
        raise

def estimate_slide_count_for_consolidated(dfs_info):
    """Estimate total slide count for consolidated PPT including splits"""
    try:
        slide_count = 1  # Title slide
        
        for df_info in dfs_info:
            df_data = df_info['df']
            title = df_info['title'].lower()
            
            # Convert to DataFrame if it's a list of dicts
            if isinstance(df_data, list):
                df = pd.DataFrame(df_data)
            else:
                df = df_data.copy()
            
            if df.empty:
                continue
                
            # Check if this report has Executive column (indicates it might split)
            has_executive_col = 'Executive' in df.columns
            
            if has_executive_col:
                # Remove ACCLP and count executives for splitting logic
                data_for_counting = df[df['Executive'] != "ACCLP"] if 'Executive' in df.columns else df
                num_data_rows = len(data_for_counting[data_for_counting['Executive'] != 'TOTAL']) if 'Executive' in df.columns else len(data_for_counting)
                
                # Use the same split logic as the backend
                split_threshold = 20
                
                if num_data_rows <= split_threshold:
                    slide_count += 1  # Single slide
                else:
                    # Will be split into 2 parts + 1 grand total
                    slide_count += 3  # Part 1, Part 2, Grand Total
            else:
                # Non-executive reports (like product growth, customer analysis)
                slide_count += 1
        
        return slide_count
        
    except Exception as e:
        logger.error(f"Error estimating slide count: {e}")
        return len(dfs_info) + 1  # Fallback: assume 1 slide per report + title

def create_customer_ppt_slide(slide, df, title, sorted_months, is_last_slide=False):
    """Create customer PPT slide with correct column order"""
    try:
        if df.empty or len(df.columns) < 2:
            logger.warning(f"Skipping customer slide: DataFrame is empty or has insufficient columns {df.columns.tolist()}")
            return
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        p.alignment = PP_ALIGN.CENTER
        
        # Use EXACT DataFrame column order
        available_columns = list(df.columns)
        
        num_rows = len(df) + 1
        num_cols = len(available_columns)
        
        table_width = Inches(12.0)
        table_height = Inches(0.3 * len(df) + 0.3)
        left = Inches(0.65)
        top = Inches(1.2)
        
        table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table
        
        # Set column widths
        col_widths = [Inches(0.5), Inches(3.0)] + [Inches(0.75)] * (len(available_columns) - 2)
        for col_idx in range(num_cols):
            if col_idx < len(col_widths):
                table.columns[col_idx].width = col_widths[col_idx]
        
        # Header row with exact order
        for col_idx, col_name in enumerate(available_columns):
            cell = table.cell(0, col_idx)
            cell.text = col_name
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Data rows with exact column order
        for row_idx, (index, row) in enumerate(df.iterrows(), start=1):
            is_total_row = index == len(df) - 1
            for col_idx, col_name in enumerate(available_columns):
                cell = table.cell(row_idx, col_idx)
                try:
                    value = row[col_name]
                    cell.text = str(value)
                except (KeyError, ValueError) as e:
                    cell.text = ""
                    logger.warning(f"Error accessing {col_name} at row {index} in customer slide: {e}")
                
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                
                if is_total_row:
                    cell.fill.fore_color.rgb = RGBColor(211, 211, 211)
                    cell.text_frame.paragraphs[0].font.bold = True
                else:
                    if (row_idx - 1) % 2 == 0:
                        cell.fill.fore_color.rgb = RGBColor(221, 235, 247)
                    else:
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
    except Exception as e:
        logger.error(f"Error creating customer PPT slide: {e}")
        raise

def create_nbc_individual_ppt(customer_df, customer_title, sorted_months, financial_year, logo_file=None):
    """Create individual PPT for NBC report only with correct column order."""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide
        create_title_slide(prs, "Number of Billed Customers Report", logo_file)
        
        # Add NBC slide with exact column order
        slide_layout = prs.slide_layouts[6]  # Blank slide
        slide1 = prs.slides.add_slide(slide_layout)
        create_customer_ppt_slide(slide1, customer_df, customer_title, sorted_months)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating NBC PPT: {e}")
        return None

def create_od_ppt_slide(slide, df, title):
    """Create OD PPT slide with exact column order"""
    try:
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), 
            Inches(12), Inches(0.8)
        )
        title_frame = title_shape.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = title
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Use EXACT DataFrame column order
        available_columns = list(df.columns)
        
        rows, cols = len(df) + 1, len(available_columns)
        table_width = Inches(8)
        table_height = Inches(len(df) * 0.4 + 0.5)
        left = Inches(2.0)
        top = Inches(1.5)
        
        table = slide.shapes.add_table(
            rows, cols, 
            left, top, 
            table_width, table_height
        ).table
        
        # Header row with exact order
        for i, col_name in enumerate(available_columns):
            header_cell = table.cell(0, i)
            header_cell.text = col_name
            header_cell.text_frame.paragraphs[0].font.bold = True
            header_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            header_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            header_cell.fill.solid()
            header_cell.fill.fore_color.rgb = RGBColor(0, 114, 188)
        
        # Data rows with exact column order
        for i in range(len(df)):
            for j, col_name in enumerate(available_columns):
                cell = table.cell(i + 1, j)
                value = df.iloc[i][col_name] if col_name in df.columns else ""
                
                if col_name == 'TARGET':
                    cell.text = f"{float(value):.2f}" if isinstance(value, (int, float)) else str(value)
                else:
                    cell.text = str(value)
                
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                row_color = RGBColor(221, 235, 247) if i % 2 == 0 else RGBColor(255, 255, 255)
                if str(df.iloc[i].get('Executive', '')) == 'TOTAL':
                    row_color = RGBColor(211, 211, 211)
                    cell.text_frame.paragraphs[0].font.bold = True
                
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_color
        
        # Set column widths
        col_width = table_width.inches / len(available_columns)
        for i in range(len(available_columns)):
            table.columns[i].width = Inches(col_width)
        
    except Exception as e:
        logger.error(f"Error creating OD PPT slide: {e}")
        raise

def create_od_individual_ppt(od_target_df, od_title, logo_file=None):
    """Create individual PPT for OD Target report only."""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide
        create_title_slide(prs, "OD Target Report", logo_file)
        
        # Add OD Target slide
        slide_layout = prs.slide_layouts[6]  # Blank slide
        slide2 = prs.slides.add_slide(slide_layout)
        create_od_ppt_slide(slide2, od_target_df, od_title)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating OD Target PPT: {e}")
        return None

def create_executive_budget_ppt(results_data, month_title=None, logo_file=None):
    """Create PowerPoint presentation for executive budget vs billed with correct column order"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide
        title = f"Monthly Review Meeting – {month_title}" if month_title else "Executive Budget vs Billed Analysis"
        create_title_slide(prs, title, logo_file)
        
        # Convert results data to DataFrames with proper column order
        budget_vs_billed_qty_df = pd.DataFrame(results_data.get('budget_vs_billed_qty', []))
        budget_vs_billed_value_df = pd.DataFrame(results_data.get('budget_vs_billed_value', []))
        overall_sales_qty_df = pd.DataFrame(results_data.get('overall_sales_qty', []))
        overall_sales_value_df = pd.DataFrame(results_data.get('overall_sales_value', []))
        
        # Add slides with correct processing and column order
        if not budget_vs_billed_qty_df.empty:
            process_df_for_slides(prs, budget_vs_billed_qty_df, "BUDGET AGAINST BILLED (Qty in Mt)", percent_cols=[3], is_consolidated=False)
        
        if not budget_vs_billed_value_df.empty:
            process_df_for_slides(prs, budget_vs_billed_value_df, "BUDGET AGAINST BILLED (Value in Lakhs)", percent_cols=[3], is_consolidated=False)
        
        if not overall_sales_qty_df.empty:
            process_df_for_slides(prs, overall_sales_qty_df, "OVERALL SALES (Qty in Mt)", percent_cols=[], is_consolidated=False)
        
        if not overall_sales_value_df.empty:
            process_df_for_slides(prs, overall_sales_value_df, "OVERALL SALES (Value in Lakhs)", percent_cols=[], is_consolidated=False)
        
        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error creating Executive Budget PPT: {str(e)}")
        raise Exception(f"Error creating Executive Budget PPT: {str(e)}")

def create_executive_od_ppt(results_data, month_title=None, logo_file=None):
    """Create PowerPoint presentation for executive OD vs Collection with correct column order"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide
        title = f"OD Target vs Collection – {month_title}" if month_title else "OD Target vs Collection Analysis"
        create_title_slide(prs, title, logo_file)
        
        # 🎯 FIX: Process OD results with exact column order matching Streamlit output
        if 'od_results' in results_data:
            od_data = results_data['od_results']
        elif isinstance(results_data, list):
            od_data = results_data
        else:
            # Handle direct DataFrame or dict input
            od_data = results_data
            
        if od_data:
            # Convert to DataFrame and maintain exact column order
            if isinstance(od_data, list):
                od_df = pd.DataFrame(od_data)
            else:
                od_df = pd.DataFrame(od_data) if not isinstance(od_data, pd.DataFrame) else od_data.copy()
            
            # 🔍 DEBUG: Log the DataFrame structure
            logger.info(f"OD DataFrame columns: {list(od_df.columns)}")
            logger.info(f"OD DataFrame shape: {od_df.shape}")
            print(f"🔍 OD DataFrame columns: {list(od_df.columns)}")
            
            # Add OD table slide with exact title and percentage columns
            slide_title = f"OD Target vs Collection – {month_title}" if month_title else "OD Target vs Collection"
            slide_title += " (Value in Lakhs)"
            
            # 🎯 STREAMLIT COLUMN ORDER EXPECTED:
            # ['Executive', 'Due Target', 'Collection Achieved', 'Overall % Achieved', 
            #  'For the month Overdue', 'For the month Collection', '% Achieved (Selected Month)']
            
            # Identify percentage columns by name (more reliable than by index)
            percentage_column_names = ['Overall % Achieved', '% Achieved (Selected Month)']
            percent_cols = []
            for i, col in enumerate(od_df.columns):
                if col in percentage_column_names:
                    percent_cols.append(i)
            
            logger.info(f"OD percentage columns identified: {percentage_column_names} at indices {percent_cols}")
            print(f"🔍 OD percentage columns: {percentage_column_names} -> indices {percent_cols}")
            
            # Use splitting logic if needed
            if 'Executive' in od_df.columns:
                process_df_for_slides(prs, od_df, slide_title, percent_cols=percent_cols, is_consolidated=False)
            else:
                add_table_slide(prs, od_df, slide_title, percent_cols=percent_cols, is_consolidated=False)
        
        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error creating Executive OD PPT: {str(e)}")
        raise Exception(f"Error creating Executive OD PPT: {str(e)}")
def validate_ppt_data(results_data):
    """Validate that results data contains the required structure for PPT generation"""
    if not isinstance(results_data, dict):
        raise ValueError("Results data must be a dictionary")
    
    # Check for budget vs billed data
    if 'budget_vs_billed_qty' in results_data:
        required_keys = ['budget_vs_billed_qty', 'budget_vs_billed_value', 'overall_sales_qty', 'overall_sales_value']
        
        for key in required_keys:
            if key not in results_data:
                raise ValueError(f"Missing required data key: {key}")
            
            data = results_data[key]
            if not isinstance(data, list) or len(data) == 0:
                raise ValueError(f"Data for {key} must be a non-empty list")
            
            # Check if first item has required columns
            first_item = data[0]
            if not isinstance(first_item, dict) or 'Executive' not in first_item:
                raise ValueError(f"Data for {key} must contain 'Executive' column")
    
    # Check for OD data
    elif 'od_results' in results_data:
        od_results = results_data['od_results']
        if not isinstance(od_results, list):
            raise ValueError("od_results must be a list")
        
        if not od_results:
            raise ValueError("od_results is empty")
        
        # Check if first row has required columns
        required_columns = [
            'Executive', 'Due Target', 'Collection Achieved', 'Overall % Achieved',
            'For the month Overdue', 'For the month Collection', '% Achieved (Selected Month)'
        ]
        
        first_row = od_results[0]
        for col in required_columns:
            if col not in first_row:
                raise ValueError(f"Missing required column: {col}")
    
    else:
        raise ValueError("Results data must contain either budget vs billed data or OD results data")
    
    return True

def validate_product_growth_ppt_data(results_data):
    """Validate data structure for Product Growth PPT generation"""
    if not isinstance(results_data, dict):
        raise ValueError("Results data must be a dictionary")
    
    if not results_data:
        raise ValueError("Results data is empty")
    
    # Check if each company has required data structure
    required_keys = ['qty_df', 'value_df']
    required_columns_qty = ['PRODUCT GROUP', 'LY_QTY', 'BUDGET_QTY', 'CY_QTY', 'ACHIEVEMENT %']
    required_columns_value = ['PRODUCT GROUP', 'LY_VALUE', 'BUDGET_VALUE', 'CY_VALUE', 'ACHIEVEMENT %']
    
    for company, data in results_data.items():
        for key in required_keys:
            if key not in data:
                raise ValueError(f"Missing required key '{key}' for company '{company}'")
            
            df_data = data[key]
            if not isinstance(df_data, list) or not df_data:
                raise ValueError(f"Data for '{key}' in company '{company}' must be a non-empty list")
            
            # Check columns
            first_row = df_data[0]
            if key == 'qty_df':
                required_cols = required_columns_qty
            else:
                required_cols = required_columns_value
            
            for col in required_cols:
                if col not in first_row:
                    raise ValueError(f"Missing required column '{col}' in {key} for company '{company}'")
    
    return True

def get_actual_slide_count(ppt_buffer):
    """Get actual slide count from generated PPT buffer"""
    try:
        ppt_buffer.seek(0)
        temp_prs = Presentation(ppt_buffer)
        slide_count = len(temp_prs.slides)
        ppt_buffer.seek(0)  # Reset buffer position
        return slide_count
    except Exception as e:
        logger.error(f"Error getting slide count: {e}")
        return 0

def log_ppt_generation_summary(report_type, slide_count, report_count=None):
    """Log PPT generation summary"""
    if report_count:
        logger.info(f"📊 {report_type} PPT generated: {report_count} reports → {slide_count} slides")
        print(f"📊 {report_type} PPT generated: {report_count} reports → {slide_count} slides")
    else:
        logger.info(f"📊 {report_type} PPT generated with {slide_count} slides")
        print(f"📊 {report_type} PPT generated with {slide_count} slides")

def handle_consolidated_ppt_request(request_data):
    """Handle consolidated PPT generation with accurate slide counting and enhanced logging"""
    try:
        reports = request_data.get('reports', [])
        title = request_data.get('title', 'Consolidated Report')
        logo_file = request_data.get('logo_file')
        
        if not reports:
            raise ValueError("No reports provided")
        
        logger.info(f"📋 Starting consolidated PPT generation with {len(reports)} reports")
        print(f"📋 Starting consolidated PPT generation with {len(reports)} reports")
        
        # Log report breakdown
        report_breakdown = {}
        for report in reports:
            category = report.get('category', 'unknown')
            report_breakdown[category] = report_breakdown.get(category, 0) + 1
        
        logger.info(f"📊 Report breakdown: {report_breakdown}")
        print(f"📊 Report breakdown: {report_breakdown}")
        
        # Estimate slide count before generation
        estimated_slides = estimate_slide_count_for_consolidated(reports)
        logger.info(f"📊 Estimated slide count: {estimated_slides}")
        print(f"📊 Estimated slide count: {estimated_slides}")
        
        # Generate PPT with enhanced column ordering
        ppt_buffer = create_consolidated_ppt(reports, logo_file, title)
        
        if ppt_buffer is None:
            raise Exception("Failed to generate PowerPoint presentation")
        
        # Get actual slide count
        actual_slides = get_actual_slide_count(ppt_buffer)
        
        # Log final summary
        log_ppt_generation_summary("Consolidated", actual_slides, len(reports))
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error handling consolidated PPT request: {str(e)}")
        raise

def handle_product_growth_ppt_request(results_data, month_title, logo_file=None, ly_month=None, cy_month=None):
    """Handle Product Growth PPT generation with enhanced logging"""
    try:
        logger.info(f"📋 Starting Product Growth PPT generation")
        print(f"📋 Starting Product Growth PPT generation")
        
        # Validate data
        validate_product_growth_ppt_data(results_data)
        
        # Generate PPT
        ppt_buffer = create_product_growth_ppt(results_data, month_title, logo_file, ly_month, cy_month)
        
        if ppt_buffer is None:
            raise Exception("Failed to generate Product Growth PowerPoint presentation")
        
        # Get actual slide count
        actual_slides = get_actual_slide_count(ppt_buffer)
        company_count = len(results_data)
        
        # Log summary (each company has 2 slides: qty + value)
        log_ppt_generation_summary("Product Growth", actual_slides, company_count * 2)
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error handling Product Growth PPT request: {str(e)}")
        raise

def handle_budget_vs_billed_ppt_request(results_data, month_title, logo_file=None):
    """Handle Budget vs Billed PPT generation with enhanced logging"""
    try:
        logger.info(f"📋 Starting Budget vs Billed PPT generation")
        print(f"📋 Starting Budget vs Billed PPT generation")
        
        # Validate data
        validate_ppt_data(results_data)
        
        # Generate PPT
        ppt_buffer = create_executive_budget_ppt(results_data, month_title, logo_file)
        
        if ppt_buffer is None:
            raise Exception("Failed to generate Budget vs Billed PowerPoint presentation")
        
        # Get actual slide count
        actual_slides = get_actual_slide_count(ppt_buffer)
        
        # Count reports (up to 4: qty budget, value budget, qty sales, value sales)
        report_count = sum(1 for key in ['budget_vs_billed_qty', 'budget_vs_billed_value', 'overall_sales_qty', 'overall_sales_value'] 
                          if key in results_data and results_data[key])
        
        # Log summary
        log_ppt_generation_summary("Budget vs Billed", actual_slides, report_count)
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error handling Budget vs Billed PPT request: {str(e)}")
        raise

def handle_od_collection_ppt_request(results_data, month_title, logo_file=None):
    """Handle OD Collection PPT generation with enhanced logging"""
    try:
        logger.info(f"📋 Starting OD Collection PPT generation")
        print(f"📋 Starting OD Collection PPT generation")
        
        # Validate data
        validate_ppt_data(results_data)
        
        # Generate PPT
        ppt_buffer = create_executive_od_ppt(results_data, month_title, logo_file)
        
        if ppt_buffer is None:
            raise Exception("Failed to generate OD Collection PowerPoint presentation")
        
        # Get actual slide count
        actual_slides = get_actual_slide_count(ppt_buffer)
        
        # Log summary
        log_ppt_generation_summary("OD Collection", actual_slides, 1)
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error handling OD Collection PPT request: {str(e)}")
        raise

def format_currency_value(value, currency_symbol="₹"):
    """Format currency values consistently"""
    if isinstance(value, (int, float)) and not pd.isna(value):
        if abs(value) >= 10000000:  # 1 crore
            return f"{currency_symbol}{value/10000000:.2f}Cr"
        elif abs(value) >= 100000:  # 1 lakh
            return f"{currency_symbol}{value/100000:.2f}L"
        elif abs(value) >= 1000:  # 1 thousand
            return f"{currency_symbol}{value/1000:.2f}K"
        else:
            return f"{currency_symbol}{value:.2f}"
    return str(value) if value is not None else ""

def format_quantity_value(value, unit="Mt"):
    """Format quantity values consistently"""
    if isinstance(value, (int, float)) and not pd.isna(value):
        if abs(value - round(value)) < 0.001:  # Check if essentially an integer
            return f"{int(round(value))} {unit}"
        else:
            return f"{value:.2f} {unit}"
    return str(value) if value is not None else ""

def clean_dataframe_for_ppt(df):
    """Clean DataFrame before PPT generation"""
    try:
        # Remove any completely empty rows
        df = df.dropna(how='all')
        
        # Fill NaN values with empty strings for text columns
        for col in df.columns:
            if df[col].dtype == 'object':
                df[col] = df[col].fillna('')
            else:
                df[col] = df[col].fillna(0)
        
        # Round numeric columns
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        for col in numeric_cols:
            df[col] = df[col].round(2)
        
        return df
        
    except Exception as e:
        logger.error(f"Error cleaning DataFrame: {e}")
        return df

def validate_consolidated_reports_structure(reports):
    """Validate structure of consolidated reports before PPT generation"""
    try:
        if not isinstance(reports, list):
            raise ValueError("Reports must be a list")
        
        if not reports:
            raise ValueError("Reports list is empty")
        
        required_fields = ['df', 'title']
        optional_fields = ['percent_cols', 'category']
        
        for i, report in enumerate(reports):
            if not isinstance(report, dict):
                raise ValueError(f"Report {i} must be a dictionary")
            
            # Check required fields
            for field in required_fields:
                if field not in report:
                    raise ValueError(f"Report {i} missing required field: {field}")
            
            # Validate df field
            df_data = report['df']
            if not isinstance(df_data, (list, pd.DataFrame)):
                raise ValueError(f"Report {i} 'df' must be a list or DataFrame")
            
            if isinstance(df_data, list) and not df_data:
                logger.warning(f"Report {i} has empty data")
            
            # Validate title
            if not isinstance(report['title'], str) or not report['title'].strip():
                raise ValueError(f"Report {i} must have a non-empty title")
        
        return True
        
    except Exception as e:
        logger.error(f"Error validating consolidated reports structure: {e}")
        raise

__all__ = [
    'create_title_slide',
    'add_table_slide', 
    'create_product_growth_ppt',
    'process_df_for_slides',
    'estimate_slide_count_for_consolidated',
    'create_consolidated_ppt',
    'create_customer_ppt_slide',
    'create_nbc_individual_ppt',
    'create_od_ppt_slide',
    'create_od_individual_ppt',
    'create_executive_budget_ppt',
    'create_executive_od_ppt',
    'validate_ppt_data',
    'validate_product_growth_ppt_data',
    'get_actual_slide_count',
    'log_ppt_generation_summary',
    'handle_consolidated_ppt_request',
    'handle_product_growth_ppt_request', 
    'handle_budget_vs_billed_ppt_request',
    'handle_od_collection_ppt_request',
    'format_currency_value',
    'format_quantity_value',
    'clean_dataframe_for_ppt',
    'validate_consolidated_reports_structure'
]