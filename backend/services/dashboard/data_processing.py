import pandas as pd
import numpy as np
import base64
import json
import re
import gc
from datetime import datetime
from io import BytesIO
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
import plotly.graph_objects as go
import plotly.utils
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from werkzeug.utils import secure_filename
import logging
from flask import current_app, jsonify  # <-- ADD THIS
# Add these imports at the top with other imports
from utils.dashboard import helpers
from utils.dashboard.helpers import (
    rename_columns,
    extract_tables,
    find_table_end,
    column_filter,
    make_jsonly_serializable,
    convert_to_numeric,
    safe_sum,
    safe_mean,
    ensure_numeric_data,
    extract_performance_column,
    extract_month_year
)       

# Add this near the top of data_processing.py
BRANCH_EXCLUDE_TERMS = [
    'CHN Total', 'ERD SALES', 'North Total', 'WEST SALES', 'GROUP COMPANIES'
]

def process_monthly_comparison(df, metric, selected_month=None, selected_year=None):
    """Process monthly comparison data for Budget, LY, Act, Gr, Ach"""
    try:
        metric_cols = []
        if metric == 'LY':
            pattern = r'^LY-\w{3}-\d{2}$'  # Exact match for LY-Apr-24
        else:
            pattern = rf'{metric}[-\s]*(?:\w{{3}}[-\s]*\d{{2,4}}|\d{{2,4}}[-\s]*\w{{3}})'  # Matches "Act Apr-23" or "Act-23-Apr"
        
        for col in df.columns:
            col_str = str(col).lower().replace(",", "").replace("–", "-")
            if re.search(pattern, col_str, re.IGNORECASE) and column_filter(col, selected_month, selected_year):
                metric_cols.append(col)
        
        if not metric_cols:
            logging.warning(f"No {metric} columns found")
            return None
        
        first_col = df.columns[0]
        comparison_data = df[[first_col] + metric_cols].copy()
        
        # Clean column names
        clean_labels = []
        for col in metric_cols:
            col_str = str(col).strip()
            if metric == 'LY' and re.match(r'^LY-\w{3}-\d{2}$', col_str, re.IGNORECASE):
                clean_labels.append(col_str)
            else:
                month_match = re.search(r'(\w{3})[-\s]*(\d{2,4})', col_str, re.IGNORECASE)
            if month_match:
                month, year = month_match.groups()
                clean_labels.append(f"{metric} {month.capitalize()}-{year[-2:]}")
            else:
                clean_labels.append(col_str)
        
        comparison_data.columns = [first_col] + clean_labels
        
        # Convert to numeric
        for col in clean_labels:
            comparison_data[col] = pd.to_numeric(
                comparison_data[col].astype(str).str.replace(',', ''),
                errors='coerce'
            )
        
        # Melt for visualization
        chart_data = comparison_data.melt(
            id_vars=first_col,
            var_name="Month",
            value_name=metric
        ).dropna()
        
        if chart_data.empty:
            logging.warning(f"No valid data for {metric} comparison")
            return None
        
        # Calculate metrics
        metrics = {
            'total': float(chart_data[metric].sum()),
            'average': float(chart_data[metric].mean()),
            'min': float(chart_data[metric].min()),
            'max': float(chart_data[metric].max())
        }
        
        return {
            'data': make_jsonly_serializable(chart_data),
            'metrics': metrics,
            'x_col': 'Month',
            'y_col': metric,
            'color_override': '#FF8C00' if metric == "Act" else '#2E86AB'
        }
        
    except Exception as e:
        logging.error(f"Error processing {metric} data: {str(e)}", exc_info=True)
        return None
    
def create_plotly_chart(data, x_col, y_col, chart_type, title, color_override=None):
    """Enhanced version with Streamlit's settings"""
    default_color = color_override if color_override else '#2E86AB'
    
    layout_config = {
        'title': {
            'text': title,
            'x': 0.5,
            'xanchor': 'center',
            'font': {'size': 20, 'family': 'Arial, sans-serif'}
        },
        'font': {'size': 14, 'family': 'Arial, sans-serif'},
        'plot_bgcolor': 'white',
        'paper_bgcolor': 'white',
        'height': 600,
        'margin': {'l': 80, 'r': 80, 't': 100, 'b': 80},
        'hovermode': 'closest',
        'showlegend': True if chart_type == 'pie' else False,
        'xaxis': {
            'tickangle': 0,  # Force straight labels
            'automargin': True  # Prevent label cutoff
        }
    }
    
    if chart_type == 'bar':
        fig = px.bar(data, x=x_col, y=y_col, 
                    color_discrete_sequence=[default_color],
                    text_auto='.2s')  # Add compact value labels
        
        fig.update_traces(
            hovertemplate='<b>%{x}</b><br>Value: %{y:,.0f}<extra></extra>',
            textposition='outside'
        )
        fig.update_layout(bargap=0.2)
        
    elif chart_type == 'line':
        fig = px.line(data, x=x_col, y=y_col, 
                     markers=True, 
                     color_discrete_sequence=[default_color])
        fig.update_traces(
            line={'width': 4}, 
            marker={'size': 10},
            hovertemplate='<b>%{x}</b><br>Value: %{y:,.0f}<extra></extra>'
        )
    elif chart_type == 'pie':
        fig = px.pie(data, values=y_col, names=x_col)
        fig.update_traces(
            textposition='inside',
            textinfo='percent+label',
            hovertemplate='<b>%{label}</b><br>Value: %{value:,.0f}<br>Percentage: %{percent:.1%}<extra></extra>',
            textfont={'size': 14, 'family': 'Arial, sans-serif'}
        )
    
    fig.update_layout(**layout_config)
    fig.update_xaxes(
        title_font={'size': 16, 'family': 'Arial, sans-serif'},
        tickfont={'size': 14},
        tickangle=0,
        automargin=True
    )
    fig.update_yaxes(
        title_font={'size': 16, 'family': 'Arial, sans-serif'},
        tickfont={'size': 14},
        tickformat=',.0f'
    )
    
    config = {
        'displayModeBar': True,
        'displaylogo': False,
        'modeBarButtonsToRemove': ['pan2d', 'lasso2d', 'select2d'],
        'toImageButtonOptions': {
            'format': 'png',
            'filename': 'chart',
            'height': 800,
            'width': 1200,
            'scale': 2
        }
    }
    
    return fig, config    
    
def process_performance_data(df, performance_type, selected_filter=None):
    """Process performance data for branches or products with comprehensive error handling"""
    if performance_type not in ['branch', 'product']:
        logging.error(f"Invalid performance type: {performance_type}")
        return None
    
    first_col = df.columns[0] if len(df.columns) > 0 else None
    if not first_col:
        logging.error("No columns found in DataFrame")
        return None
    
    # Define exclusion terms based on type
    if performance_type == 'product':
        exclude_terms = BRANCH_EXCLUDE_TERMS + ['TOTAL SALES', 'GRAND TOTAL', 'OVERALL TOTAL', 'REGIONS', 'NORTH TOTAL','SOUTH TOTAL', 'EAST TOTAL']
    else:  # product
        exclude_terms = ['TOTAL', 'GRAND TOTAL', 'OVERALL TOTAL', 'TOTAL SALES']
    
    try:
        # Filter out excluded rows with regex pattern
        pattern = '|'.join([re.escape(term) for term in exclude_terms])
        filtered_df = df[
            ~df[first_col].astype(str).str.contains(pattern, case=False, na=False, regex=True)
        ].copy()
        
        if filtered_df.empty:
            logging.warning(f"All data filtered out by exclude terms for {performance_type}")
            return None
        
        # Find and validate performance column with is_product flag
        performance_col = extract_performance_column(filtered_df, is_product=(performance_type == 'product'))
        if not performance_col:
            logging.warning(f"No valid performance column found for {performance_type}")
            return None
        
        # Convert to numeric with comprehensive cleaning
        filtered_df[performance_col] = pd.to_numeric(
            filtered_df[performance_col].astype(str).str.replace(r'[^\d.-]', '', regex=True),
            errors='coerce'
        )
        filtered_df = filtered_df.dropna(subset=[performance_col])
        
        if filtered_df.empty:
            logging.warning(f"No valid numeric data after conversion for {performance_type}")
            return None
        
        # Apply selected filter if provided
        if selected_filter and selected_filter != "Select All":
            filtered_df = filtered_df[filtered_df[first_col].astype(str) == selected_filter]
            if filtered_df.empty:
                logging.warning(f"No data matches selected filter: {selected_filter}")
                return None
        
        # Prepare performance data with validation
        performance_data = []
        for _, row in filtered_df.iterrows():
            try:
                name = str(row[first_col]).strip()
                value = float(row[performance_col])
                if name and name != 'nan' and not pd.isna(value):
                    performance_data.append({
                        'Name': name,
                        'Performance': value
                    })
            except Exception as e:
                logging.warning(f"Error processing row {row}: {e}")
                continue
        
        if not performance_data:
            logging.warning("No valid performance data after processing")
            return None
        
        # Create DataFrame and sort
        result_df = pd.DataFrame(performance_data)
        result_df = result_df.sort_values('Performance', ascending=False)
        result_df = make_jsonly_serializable(result_df)
        
        # Calculate metrics with error handling
        try:
            total = float(result_df['Performance'].sum())
            avg = float(result_df['Performance'].mean())
            top_item = result_df.iloc[0].to_dict() if len(result_df) > 0 else {'Name': 'N/A', 'Performance': 0}
            top_5 = result_df.head(5).replace({np.nan: None}).to_dict(orient='records')
            bottom_5 = result_df.tail(5).replace({np.nan: None}).to_dict(orient='records')
        except Exception as e:
            logging.error(f"Error calculating metrics: {e}")
            return None
        
        return {
            'data': result_df,
            'metrics': {
                'total': total,
                'average': avg,
                'top_item': top_item,
                'top_5': top_5,
                'bottom_5': bottom_5
            },
            'x_col': 'Name',
            'y_col': 'Performance'
        }
        
    except Exception as e:
        logging.error(f"Error processing {performance_type} data: {str(e)}", exc_info=True)
        return None

def create_plotly_chart(data, x_col, y_col, chart_type, title, color_override=None):
    """Enhanced version with Streamlit's settings"""
    default_color = color_override if color_override else '#2E86AB'
    
    layout_config = {
        'title': {
            'text': title,
            'x': 0.5,
            'xanchor': 'center',
            'font': {'size': 20, 'family': 'Arial, sans-serif'}
        },
        'font': {'size': 14, 'family': 'Arial, sans-serif'},
        'plot_bgcolor': 'white',
        'paper_bgcolor': 'white',
        'height': 600,
        'margin': {'l': 80, 'r': 80, 't': 100, 'b': 80},
        'hovermode': 'closest',
        'showlegend': True if chart_type == 'pie' else False,
        'xaxis': {
            'tickangle': 0,  # Force straight labels
            'automargin': True  # Prevent label cutoff
        }
    }
    
    if chart_type == 'bar':
        fig = px.bar(data, x=x_col, y=y_col, 
                    color_discrete_sequence=[default_color],
                    text_auto='.2s')  # Add compact value labels
        
        fig.update_traces(
            hovertemplate='<b>%{x}</b><br>Value: %{y:,.0f}<extra></extra>',
            textposition='outside'
        )
        fig.update_layout(bargap=0.2)
        
    elif chart_type == 'line':
        fig = px.line(data, x=x_col, y=y_col, 
                     markers=True, 
                     color_discrete_sequence=[default_color])
        fig.update_traces(
            line={'width': 4}, 
            marker={'size': 10},
            hovertemplate='<b>%{x}</b><br>Value: %{y:,.0f}<extra></extra>'
        )
    elif chart_type == 'pie':
        fig = px.pie(data, values=y_col, names=x_col)
        fig.update_traces(
            textposition='inside',
            textinfo='percent+label',
            hovertemplate='<b>%{label}</b><br>Value: %{value:,.0f}<br>Percentage: %{percent:.1%}<extra></extra>',
            textfont={'size': 14, 'family': 'Arial, sans-serif'}
        )
    
    fig.update_layout(**layout_config)
    fig.update_xaxes(
        title_font={'size': 16, 'family': 'Arial, sans-serif'},
        tickfont={'size': 14},
        tickangle=0,
        automargin=True
    )
    fig.update_yaxes(
        title_font={'size': 16, 'family': 'Arial, sans-serif'},
        tickfont={'size': 14},
        tickformat=',.0f'
    )
    
    config = {
        'displayModeBar': True,
        'displaylogo': False,
        'modeBarButtonsToRemove': ['pan2d', 'lasso2d', 'select2d'],
        'toImageButtonOptions': {
            'format': 'png',
            'filename': 'chart',
            'height': 800,
            'width': 1200,
            'scale': 2
        }
    }
    
    return fig, config

def create_matplotlib_chart(chart_data, x_col, y_col, chart_type, title, color_override=None, is_product_chart=False):
    try:
        plt.style.use('seaborn')
        plt.figure(figsize=(10, 6))
        plt.rcParams['xtick.major.pad'] = 10
        plt.rcParams['axes.labelpad'] = 10

        if chart_type == 'bar':
            if 'Metric' in chart_data.columns:
                pivot_df = chart_data.pivot(index=x_col, columns='Metric', values=y_col)
                pivot_df.plot(kind='bar', ax=plt.gca(), color=['#2E86AB', '#FF8C00'], width=0.8)
            else:
                plt.bar(chart_data[x_col], chart_data[y_col], color=color_override or '#2E86AB', width=0.6)
            
            # Conditional label rotation
            if is_product_chart:
                plt.xticks(rotation=45, ha='right', fontsize=12)  # Slanted for products
            else:
                plt.xticks(rotation=0, ha='center', fontsize=12)  # Straight for others

        elif chart_type == 'line':
            if 'Metric' in chart_data.columns:
                for name, group in chart_data.groupby('Metric'):
                    plt.plot(
                        group[x_col],
                        group[y_col],
                        marker='o',
                        markersize=8,
                        linewidth=2,
                        label=name,
                        color='#FF8C00' if 'Actual' in str(name) else '#2E86AB'
                    )
                plt.legend(fontsize=12)
            else:
                plt.plot(
                    chart_data[x_col],
                    chart_data[y_col],
                    marker='o',
                    markersize=8,
                    linewidth=2,
                    color=color_override or '#2E86AB'
                )
            
            # Conditional label rotation
            if is_product_chart:
                plt.xticks(rotation=45, ha='right', fontsize=12)
            else:
                plt.xticks(rotation=0, ha='center', fontsize=12)

        elif chart_type == 'pie':
            plt.pie(
                chart_data[y_col],
                labels=chart_data[x_col],
                autopct='%1.1f%%',
                startangle=90,
                textprops={'fontsize': 12},
                colors=['#2E86AB', '#FF8C00', '#A23B72', '#F18F01', '#C73E1D']
            )
            plt.axis('equal')

        plt.title(title, fontsize=16, pad=20)
        plt.tight_layout()

        img_buffer = BytesIO()
        plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
        img_buffer.seek(0)
        plt.close()
        return img_buffer

    except Exception as e:
        current_app.logging.error(f"Error creating matplotlib chart: {str(e)}", exc_info=True)
        return None

def create_ppt_with_chart(title, chart_data, x_col, y_col, chart_type='bar', color_override=None, selected_filter=None, is_product_chart=False):
    try:
        current_app.logging.debug(f"Creating PPT: title={title}, chart_type={chart_type}, x_col={x_col}, y_col={y_col}")
        
        # Validate inputs
        if not isinstance(chart_data, pd.DataFrame):
            chart_data = pd.DataFrame(chart_data)
        if chart_data.empty:
            current_app.logging.error("Chart data is empty")
            return None
        if x_col not in chart_data.columns or y_col not in chart_data.columns:
            current_app.logging.error(f"Invalid columns: x_col={x_col}, y_col={y_col}, available={list(chart_data.columns)}")
            return None

        # Ensure numeric data
        chart_data = make_jsonly_serializable(chart_data)
        if not ensure_numeric_data(chart_data, y_col):
            current_app.logging.error(f"No valid numeric data in {y_col}")
            return None

        # Handle pie chart positive values
        if chart_type == 'pie':
            chart_data = chart_data[chart_data[y_col] > 0].copy()
            if chart_data.empty:
                current_app.logging.error("No positive values for pie chart")
                return None

        # Create PPT
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
        if selected_filter and selected_filter != "Select All":
            title = f"{title} - {selected_filter}"
        
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
        else:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = title
        
        # Generate chart
        img_buffer = create_matplotlib_chart(chart_data, x_col, y_col, chart_type, title, color_override, is_product_chart=is_product_chart)

        if img_buffer:
            slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
        else:
            current_app.logging.warning("Matplotlib failed, trying Plotly fallback")
            try:
                fig, _ = create_plotly_chart(chart_data, x_col, y_col, chart_type, title, color_override)
                img_bytes = fig.to_image(format="png", width=1200, height=800, scale=2)
                img_buffer = BytesIO(img_bytes)
                slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
            except Exception as plotly_error:
                current_app.logging.error(f"Plotly fallback failed: {str(plotly_error)}")
                text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
                text_frame = text_box.text_frame
                text_frame.text = f"Chart: {title}\n(Image generation failed)"
        
        ppt_bytes = BytesIO()
        ppt.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes
    except Exception as e:
        current_app.logging.error(f"Could not create PPT: {str(e)}", exc_info=True)
        return None


def create_master_ppt_with_matplotlib(all_data, table_name, selected_sheet, visual_type, selected_filter=None):
    """Create master PPT with all visualizations using matplotlib"""
    try:
        master_ppt = Presentation()
        
        # Add title slide
        title_slide_layout = master_ppt.slide_layouts[0]
        title_slide = master_ppt.slides.add_slide(title_slide_layout)
        
        # Add title with filter info if provided
        title_text = f"Complete Analysis Report - {table_name}"
        if selected_filter and selected_filter != "Select All":
            title_text += f" - {selected_filter}"
        title_slide.shapes.title.text = title_text
        
        # Add subtitle with sheet info and date
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = f"Sheet: {selected_sheet}\nGenerated on: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        # Process each visualization in order
        for item in all_data:
            label = item['label']
            data = item['data']
            
            try:
                # Skip if no data
                if not data:
                    current_app.logging.debug(f"Skipping {label}: No data provided")
                    continue
                    
                # Convert to DataFrame if not already
                if not isinstance(data, pd.DataFrame):
                    try:
                        data = pd.DataFrame(data)
                    except Exception as e:
                        current_app.logging.warning(f"Could not convert {label} data to DataFrame: {str(e)}")
                        continue
                
                # Clean and validate data
                data = make_jsonly_serializable(data)
                if data.empty:
                    current_app.logging.warning(f"Skipping {label}: Data is empty after cleaning")
                    continue
                
                # Determine columns with flexible fallback logic
                x_col = y_col = None
                chart_type = visual_type.lower().replace(" chart", "")
                
                # Special handling for Budget vs Actual pie chart
                if label == "Budget vs Actual" and item.get('is_pie_chart'):
                    # Create slide
                    slide = master_ppt.slides.add_slide(master_ppt.slide_layouts[5])
                    slide_title = f"{label} - {table_name}"
                    if selected_filter and selected_filter != "Select All":
                        slide_title += f" - {selected_filter}"
                    slide.shapes.title.text = slide_title
                    
                    # Create pie chart
                    plt.figure(figsize=(10, 6))
                    plt.pie(
                        data['Value'],
                        labels=data['Metric'],
                        autopct='%1.1f%%',
                        startangle=90,
                        colors=['#2E86AB', '#FF8C00'],
                        textprops={'fontsize': 12}
                    )
                    plt.axis('equal')
                    plt.title(slide_title, pad=20)
                    
                    # Save plot to buffer
                    img_buffer = BytesIO()
                    plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
                    img_buffer.seek(0)
                    plt.close()
                    
                    # Add image to slide
                    slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
                    continue
                
                # Performance charts - needs Name, Performance
                elif label in ["Branch Performance", "Product Performance"]:
                    if len(data.columns) >= 2:
                        x_col = data.columns[0]
                        y_col = 'Performance' if 'Performance' in data.columns else data.columns[1]
                
                # Monthwise charts - needs Month, Value
                elif label in ["Branch Monthwise", "Product Monthwise"]:
                    if len(data.columns) >= 2:
                        x_col = 'Month' if 'Month' in data.columns else data.columns[0]
                        y_col = 'Value' if 'Value' in data.columns else data.columns[1]
                
                # YTD charts - needs Period, [metric]
                elif "YTD" in label:
                    y_col_name = label.replace("YTD ", "")
                    if len(data.columns) >= 2:
                        x_col = 'Period' if 'Period' in data.columns else data.columns[0]
                        y_col = y_col_name if y_col_name in data.columns else data.columns[1]
                
                # Monthly comparisons (Budget, LY, Act, Gr, Ach)
                elif label in ["Budget", "LY", "Act", "Gr", "Ach"]:
                    if len(data.columns) >= 2:
                        x_col = data.columns[0]
                        y_col = label if label in data.columns else data.columns[1]
                
                # Fallback for all other cases
                if x_col is None or y_col is None:
                    if len(data.columns) >= 2:
                        x_col = data.columns[0]
                        y_col = data.columns[1]
                    else:
                        current_app.logging.warning(f"Skipping {label}: Insufficient columns")
                        continue
                
                # Ensure numeric data
                try:
                    data[y_col] = pd.to_numeric(data[y_col].astype(str).str.replace(',', ''), errors='coerce')
                    data = data.dropna(subset=[y_col])
                except Exception as e:
                    current_app.logging.warning(f"Could not convert {y_col} to numeric for {label}: {str(e)}")
                    continue
                
                if data.empty:
                    current_app.logging.warning(f"Skipping {label}: No numeric data after conversion")
                    continue
                
                # Skip pie chart filtering (allow all values)
                if chart_type == 'pie':
                    data = data[data[y_col] > 0].copy()
                    if data.empty:
                        current_app.logging.warning(f"Skipping {label}: No positive values for pie chart")
                        continue
                
                # Create slide
                slide_layout = master_ppt.slide_layouts[6]  # Blank layout
                slide = master_ppt.slides.add_slide(slide_layout)
                
                # Set slide title
                slide_title = f"{label} Analysis - {table_name}"
                if selected_filter and selected_filter != "Select All":
                    slide_title += f" - {selected_filter}"
                
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title_frame = title_shape.text_frame
                title_frame.text = slide_title
                title_frame.paragraphs[0].font.size = Pt(24)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Generate chart image
                try:
                    img_buffer = create_matplotlib_chart(
                        data, 
                        x_col, 
                        y_col, 
                        chart_type, 
                        slide_title,
                        '#FF8C00' if 'Act' in label else '#2E86AB'
                    )
                    
                    if img_buffer:
                        slide.shapes.add_picture(
                            img_buffer, 
                            Inches(0.5), 
                            Inches(1.2), 
                            width=Inches(9), 
                            height=Inches(6.5)
                        )
                    else:
                        text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
                        text_frame = text_box.text_frame
                        text_frame.text = f"Chart: {label}\n(Image generation failed)"
                        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                except Exception as chart_error:
                    current_app.logging.error(f"Chart generation error for {label}: {str(chart_error)}")
                    continue
                
            except Exception as e:
                current_app.logging.error(f"Error creating chart for {label}: {str(e)}", exc_info=True)
                continue
        
        # Save to buffer
        master_ppt_buffer = BytesIO()
        master_ppt.save(master_ppt_buffer)
        master_ppt_buffer.seek(0)
        
        return master_ppt_buffer.getvalue()
        
    except Exception as e:
        current_app.logging.error(f"Could not create master PPT: {str(e)}", exc_info=True)
        return None
    

def process_branch_performance(df, visual_type, selected_branch=None):
    """
    Process branch performance data for visualization.
    
    Args:
        df: Input pandas DataFrame
        visual_type: Type of chart ('Bar Chart', 'Line Chart', 'Pie Chart')
        selected_branch: Optional branch filter
    
    Returns:
        Tuple of (JSON-serializable dict, HTTP status code)
    """
    try:
        # Define exclusion terms to match frontend
        BRANCH_EXCLUDE_TERMS = [
            'CHN Total', 'ERD SALES', 'WEST SALES', 'GROUP COMPANIES',
            'TOTAL SALES', 'GRAND TOTAL', 'OVERALL TOTAL',
            'NORTH TOTAL', 'SOUTH TOTAL', 'EAST TOTAL', 'REGIONS'
        ]

        # Check if sheet is for branch analysis
        is_branch_analysis = any(col.lower().startswith('act') for col in df.columns) and \
                            'SALES in MT' in df.columns

        if not is_branch_analysis:
            logging.warning("Branch analysis not applicable for this sheet")
            return jsonify({
                'error': "Branch analysis not applicable for this sheet"
            }), 400

        # Find YTD Act column
        ytd_act_col = None
        for col in df.columns:
            col_str = str(col).strip().lower().replace("–", "-")
            if re.search(r'(act-ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)|ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)\s*act)', col_str, re.IGNORECASE):
                if column_filter(col):
                    ytd_act_col = col
                    break

        # Fallback to any Act column
        if not ytd_act_col:
            for col in df.columns:
                col_str = str(col).strip().lower().replace("–", "-")
                if col_str.startswith('act') and not 'ytd' in col_str and column_filter(col):
                    ytd_act_col = col
                    break

        # Fallback to first numeric column (skip first column)
        if not ytd_act_col:
            for col in df.columns[1:]:
                if column_filter(col):
                    try:
                        if df[col].dtype in ['float64', 'int64'] or \
                           df[col].astype(str).str.replace(r'[^\d.-]', '', regex=True).str.isnumeric().any():
                            ytd_act_col = col
                            break
                    except:
                        continue

        if not ytd_act_col:
            logging.warning(f"Could not find valid performance column. Available columns: {list(df.columns)}")
            return jsonify({
                'error': "Could not find valid performance column",
                'columns': list(df.columns)
            }), 400

        first_col = df.columns[0]
        regions_df = df[~df[first_col].str.contains('|'.join(BRANCH_EXCLUDE_TERMS), na=False, case=False)].copy()
        regions_df = regions_df[~regions_df[first_col].str.contains(
            r'^(?:TOTAL SALES|GRAND TOTAL|OVERALL TOTAL|NORTH TOTAL|SOUTH TOTAL|EAST TOTAL|REGIONS)', 
            na=False, case=False, regex=True
        )].copy()
        regions_df = regions_df.dropna(subset=[first_col, ytd_act_col])

        if selected_branch and selected_branch != "Select All":
            regions_df = regions_df[regions_df[first_col].astype(str).str.strip() == selected_branch]

        if regions_df.empty:
            logging.warning(f"No branch data available after filtering. First col: {first_col}, Performance col: {ytd_act_col}")
            return jsonify({
                'error': "No branch data available after filtering",
                'columns': list(df.columns)
            }), 400

        regions_df[ytd_act_col] = pd.to_numeric(
            regions_df[ytd_act_col].astype(str).str.replace(r'[^\d.-]', '', regex=True),
            errors='coerce'
        )
        regions_df = regions_df.dropna(subset=[ytd_act_col])

        if regions_df.empty:
            logging.warning("No numeric data available for branch performance")
            return jsonify({
                'error': "No numeric data available for branch performance"
            }), 400

        regions_df = regions_df.sort_values(by=ytd_act_col, ascending=False)
        clean_regions_df = regions_df[[first_col, ytd_act_col]].copy()
        clean_regions_df.columns = ['Branch', 'Performance']

        chart_data = make_jsonly_serializable(clean_regions_df)

        metrics = {
            'top_performer': {
                'branch': chart_data.iloc[0]['Branch'] if not chart_data.empty else None,
                'performance': float(chart_data.iloc[0]['Performance']) if not chart_data.empty else 0
            },
            'total_performance': float(chart_data['Performance'].sum()),
            'average_performance': float(chart_data['Performance'].mean()),
            'top_5': chart_data.head(5).to_dict(orient='records'),
            'bottom_5': chart_data.tail(5).to_dict(orient='records')
        }

        ppt_type = 'bar' if visual_type == "Bar Chart" else 'line' if visual_type == "Line Chart" else 'pie'

        response = {
            'data': chart_data.to_dict(orient='records'),
            'metrics': metrics,
            'x_col': 'Branch',
            'y_col': 'Performance',
            'chart_type': ppt_type
        }

        logging.debug(f"Processed branch performance data: {len(chart_data)} rows")
        return jsonify(response), 200

    except Exception as e:
        logging.error(f"Error processing branch performance: {str(e)}", exc_info=True)
        return jsonify({
            'error': f"Error processing branch performance: {str(e)}"
        }), 500

def process_branch_monthwise(df, visual_type, selected_branch=None):
    """Process monthwise branch data specifically for Region Wise Analysis sheets"""
    try:
        # Find month-wise columns
        month_cols = []
        months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
        
        for col in df.columns:
            col_str = str(col)
            if any(month in col_str for month in months):
                if column_filter(col, None, None):  # No month/year filter for this
                    month_cols.append(col)
        
        if not month_cols:
            logging.debug("No month-wise columns found for branch analysis")
            return None
        
        # Prepare data
        first_col = df.columns[0]
        month_data = []
        
        for col in month_cols:
            month = extract_month_year(col)
            try:
                value = safe_sum(df[col])
                month_data.append({
                    'Month': month,
                    'Value': float(value)
                })
            except:
                continue
        
        if not month_data:
            logging.debug("No valid month-wise data after processing")
            return None
        
        # Create visualization
        fig, config = create_plotly_chart(
            pd.DataFrame(month_data),
            'Month',
            'Value',
            visual_type,
            'Branch Monthwise Analysis'
        )
        
        # Calculate metrics
        metrics = {
            'best_month': max(month_data, key=lambda x: x['Value']) if month_data else None,
            'total': sum(d['Value'] for d in month_data),
            'average': sum(d['Value'] for d in month_data) / len(month_data) if month_data else 0
        }
        
        return {
            'figure': fig,
            'data': month_data,
            'metrics': metrics
        }
        
    except Exception as e:
        logging.error(f"Error processing branch monthwise data: {str(e)}", exc_info=True)
        return None    
    
def process_budget_vs_actual(df, visual_type):
    """Process budget vs actual comparison data"""
    try:
        budget_cols = [col for col in df.columns if 'budget' in str(col).lower()]
        actual_cols = [col for col in df.columns if 'act' in str(col).lower() and 'budget' not in str(col).lower()]
        
        if not budget_cols or not actual_cols:
            return None
        
        # Find matching months between budget and actual
        month_data = []
        months_found = set()
        
        for col in budget_cols + actual_cols:
            month = extract_month_year(col)
            if month not in months_found:
                months_found.add(month)
                is_budget = 'budget' in str(col).lower()
                
                try:
                    value = df[col].sum() if col in df.columns else 0
                    month_data.append({
                        'Month': month,
                        'Value': float(value),
                        'Metric': 'Budget' if is_budget else 'Actual'
                    })
                except:
                    continue
        
        if not month_data:
            return None
        
        chart_data = pd.DataFrame(month_data)
        chart_data = make_jsonly_serializable(chart_data)
        
        fig, config = create_plotly_chart(
            chart_data,
            'Month',
            'Value',
            visual_type,
            'Budget vs Actual Analysis'
        )
        
        return {
            'figure': json.loads(json.dumps(fig.to_dict(), cls=plotly.utils.PlotlyJSONEncoder)),
            'data': chart_data.to_dict(orient='records'),
            'x_col': 'Month',
            'y_col': 'Value',
            'ppt_type': visual_type,
            'color_override': None
        }
        
    except Exception as e:
        logging.error(f"Error processing budget vs actual: {str(e)}", exc_info=True)
        return None
           


def is_monthwise_sheet(sheet_name):
    """Check if the sheet is a Sales Analysis Month-wise sheet"""
    sheet_lower = sheet_name.lower().strip()
    return any(
        sheet_lower == term.lower() 
        for term in [
            "Sales Analysis Month wise",  # Exact match for your sheet name
            "Sales Analysis Month-wise",
            "Month-wise Sales"
        ]
    )

def identify_ytd_columns(df, metric):
    """Improved YTD column identification matching Streamlit logic"""
    ytd_cols = []
    patterns = [
        rf'{metric}-YTD-\d{{2}}-\d{{2}}\s*\(.*?\)',  # Act-YTD-23-24 (Apr-Jun)
        rf'YTD-\d{{2}}-\d{{2}}\s*\(.*?\)\s*{metric}', # YTD-23-24 (Apr-Jun) Act
        rf'{metric}-YTD\s*\d{{4}}\s*\(.*?\)',         # Act-YTD 2024 (Apr-Jun)
        rf'YTD\s*\d{{4}}\s*\(.*?\)\s*{metric}'        # YTD 2024 (Apr-Jun) Act
    ]
    
    for col in df.columns:
        col_str = str(col).replace(",", "").replace("–", "-")
        if any(re.search(p, col_str, re.IGNORECASE) for p in patterns):
            ytd_cols.append(col)
    
    return ytd_cols

def add_ytd_slide(ppt, title, chart_data, x_col, y_col, chart_type):
    """Add a YTD slide with proper formatting"""
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    
    # Title with YTD indicator
    title_shape = slide.shapes.title
    title_shape.text = f"YTD {title}"
    title_shape.text_frame.paragraphs[0].font.size = Pt(24)
    
    # Create chart with fiscal year ordering
    fig, ax = plt.subplots(figsize=(10, 6))
    
    if chart_type == 'bar':
        chart_data.plot.bar(
            x=x_col, 
            y=y_col, 
            ax=ax,
            color=['#FF8C00' if 'Act' in title else '#2E86AB']
        )
        plt.xticks(rotation=45, ha='right')
    elif chart_type == 'line':
        chart_data.plot.line(
            x=x_col,
            y=y_col,
            ax=ax,
            marker='o',
            color='#FF8C00' if 'Act' in title else '#2E86AB',
            linewidth=3
        )
    
    ax.set_title(f"YTD {title}", pad=20)
    ax.grid(True, alpha=0.3)
    
    # Save to buffer and add to slide
    img_buffer = BytesIO()
    plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
    img_buffer.seek(0)
    slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
    plt.close()


def clean_ytd_column_name(col_name):
    """
    Standardizes YTD column names while preserving metric type (Budget/LY/Act/Gr/Ach)
    """
    try:
        if pd.isna(col_name):
            return "Unnamed"
            
        col_str = str(col_name).strip().replace("–", "-")  # Normalize dashes
        
        # Match all YTD patterns with metric types
        ytd_match = re.search(
            r'(?P<prefix>Budget|LY|Act|Gr|Ach)?[-\s]*(?P<ytd>YTD[-\s]*(?P<start>\d{2,4})[-\s]*(?P<end>\d{2,4})\s*\((?P<range>.*?)\)',
            col_str,
            re.IGNORECASE
        )
        
        if ytd_match:
            groups = ytd_match.groupdict()
            prefix = groups.get('prefix', '').strip()
            ytd_str = f"YTD-{groups['start']}-{groups['end']}({groups['range']})"
            
            # Only add prefix if it exists and isn't already part of the YTD string
            if prefix and prefix not in ytd_str:
                return f"{prefix} {ytd_str}"
            return ytd_str
            
        return col_str  # Return original if no YTD pattern matched
        
    except Exception as e:
        logging.warning(f"Could not clean YTD column '{col_name}': {str(e)}")
        return str(col_name)
    

    
def process_ytd_comparison(df, metric, selected_year=None):
    """
    Processes YTD data exactly matching Streamlit logic with Flask adaptations
    Handles: Budget, LY, Act, Gr, Ach with fiscal year sorting and filtering
    """
    try:
        # 1. Column Identification (matches Streamlit's regex patterns)
        ytd_cols = []
        patterns = [
            rf'{metric}-YTD-\d{{2,4}}-\d{{2,4}}\s*\(.*?\)',
        rf'YTD-\d{{2,4}}-\d{{2,4}}\s*\(.*?\)\s*{metric}',
        rf'{metric}-YTD\s*\d{{4}}\s*\(.*?\)',
        rf'YTD\s*\d{{4}}\s*\(.*?\)\s*{metric}',
        # More flexible patterns
        rf'\b{metric}\b.*YTD',
        rf'YTD.*\b{metric}\b',
        rf'\b{metric}\b.*\d{{4}}'
        ]
        
        for col in df.columns:
            col_str = str(col).replace(",", "").replace("–", "-")
            if any(re.search(p, col_str, re.IGNORECASE) for p in patterns):
                if selected_year and selected_year != "Select All":
                    if f"-{selected_year[-2:]}" in col_str or f"-{selected_year}" in col_str:
                        ytd_cols.append(col)
                else:
                    ytd_cols.append(col)

        if not ytd_cols:
            logging.debug(f"No YTD columns found for {metric}")
            return None

        # 2. Data Extraction
        first_col = df.columns[0]
        comparison_data = df[[first_col] + ytd_cols].copy()

        # 3. Column Label Cleaning (matches Streamlit format)
        clean_labels = []
        for col in ytd_cols:
            col_str = str(col)
            # Extract fiscal year and period
            match = re.search(r'(\d{2,4})[-–](\d{2,4})\s*\((.*?)\)', col_str)
            if match:
                start, end, months = match.groups()
                clean_labels.append(f"YTD {start}-{end} ({months})")
            else:
                clean_labels.append(col_str)
        
        comparison_data.columns = [first_col] + clean_labels

        # 4. Numeric Conversion (with Streamlit's comma handling)
        for col in clean_labels:
            comparison_data[col] = pd.to_numeric(
                comparison_data[col].astype(str).str.replace(r'[^\d.-]', '', regex=True),
                errors='coerce'
            )

        # 5. Melt and Clean (identical to Streamlit)
        chart_data = comparison_data.melt(
            id_vars=first_col,
            var_name="Period",
            value_name=metric
        ).dropna()

        if chart_data.empty:
            logging.debug(f"No valid data for YTD {metric} after filtering")
            return None

        # 6. Fiscal Year Sorting (matches Streamlit's Apr-Mar order)
        def fiscal_sort_key(period):
            month_map = {'Apr':1, 'May':2, 'Jun':3, 'Jul':4, 'Aug':5, 'Sep':6,
                        'Oct':7, 'Nov':8, 'Dec':9, 'Jan':10, 'Feb':11, 'Mar':12}
            match = re.search(r'\((Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec|Jan|Feb|Mar)', str(period))
            return month_map.get(match.group(1) if match else 99, 99)

        chart_data = chart_data.sort_values(by='Period', key=lambda x: x.map(fiscal_sort_key))

        # 7. Metrics Calculation (identical to Streamlit)
        metrics = {
            'total': float(chart_data[metric].sum()),
            'average': float(chart_data[metric].mean()),
            'min': float(chart_data[metric].min()),
            'max': float(chart_data[metric].max()),
            'period_count': len(chart_data['Period'].unique())
        }

        return {
            'data': make_jsonly_serializable(chart_data),
            'metrics': metrics,
            'x_col': 'Period',
            'y_col': metric,
            'color_override': '#FF8C00' if metric == "Act" else '#2E86AB'  # Matches Streamlit
        }

    except Exception as e:
        logging.error(f"YTD {metric} processing failed: {str(e)}", exc_info=True)
        return None
    
def create_ytd_visualizations(df, visual_type, selected_year=None):
    """Create all YTD visualizations (Budget, LY, Act, Gr, Ach)"""
    ytd_visualizations = {}
    metrics = ['Budget', 'LY', 'Act', 'Gr', 'Ach']
    
    for metric in metrics:
        result = process_ytd_comparison(df, metric, selected_year)
        if result:
            fig, config = create_plotly_chart(
                result['data'],
                result['x_col'],
                result['y_col'],
                visual_type,
                f'YTD {metric} Analysis',
                result['color_override']
            )
            
            ytd_visualizations[f'YTD {metric}'] = {
                'figure': json.loads(json.dumps(fig.to_dict(), cls=plotly.utils.PlotlyJSONEncoder)),
                'data': result['data'].to_dict(orient='records'),
                'metrics': result['metrics'],
                'x_col': result['x_col'],
                'y_col': result['y_col'],
                'color_override': result['color_override']
            }
    
    return ytd_visualizations

def create_budget_vs_actual_ppt(data, table_name, selected_filter=None):
    """Special handling for Budget vs Actual PPT generation"""
    try:
        # Create presentation
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Title and content
        
        # Set title
        title = f"Budget vs Actual - {table_name}"
        if selected_filter and selected_filter != "Select All":
            title += f" ({selected_filter})"
        slide.shapes.title.text = title
        
        # Prepare data - ensure we have both Budget and Actual traces
        if not isinstance(data, list) or len(data) < 2:
            raise ValueError("Budget vs Actual requires both Budget and Actual data")
            
        budget_data = next((d for d in data if d.get('name') == 'Budget'), None)
        actual_data = next((d for d in data if d.get('name') == 'Actual'), None)
        
        if not budget_data or not actual_data:
            raise ValueError("Missing Budget or Actual data in chart data")
            
        # Create DataFrame for plotting
        df = pd.DataFrame({
            'Month': budget_data['x'],
            'Budget': budget_data['y'],
            'Actual': actual_data['y']
        })
        
        # Create plot with matplotlib
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # Set positions and width for bars
        bar_width = 0.35
        x_pos = np.arange(len(df['Month']))
        
        # Plot bars
        budget_bars = ax.bar(
            x_pos - bar_width/2,
            df['Budget'],
            bar_width,
            label='Budget',
            color='#2E86AB'
        )
        
        actual_bars = ax.bar(
            x_pos + bar_width/2,
            df['Actual'],
            bar_width,
            label='Actual',
            color='#FF8C00'
        )
        
        # Customize plot
        ax.set_xlabel('Month')
        ax.set_ylabel('Value')
        ax.set_title(title)
        ax.set_xticks(x_pos)
        ax.set_xticklabels(df['Month'], rotation=45, ha='right')
        ax.legend()
        
        plt.tight_layout()
        
        # Save plot to buffer
        img_buffer = BytesIO()
        plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
        img_buffer.seek(0)
        plt.close()
        
        # Add image to slide
        slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
        
        # Save PPT to buffer
        ppt_buffer = BytesIO()
        ppt.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
        
    except Exception as e:
        logging.error(f"Error creating Budget vs Actual PPT: {str(e)}", exc_info=True)
        raise

