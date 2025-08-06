from flask import Blueprint, jsonify, request, Response, current_app, send_file
import utils.helpers as helpers
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
import plotly.graph_objects as go
import plotly.utils
from io import BytesIO
import re
import csv
import gc
import io
from datetime import datetime
import os
import json
import base64
import logging
import numpy as np
import pandas as pd
from werkzeug.utils import secure_filename
from services.data_processing import process_ytd_comparison, create_plotly_chart
from flask import current_app  # <-- ADD THIS
# Add these imports at the top with other imports
from utils.helpers import (
    rename_columns,
    extract_tables,
    find_table_end,
    column_filter,
    make_jsonly_serializable,
    convert_to_numeric,
    safe_sum,
    safe_mean,
    is_monthwise_sheet,
    extract_month_year
)
from utils import helpers
from utils.helpers import make_jsonly_serializable
import services.data_processing as data_processing

# Add this near the top of data_processing.py
BRANCH_EXCLUDE_TERMS = [
    'CHN Total', 'ERD SALES', 'North Total', 'WEST SALES', 'GROUP COMPANIES'
]

#from services.data_processing import handle_upload, handle_process_sheet
main_bp = Blueprint('main', __name__)

@main_bp.route('/test', methods=['GET'])
def test():
    return jsonify({'message': 'Backend is running successfully!', 'status': 'ok'})

# Routes
@main_bp.route('/test', methods=['GET'])
def test_connection():
    return jsonify({'message': 'Backend is running successfully!', 'status': 'ok'})

@main_bp.errorhandler(500)
def internal_error(error):
    logging.error(f"Internal server error: {str(error)}", exc_info=True)
    return jsonify({'error': 'Internal server error', 'message': str(error)}), 500

@main_bp.errorhandler(404)
def not_found(error):
    return jsonify({'error': 'Endpoint not found'}), 404

@main_bp.route('/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if not file.filename.endswith('.xlsx'):
            return jsonify({'error': 'Invalid file format. Please upload an Excel file (.xlsx)'}), 400

        file_buffer = BytesIO(file.read())
        
        try:
            xls = pd.ExcelFile(file_buffer, engine='openpyxl')
        except:
            file_buffer.seek(0)
            xls = pd.ExcelFile(file_buffer)
        
        sheet_names = xls.sheet_names
        
        file_buffer.seek(0)
        df_sheet = pd.read_excel(file_buffer, sheet_name=sheet_names[0], header=None, nrows=1000)
        
        # Determine sheet type for table ending logic
        sheet_name = sheet_names[0].lower()
        is_sales_analysis = 'sales analysis month-wise' in sheet_name
        is_branch_analysis = 'region wise analysis' in sheet_name
        is_product_analysis = 'product' in sheet_name or 'ts-pw' in sheet_name or 'ero-pw' in sheet_name

        # Alternative reading method for narrow sheets with long text
        if df_sheet.shape[1] < 10 and df_sheet.iloc[:, 0].astype(str).str.len().max() > 200:
            try:
                file_buffer.seek(0)
                df_sheet_alt = pd.read_excel(file_buffer, sheet_name=sheet_names[0], header=None, engine='openpyxl', nrows=1000)
                if df_sheet_alt.shape[1] > df_sheet.shape[1]:
                    df_sheet = df_sheet_alt
            except:
                pass
        
        # Process sheets with fewer than 20 columns (align with Streamlit logic)
        if df_sheet.shape[1] < 20:
            new_data = []
            months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
            metrics = ['Budget', 'LY', 'Act', 'Gr', 'Ach', 'YTD']
            year_pattern = r'\d{2,4}(?:[-–]\d{2,4})?'
            
            for idx, row in df_sheet.iterrows():
                if pd.notna(row.iloc[0]):
                    row_text = str(row.iloc[0]).strip()
                    if any(metric in row_text for metric in metrics) or re.search(r'SALES\s*(in\s*(MT|Value|Ton[n]?age))?', row_text, re.IGNORECASE):
                        patterns = []
                        patterns.append(r'SALES\s*in\s*(MT|Value|Ton[n]?age)', re.IGNORECASE)
                        for metric in metrics:
                            for month in months:
                                patterns.append(rf'{metric}[-–\s]*{month}[-–\s]*{year_pattern}', re.IGNORECASE)
                            patterns.append(rf'{metric}[-–\s]*YTD[-–\s]*{year_pattern}\s*\([^)]*\)', re.IGNORECASE)
                            patterns.append(rf'YTD[-–\s]*{year_pattern}\s*\([^)]*\)\s*{metric}', re.IGNORECASE)
                        
                        positions = []
                        for pattern in patterns:
                            for match in re.finditer(pattern, row_text):
                                positions.append((match.start(), match.group()))
                        positions.sort()
                        parts = [item[1].strip() for item in positions]
                        
                        # Fallback splitting if regex doesn't capture enough parts
                        if len(parts) < 5:
                            parts = [part.strip() for part in row_text.split() if part.strip()]
                        
                        new_data.append(parts)
                    else:
                        new_data.append(row_text.split())
                else:
                    new_data.append([])
            
            if new_data:
                max_cols = max(len(row) for row in new_data)
                for row in new_data:
                    while len(row) < max_cols:
                        row.append(None)
                df_sheet = pd.DataFrame(new_data)
        
        # Clean and prepare preview data
        df_sheet = df_sheet.dropna(how='all').reset_index(drop=True)
        
        # Find table end based on sheet type
        table_end_idx = len(df_sheet)
        for i in range(len(df_sheet)):
            row_text = ' '.join(str(cell) for cell in df_sheet.iloc[i].values if pd.notna(cell)).upper()
            if is_branch_analysis and 'GRAND TOTAL' in row_text:
                table_end_idx = i + 1  # Include Grand Total
                break
            elif (is_sales_analysis or is_product_analysis) and 'TOTAL SALES' in row_text:
                table_end_idx = i + 1  # Include Total Sales
                break
        
        # Truncate DataFrame to table end
        df_sheet = df_sheet.iloc[:table_end_idx].reset_index(drop=True)
        
        # Detect header row
        header_row_idx = None
        for i in range(min(5, len(df_sheet))):
            row_text = ' '.join(str(cell) for cell in df_sheet.iloc[i].values if pd.notna(cell))
            if re.search(r'\b(?:budget|ly|act|gr|ach)\b', row_text, re.IGNORECASE):
                header_row_idx = i
                break
        
        if header_row_idx is not None:
            header_row = df_sheet.iloc[header_row_idx]
            new_columns = [str(val).strip() if pd.notna(val) else f'Unnamed_{i}' 
                          for i, val in enumerate(header_row)]
            df_sheet.columns = new_columns
            df_sheet = df_sheet.iloc[header_row_idx + 1:].reset_index(drop=True)
            
            # Remove specific rows (e.g., 'REGIONS' for first sheet)
            if sheet_names.index(sheet_names[0]) == 0:
                df_sheet = df_sheet[~df_sheet[df_sheet.columns[0]].str.contains('REGIONS', case=False, na=False, regex=True)].reset_index(drop=True)
        
        # Serialize data for JSON response
        preview_data = make_jsonly_serializable(df_sheet)
        preview_columns = [
            {'field': str(col), 'headerName': str(col), 'width': 150}
            for col in preview_data.columns
        ]
        
        response = {
            'preview_data': preview_data.to_dict(orient='records'),
            'preview_columns': preview_columns,
            'sheet_names': sheet_names
        }
        
        return jsonify(response)
    
    except Exception as e:
        return jsonify({'error': f'Error processing file: {str(e)}'}), 500

@main_bp.route('/process-sheet', methods=['POST'])
def process_sheet():
    try:
        # Initial setup and validation
        data = request.get_json()
        file_data = data.get('file_data')
        sheet_name = data.get('sheet_name')
        table_choice = data.get('table_choice')
        
        if not file_data or not sheet_name:
            return jsonify({'error': 'Missing file data or sheet name'}), 400

        # Decode and read Excel file
        file_content = base64.b64decode(file_data)
        file_buffer = BytesIO(file_content)
        
        try:
            xls = pd.ExcelFile(file_buffer, engine='openpyxl')
        except Exception as e:
            logging.warning(f"Openpyxl failed, trying default engine: {str(e)}")
            file_buffer.seek(0)
            xls = pd.ExcelFile(file_buffer)
        
        if sheet_name not in xls.sheet_names:
            return jsonify({'error': f'Sheet "{sheet_name}" not found in file'}), 404

        # Read sheet data with initial processing
        df_sheet = pd.read_excel(xls, sheet_name=sheet_name, header=None, nrows=1000)
        original_shape = df_sheet.shape
        logging.info(f"Processing sheet: {sheet_name}, initial shape: {original_shape}")

        # Data restructuring for poorly formatted files
        if df_sheet.shape[1] < 20:
            new_data = []
            months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 
                     'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
            metrics = ['Budget', 'LY', 'Act', 'Gr', 'Ach', 'YTD']
            year_pattern = r'\d{2,4}(?:[-–]\d{2,4})?'
            
            for idx, row in df_sheet.iterrows():
                if pd.notna(row.iloc[0]):
                    row_text = str(row.iloc[0]).strip()
                    
                    if any(metric in row_text for metric in metrics) or re.search(r'SALES\s*(in\s*(MT|Value|Ton[n]?age))?', row_text, re.IGNORECASE):
                        patterns = []
                        patterns.append(r'SALES\s*in\s*(MT|Value|Ton[n]?age)', re.IGNORECASE)
                        for metric in metrics:
                            for month in months:
                                patterns.append(
                                    rf'{metric}[-–\s]*{month}[-–\s]*{year_pattern}', 
                                    re.IGNORECASE
                                )
                            patterns.append(
                                rf'{metric}[-–\s]*YTD[-–\s]*{year_pattern}\s*\([^)]*\)', 
                                re.IGNORECASE
                            )
                            patterns.append(
                                rf'YTD[-–\s]*\d{2,4}\s*\([^)]*\)\s*{metric}', 
                                re.IGNORECASE
                            )
                        
                        positions = []
                        for pattern in patterns:
                            for match in re.finditer(pattern, row_text):
                                positions.append((match.start(), match.group()))
                        positions.sort()
                        
                        parts = [item[1].strip() for item in positions]
                        
                        if len(parts) < 5:
                            parts = [part.strip() for part in row_text.split() if part.strip()]
                        
                        new_data.append(parts)
                    else:
                        new_data.append(row_text.split())
                else:
                    new_data.append([])
             
            if new_data:
                max_cols = max(len(row) for row in new_data)
                for row in new_data:
                    while len(row) < max_cols:
                        row.append(None)
                df_sheet = pd.DataFrame(new_data)

        # Determine sheet type
        sheet_index = xls.sheet_names.index(sheet_name)
        is_first_sheet = sheet_index == 0
        is_branch_analysis = 'region wise analysis' in sheet_name.lower()
        is_product_analysis = ('product' in sheet_name.lower() or 
                             'ts-pw' in sheet_name.lower() or 
                             'ero-pw' in sheet_name.lower())

        # Enhanced table processing function
        def process_table(start_idx, end_idx, is_branch=False, is_product=False, is_table1=False):
            if start_idx is None:
                logging.warning("No start index for table")
                return None
            
            table = df_sheet.iloc[start_idx:end_idx].dropna(how='all')
            if table.empty:
                logging.warning("Extracted table is empty")
                return None
            
            # Enhanced header detection
            header_row_idx = None
            for i in range(min(10, len(table))):
                row_text = ' '.join(str(cell) for cell in table.iloc[i].values if pd.notna(cell)).upper()
                metric_count = sum(1 for m in ['BUDGET', 'LY', 'ACT', 'GR', 'ACH'] if m in row_text)
                month_count = sum(1 for m in ['JAN', 'FEB', 'MAR', 'APR', 'MAY', 'JUN', 
                                            'JUL', 'AUG', 'SEP', 'OCT', 'NOV', 'DEC'] if m in row_text)
                
                if metric_count >= 2 or month_count >= 3:
                    header_row_idx = i
                    break
            
            if header_row_idx is None:
                return None
                
            table.columns = [str(val).strip() if pd.notna(val) else f'Unnamed_{i}' 
                           for i, val in enumerate(table.iloc[header_row_idx])]
            table = table.iloc[header_row_idx+1:].reset_index(drop=True)
            
            # Clean and standardize column names
            table.columns = rename_columns(table.columns)
            table = table.loc[:, ~table.columns.duplicated()]
            table = table.replace([np.nan, None], 0.0)
            
            # Filtering logic
            if is_branch and table.columns[0]:
                exclude_terms = current_app.config.get('BRANCH_EXCLUDE_TERMS', [])
                if exclude_terms:
                    pattern = '|'.join(exclude_terms)
                    table = table[~table[table.columns[0]].astype(str).str.upper().str.contains(pattern, na=False)]
            
            if is_table1 and table.columns[0]:
                table = table[~(
                    table[table.columns[0]].astype(str).str.contains('SALES', case=False, na=False) &
                    (table.iloc[:, 1:].eq(0.0) | table.iloc[:, 1:].isna()).all(axis=1)
                )]
            # Extra filter for product sheets to remove summary rows like:
# "TS-PW Monthly Budget and Actual Values North [FY 25-26]"
# or "ERO-PW Monthly Budget and Actual Values West [FY 25-26]"
            if is_product and table.columns[0]:
                first_col_series = table[table.columns[0]].astype(str)
                # Normalize multiple spaces to single space before matching
                first_col_series = first_col_series.str.replace(r'\s+', ' ', regex=True)
                table = table[
                    ~first_col_series.str.contains(
                        r'Monthly\s+Budget\s+and\s+Actual',
                        case=False,
                        na=False,
                        regex=True
                    )
                ]

            return make_jsonly_serializable(table)

        # First sheet processing
        if is_first_sheet:
            # Enhanced table detection
            table1_start = table2_start = None
            for i in range(len(df_sheet)):
                row_text = ' '.join(str(cell) for cell in df_sheet.iloc[i].values if pd.notna(cell)).upper()
                
                if table1_start is None and any(t in row_text for t in ['SALES IN MT', 'SALES(MT)', 'SALES IN TONNAGE', 'SALES IN TONAGE']):
                    table1_start = i
                
                if table2_start is None and any(t in row_text for t in ['SALES IN VALUE', 'SALES(VALUE)']):
                    table2_start = i
            
            table_options = []
            if table1_start is not None:
                table_options.append("Table 1: SALES IN MT")  # Always use "MT" here
            if table2_start is not None:
                table_options.append("Table 2: SALES IN VALUE")
            
            if not table_options:
                return jsonify({
                    'tables': [],
                    'data': make_jsonly_serializable(df_sheet).to_dict(orient='records'),
                    'csv_filename': 'raw_data.csv',
                    'is_first_sheet': True,
                    'columns': list(df_sheet.columns)
                })

            if not table_choice:
                table_choice = table_options[0]
            elif table_choice not in table_options:
                # Handle case where frontend expects "MT" but we found "TONAGE"
                if table_choice == "Table 1: SALES IN MT" and "Table 1: SALES IN TONAGE" in table_options:
                    table_choice = "Table 1: SALES IN TONAGE"
                else:
                    return jsonify({
                        'tables': table_options,
                        'error': f"Selected table '{table_choice}' is invalid. Available options: {table_options}",
                        'is_first_sheet': True
                    }), 400
            
            # Process selected table
            # Process selected table
            if "Table 1" in table_choice and table1_start is not None:  # Handles both MT and TONAGE cases
                table1_end = table2_start if table2_start is not None else len(df_sheet)
                table_data = process_table(
                    table1_start,
                    table1_end,
                    is_branch=is_branch_analysis,
                    is_product=is_product_analysis,
                    is_table1=True
                )
                csv_filename = 'sales_in_mt.csv'

            elif table_choice == "Table 2: SALES IN VALUE" and table2_start is not None:
                table2_end = len(df_sheet)
                table_data = process_table(
                    table2_start,
                    table2_end,
                    is_branch=is_branch_analysis,
                    is_product=is_product_analysis
                )
                csv_filename = 'sales_in_value.csv'

            else:
                return jsonify({
                    'tables': table_options,
                    'error': f"Selected table '{table_choice}' not found or invalid.",
                    'is_first_sheet': True
                }), 400

            
            if table_data is None or table_data.empty:
                return jsonify({
                    'tables': table_options,
                    'error': f"No valid data in selected table '{table_choice}'.",
                    'debug_info': {
                        'table_start': table1_start if "Table 1" in table_choice else table2_start,
                        'sheet_name': sheet_name,
                        'sample_data': df_sheet.head(5).to_dict(orient='records')
                    },
                    'is_first_sheet': True
                }), 400
            
            return jsonify({
                'tables': table_options,
                'data': table_data.to_dict(orient='records'),
                'csv_filename': csv_filename,
                'is_first_sheet': True,
                'columns': list(table_data.columns)
            })
        
        # Non-first sheet processing - handle both "MT" and "Tonage" for branch/product analysis
        if is_branch_analysis or is_product_analysis:
            # Try both "MT" and "Tonage" variants
            table1_header = "Sales in MT"
            table2_header = "Sales in Value"
            
            # First try with "MT"
            idx1, idx2 = extract_tables(df_sheet, table1_header, table2_header)
            
            # If "MT" not found, try "Tonage"
            if idx1 is None:
                table1_header_alt = "Sales in Tonage"
                idx1, idx2 = extract_tables(df_sheet, table1_header_alt, table2_header)
                if idx1 is None:
                    return jsonify({'error': f"Could not locate sales table header in sheet."}), 400
        else:
            # For other sheets, only use "MT"
            table1_header = "Sales in MT"
            table2_header = "Sales in Value"
            idx1, idx2 = extract_tables(df_sheet, table1_header, table2_header)
            if idx1 is None:
                return jsonify({'error': f"Could not locate sales table header in sheet."}), 400
        
        # Process the selected table
        table1_end = idx2 if idx2 is not None else find_table_end(df_sheet, idx1 + 1, is_branch_analysis, is_product_analysis)
        table1 = process_table(idx1, table1_end, is_branch_analysis, is_product_analysis, is_table1=True)

        table2 = None
        if idx2 is not None:
            table2_end = find_table_end(df_sheet, idx2 + 1, is_branch_analysis, is_product_analysis)
            table2 = process_table(idx2, table2_end, is_branch_analysis, is_product_analysis)

        # Set table options - always use "MT" in the label even if we found "Tonage"
        table_options = ["Table 1: SALES IN MT"]
        if idx2 is not None:
            table_options.append("Table 2: SALES IN VALUE")

        # Validate table choice - handle both "MT" and "Tonage" cases
        if not table_choice:
            table_choice = table_options[0]
        elif table_choice not in table_options:
            # Special handling for frontend sending "MT" when we have "Tonage"
            if table_choice == "Table 1: SALES IN MT":
                table_choice = "Table 1: SALES IN MT"  # Still use this label
            else:
                return jsonify({
                    'tables': table_options,
                    'error': f"Selected table '{table_choice}' is invalid. Available options: {table_options}",
                    'is_first_sheet': False
                }), 400
        
        # Select the appropriate table data
        table_df = table1 if "Table 1" in table_choice else (table2 if idx2 is not None else table1)
        
        if table_df is None or table_df.empty:
            return jsonify({
                'tables': table_options,
                'error': f"Selected table '{table_choice}' is empty or invalid.",
                'is_first_sheet': False
            }), 400
        
        # Prepare response data
        first_col = table_df.columns[0]
        branch_list = []
        product_list = []
        
        if is_branch_analysis:
            valid_rows = table_df[table_df[first_col].notna()]
            branch_list = sorted(
                valid_rows[first_col].astype(str).str.strip().unique(),
                key=lambda x: x.lower()
            )
            exclude_terms = current_app.config.get('BRANCH_EXCLUDE_TERMS', [])
            branch_list = [b for b in branch_list if not any(ex in b.upper() for ex in exclude_terms)]
        
        elif is_product_analysis:
            valid_rows = table_df[table_df[first_col].notna()]
            product_list = sorted(
                valid_rows[first_col].astype(str).str.strip().unique(),
                key=lambda x: x.lower()
            )
        
        all_columns = ' '.join(str(col) for col in table_df.columns)
        months = sorted(
            set(re.findall(
                r'\b(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)\b',
                all_columns,
                re.IGNORECASE
            )),
            key=lambda x: ['jan','feb','mar','apr','may','jun',
                          'jul','aug','sep','oct','nov','dec'].index(x.lower())
        )
        
        years = sorted(
            set(re.findall(r'[-–](\d{2,4})\b', all_columns)),
            reverse=True
        )

        response_data = {
            'tables': table_options,
            'data': table_df.to_dict(orient='records'),
            'csv_filename': 'filtered_data.csv',
            'is_first_sheet': False,
            'is_branch_analysis': is_branch_analysis,
            'is_product_analysis': is_product_analysis,
            'months': ['Select All'] + months,
            'years': ['Select All'] + years,
            'columns': list(table_df.columns)
        }
        
        if is_branch_analysis:
            response_data['branches'] = ['Select All'] + branch_list
            response_data['products'] = []
        elif is_product_analysis:
            response_data['products'] = ['Select All'] + product_list
            response_data['branches'] = []
        else:
            response_data['branches'] = []
            response_data['products'] = []

        return jsonify(response_data)

    except Exception as e:
        logging.error(f"Error in process_sheet: {str(e)}", exc_info=True)
        return jsonify({
            'error': 'Internal server error',
            'details': str(e)
        }), 500
    
@main_bp.route('/process-ytd', methods=['POST'])
def process_ytd():
    try:
        data = request.get_json()
        if not isinstance(data.get('table_data'), list):
            return jsonify({'error': 'table_data must be a list of records'}), 400
        
        table_data = pd.DataFrame(data['table_data'])
        metric = data.get('metric')
        selected_year = data.get('selected_year')
        
        if not metric:
            return jsonify({'error': 'Metric is required'}), 400
        
        # Process YTD data
        result = process_ytd_comparison(table_data, metric, selected_year)
        if not result:
            return jsonify({'error': f'No valid YTD {metric} data found'}), 404
        
        # Create visualization
        fig, config = create_plotly_chart(
            result['data'],
            result['x_col'],
            result['y_col'],
            data.get('visual_type', 'bar'),
            f'YTD {metric} Analysis',
            result['color_override']
        )
        
        return jsonify({
            'figure': json.loads(json.dumps(fig.to_dict(), cls=plotly.utils.PlotlyJSONEncoder)),
            'data': result['data'].to_dict(orient='records'),
            'metrics': result['metrics'],
            'status': 'success'
        })
        
    except Exception as e:
        logging.error(f"Error in YTD processing: {str(e)}", exc_info=True)
        return jsonify({
            'error': 'Internal server error',
            'details': str(e)
        }), 500    
        
@main_bp.route('/visualizations', methods=['POST'])
def get_visualizations():
    try:
        # Validate and parse request data
        data = request.get_json()
        if not isinstance(data.get('table_data'), list):
            logging.error("Invalid table_data format")
            return jsonify({'error': 'table_data must be a list of records'}), 400
        
        selected_sheet = data.get('selected_sheet', '').lower()
        
        # Skip visualization generation for Sales Analysis Month-wise
        if is_monthwise_sheet(selected_sheet):
            return jsonify({
                'status': 'skipped',
                'message': 'Visualizations blocked for Sales Month-wise sheet',
                'display_data': make_jsonly_serializable(pd.DataFrame(data['table_data'])).to_dict(orient='records')
            }), 200

        table_data = pd.DataFrame(data['table_data'])
        table_name = data.get('table_name', '')
        visual_type = data.get('visual_type', 'bar').lower().replace(" chart", "")
        selected_month = data.get('selected_month')
        selected_year = data.get('selected_year')
        selected_branch = data.get('selected_branch')
        selected_product = data.get('selected_product')
        is_branch_analysis = data.get('is_branch_analysis', False)
        is_product_analysis = data.get('is_product_analysis', False)
        
        # Validate required fields
        if not all([table_name, selected_sheet]):
            logging.error("Missing required fields")
            return jsonify({'error': 'Missing required fields: table_name and selected_sheet are required'}), 400

        # Check if this is a Region Wise Analysis sheet
        is_region_wise_analysis = 'region wise analysis' in selected_sheet.lower()

        # Initialize response structure
        visualizations = {}
        display_data = None
        selected_filter = None

        # Process the data frame
        formatted_df = make_jsonly_serializable(table_data)
        
        # Convert numeric columns
        numeric_cols = []
        for col in formatted_df.columns:
            if col == formatted_df.columns[0]:  # Skip first column (names)
                continue
            formatted_df[col] = convert_to_numeric(formatted_df[col])
            if pd.api.types.is_numeric_dtype(formatted_df[col]):
                numeric_cols.append(col)
                formatted_df[col] = formatted_df[col].round(2)

        # Set selected filter for title
        if is_branch_analysis and selected_branch and selected_branch != "Select All":
            selected_filter = selected_branch
        elif selected_product and selected_product != "Select All":
            selected_filter = selected_product
        elif selected_month and selected_month != "Select All":
            selected_filter = selected_month
        elif selected_year and selected_year != "Select All":
            selected_filter = selected_year

        # Process Budget vs Actual
        budget_cols = [col for col in formatted_df.columns if 'budget' in str(col).lower()]
        actual_cols = [col for col in formatted_df.columns if 'act' in str(col).lower() and 'budget' not in str(col).lower()]
        
        if budget_cols and actual_cols:
            # Find matching months between budget and actual
            month_data = []
            months_found = set()
            
            for col in budget_cols + actual_cols:
                month = extract_month_year(col)
                if month not in months_found:
                    months_found.add(month)
                    is_budget = 'budget' in str(col).lower()
                    
                    try:
                        value = safe_sum(formatted_df[col]) if col in formatted_df.columns else 0
                        month_data.append({
                            'Month': month,
                            'Value': float(value),
                            'Metric': 'Budget' if is_budget else 'Actual'
                        })
                    except:
                        continue
            
            if month_data:
                chart_data = pd.DataFrame(month_data)
                chart_data = make_jsonly_serializable(chart_data)
                
                fig, config = create_plotly_chart(
                    chart_data,
                    'Month',
                    'Value',
                    visual_type,
                    'Budget vs Actual Analysis'
                )
                
                visualizations['Budget vs Actual'] = {
                    'figure': json.loads(json.dumps(fig.to_dict(), cls=plotly.utils.PlotlyJSONEncoder)),
                    'data': chart_data.to_dict(orient='records'),
                    'x_col': 'Month',
                    'y_col': 'Value',
                    'ppt_type': visual_type,
                    'color_override': None
                }

        # Process YTD Comparisons (Budget, LY, Act, Gr, Ach)
        ytd_metrics = ['Budget', 'LY', 'Act', 'Gr', 'Ach']
        for metric in ytd_metrics:
            ytd_cols = []
            pattern = rf'(?:{metric}[-\s]*)?ytd[-\s]*\d{{2,4}}[-\s]*\d{{2,4}}\s*\([^)]*\)'
            
            for col in formatted_df.columns:
                col_str = str(col).lower().replace(",", "").replace("–", "-")
                if re.search(pattern, col_str, re.IGNORECASE) and column_filter(col, None, selected_year):
                    ytd_cols.append(col)
            
            if not ytd_cols:
                logging.debug(f"No YTD columns found for {metric}")
                continue
            
            # Prepare YTD data
            first_col = formatted_df.columns[0]
            comparison_data = formatted_df[[first_col] + ytd_cols].copy()
            
            # Clean column names
            clean_labels = []
            for col in ytd_cols:
                col_str = str(col)
                ytd_match = re.search(r'ytd[-\s]*(\d{2,4})[-\s]*(\d{2,4})\s*\((.*?)\)', col_str, re.IGNORECASE)
                if ytd_match:
                    start_year, end_year, month_range = ytd_match.groups()
                    clean_labels.append(f"YTD {start_year}-{end_year} ({month_range})")
                else:
                    clean_labels.append(col_str)
            
            comparison_data.columns = [first_col] + clean_labels
            
            # Convert to numeric
            for col in clean_labels:
                comparison_data[col] = pd.to_numeric(
                    comparison_data[col].astype(str).str.replace(',', ''),
                    errors='coerce'
                )
            
            # Melt for visualization
            chart_data = comparison_data.melt(
                id_vars=first_col,
                var_name="Period",
                value_name=metric
            ).dropna()
            
            if chart_data.empty:
                logging.debug(f"No valid data for YTD {metric} comparison")
                continue
            
            # Calculate metrics
            total = safe_sum(chart_data[metric])
            average = safe_mean(chart_data[metric])
            min_val = float(chart_data[metric].min()) if len(chart_data) > 0 else 0
            max_val = float(chart_data[metric].max()) if len(chart_data) > 0 else 0
            
            # Create visualization
            fig, config = create_plotly_chart(
                chart_data,
                'Period',
                metric,
                visual_type,
                f'YTD {metric} Analysis',
                '#FF8C00' if metric == "Act" else '#2E86AB'
            )
            
            visualizations[f'YTD {metric}'] = {
                'figure': json.loads(json.dumps(fig.to_dict(), cls=plotly.utils.PlotlyJSONEncoder)),
                'data': chart_data.to_dict(orient='records'),
                'x_col': 'Period',
                'y_col': metric,
                'ppt_type': visual_type,
                'color_override': '#FF8C00' if metric == "Act" else '#2E86AB',
                'metrics': {
                    'total': total,
                    'average': average,
                    'min': min_val,
                    'max': max_val
                }
            }

        # Process regular metrics (Act, Gr, Ach)
        metrics = ['Act', 'Gr', 'Ach']
        for metric in metrics:
            # Find the appropriate column
            metric_col = None
            for col in formatted_df.columns:
                col_str = str(col).lower().replace(",", "").replace("–", "-")
                if metric.lower() in col_str and 'ytd' not in col_str:
                    if column_filter(col, selected_month, selected_year):
                        metric_col = col
                        break
            
            if not metric_col:
                logging.warning(f"No {metric} column found")
                continue
                
            # Prepare data
            first_col = formatted_df.columns[0]
            chart_data = formatted_df[[first_col, metric_col]].copy()
            chart_data.columns = [first_col, metric]
            
            # Convert to numeric
            chart_data[metric] = pd.to_numeric(
                chart_data[metric].astype(str).str.replace(',', ''),
                errors='coerce'
            )
            chart_data = chart_data.dropna()
            
            if chart_data.empty:
                continue
                
            # Create visualization
            fig, config = create_plotly_chart(
                chart_data,
                first_col,
                metric,
                visual_type,
                f"{metric} Analysis",
                '#FF8C00' if metric == "Act" else '#2E86AB'
            )
            
            visualizations[metric] = {
                'figure': json.loads(json.dumps(fig.to_dict(), cls=plotly.utils.PlotlyJSONEncoder)),
                'data': chart_data.to_dict(orient='records'),
                'x_col': first_col,
                'y_col': metric,
                'ppt_type': visual_type,
                'color_override': '#FF8C00' if metric == "Act" else '#2E86AB'
            }

        # Process Branch Performance
        if is_region_wise_analysis:
            # Find performance column
            performance_col = None
            for col in formatted_df.columns:
                col_str = str(col).lower().replace(",", "").replace("–", "-")
                if ('act' in col_str or 'performance' in col_str) and 'budget' not in col_str:
                    performance_col = col
                    break
            if not performance_col:
                for col in formatted_df.columns[1:]:
                    try:
                        if pd.api.types.is_numeric_dtype(formatted_df[col]):
                            performance_col = col
                            break
                    except:
                        continue    
            
            if performance_col:
                # Filter out unwanted rows
                first_col = formatted_df.columns[0]
                branch_data = formatted_df[
                    ~formatted_df[first_col].str.contains('|'.join(BRANCH_EXCLUDE_TERMS), case=False, na=False)
                ].copy()
                
                if selected_branch and selected_branch != "Select All":
                    branch_data = branch_data[branch_data[first_col].astype(str) == selected_branch]
                
                # Convert to numeric
                branch_data[performance_col] = pd.to_numeric(
                    branch_data[performance_col].astype(str).str.replace(',', ''),
                    errors='coerce'
                )
                branch_data = branch_data.dropna()
                
                if not branch_data.empty:
                    # Prepare data
                    chart_data = branch_data[[first_col, performance_col]].copy()
                    chart_data.columns = ['Branch', 'Performance']
                    chart_data = chart_data.sort_values('Performance', ascending=False)
                    chart_data = make_jsonly_serializable(chart_data)
                    
                    # Calculate metrics
                    total_performance = float(safe_sum(chart_data['Performance']))
                    avg_performance = float(safe_mean(chart_data['Performance']))
                    top_performer = chart_data.iloc[0].to_dict() if len(chart_data) > 0 else {'Branch': 'N/A', 'Performance': 0}
                    top_5 = chart_data.head(5).to_dict(orient='records')
                    bottom_5 = chart_data.tail(5).to_dict(orient='records')
                    
                    # Create visualization
                    fig, config = create_plotly_chart(
                        chart_data,
                        'Branch',
                        'Performance',
                        visual_type,
                        'Branch Performance Analysis'
                    )
                    
                    visualizations['Branch Performance'] = {
                        'figure': json.loads(json.dumps(fig.to_dict(), cls=plotly.utils.PlotlyJSONEncoder)),
                        'data': chart_data.to_dict(orient='records'),
                        'x_col': 'Branch',
                        'y_col': 'Performance',
                        'ppt_type': visual_type,
                        'color_override': None,
                        'metrics': {
                            'topPerformer': {
                                'name': str(top_performer['Branch']),
                                'value': float(top_performer['Performance'])
                            },
                            'totalPerformance': total_performance,
                            'avgPerformance': avg_performance,
                            'top5': [{
                                'name': str(item['Branch']),
                                'value': float(item['Performance'])
                            } for item in top_5],
                            'bottom5': [{
                                'name': str(item['Branch']),
                                'value': float(item['Performance'])
                            } for item in bottom_5]
                        }
                    }
        # Process Branch Monthwise ONLY for Region Wise Analysis
            month_cols = []
            months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
            
            for col in formatted_df.columns:
                col_str = str(col)
                if any(month in col_str for month in months):
                    if column_filter(col, None, None):  # No month/year filter for this
                        month_cols.append(col)
            
            if month_cols:
                # Prepare data
                first_col = formatted_df.columns[0]
                month_data = []
                
                for col in month_cols:
                    month = extract_month_year(col)
                    try:
                        value = safe_sum(formatted_df[col]) if col in formatted_df.columns else 0
                        month_data.append({
                            'Month': month,
                            'Value': float(value)
                        })
                    except:
                        continue
                
                if month_data:
                    chart_data = pd.DataFrame(month_data)
                    chart_data = make_jsonly_serializable(chart_data)
                    
                    # Create visualization
                    fig, config = create_plotly_chart(
                        chart_data,
                        'Month',
                        'Value',
                        visual_type,
                        'Branch Monthwise Analysis'
                    )
                    
                    # Calculate metrics
                    metrics = {
                        'bestMonth': max(month_data, key=lambda x: x['Value']) if month_data else None,
                        'total': sum(d['Value'] for d in month_data),
                        'average': sum(d['Value'] for d in month_data) / len(month_data) if month_data else 0
                    }
                    
                    visualizations['Branch Monthwise'] = {
                        'figure': json.loads(json.dumps(fig.to_dict(), cls=plotly.utils.PlotlyJSONEncoder)),
                        'data': chart_data.to_dict(orient='records'),
                        'x_col': 'Month',
                        'y_col': 'Value',
                        'ppt_type': visual_type,
                        'color_override': None,
                        'metrics': metrics
                    }            

        # Process Product Performance
        if is_product_analysis and len(formatted_df) > 0:
            # Find performance column
            performance_col = None
            for col in formatted_df.columns:
                col_str = str(col).lower().replace(",", "").replace("–", "-")
                if (('act' in col_str or 'sales' in col_str) and 'budget' not in col_str and 
                ('mt' in col_str or 'tonnage' in col_str or 'tonage' in col_str)):
                    performance_col = col
                    break
            if not performance_col:
        # Fallback to any numeric column
                for col in formatted_df.columns[1:]:
                    try:
                        if pd.api.types.is_numeric_dtype(formatted_df[col]):
                             performance_col = col
                             break
                    except:
                        continue    
            
            if performance_col:
                # Filter out unwanted rows
                first_col = formatted_df.columns[0]
                product_data = formatted_df[
                    ~formatted_df[first_col].str.contains('TOTAL|GRAND TOTAL|OVERALL TOTAL', case=False, na=False)
                ].copy()
                
                if selected_product and selected_product != "Select All":
                    product_data = product_data[product_data[first_col].astype(str) == selected_product]
                
                # Convert to numeric
                product_data[performance_col] = pd.to_numeric(
                    product_data[performance_col].astype(str).str.replace(',', ''),
                    errors='coerce'
                )
                product_data = product_data.dropna()
                
                if not product_data.empty:
                    # Prepare data
                    chart_data = product_data[[first_col, performance_col]].copy()
                    chart_data.columns = ['Product', 'Performance']
                    chart_data = chart_data.sort_values('Performance', ascending=False)
                    chart_data = make_jsonly_serializable(chart_data)
                    
                    # Calculate metrics
                    total_performance = float(safe_sum(chart_data['Performance']))
                    avg_performance = float(safe_mean(chart_data['Performance']))
                    top_performer = chart_data.iloc[0].to_dict() if len(chart_data) > 0 else {'Product': 'N/A', 'Performance': 0}
                    top_5 = chart_data.head(5).to_dict(orient='records')
                    bottom_5 = chart_data.tail(5).to_dict(orient='records')
                    
                    # Create visualization
                    fig, config = create_plotly_chart(
                        chart_data,
                        'Product',
                        'Performance',
                        visual_type,
                        'Product Performance Analysis'
                    )
                    
                    visualizations['Product Performance'] = {
                        'figure': json.loads(json.dumps(fig.to_dict(), cls=plotly.utils.PlotlyJSONEncoder)),
                        'data': chart_data.to_dict(orient='records'),
                        'x_col': 'Product',
                        'y_col': 'Performance',
                        'ppt_type': visual_type,
                        'color_override': None,
                        'metrics': {
                            'topPerformer': {
                                'name': str(top_performer['Product']),
                                'value': float(top_performer['Performance'])
                            },
                            'totalPerformance': total_performance,
                            'avgPerformance': avg_performance,
                            'top5': [{
                                'name': str(item['Product']),
                                'value': float(item['Performance'])
                            } for item in top_5],
                            'bottom5': [{
                                'name': str(item['Product']),
                                'value': float(item['Performance'])
                            } for item in bottom_5]
                        }
                    }

        # Process Month-wise Data (Branch and Product)
        if (is_branch_analysis or is_product_analysis) and len(formatted_df) > 0:
            # Find month-wise columns
            month_cols = []
            months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
            
            for col in formatted_df.columns:
                col_str = str(col)
                if any(month in col_str for month in months):
                    if column_filter(col, selected_month, selected_year):
                        month_cols.append(col)
            
            if month_cols:
                # Prepare data
                first_col = formatted_df.columns[0]
                month_data = []
                
                for col in month_cols:
                    month = extract_month_year(col)
                    try:
                        value = safe_sum(formatted_df[col])
                        month_data.append({
                            'Month': month,
                            'Value': float(value)
                        })
                    except:
                        continue
                
                if month_data:
                    chart_data = pd.DataFrame(month_data)
                    chart_data = make_jsonly_serializable(chart_data)
                    
                    # Calculate metrics
                    total_performance = float(safe_sum(chart_data['Value']))
                    avg_monthly = float(safe_mean(chart_data['Value']))
                    
                    if len(chart_data) > 0:
                        max_idx = chart_data['Value'].idxmax()
                        best_month = chart_data.loc[max_idx].to_dict()
                    else:
                        best_month = {'Month': 'N/A', 'Value': 0}
                    
                    metrics = {
                        'bestMonth': {
                            'month': str(best_month['Month']),
                            'value': float(best_month['Value'])
                        },
                        'totalPerformance': total_performance,
                        'avgMonthly': avg_monthly
                    }
                    
                    # Create visualization
                    label = 'Branch Monthwise' if is_branch_analysis else 'Product Monthwise'
                    fig, config = create_plotly_chart(
                        chart_data,
                        'Month',
                        'Value',
                        visual_type,
                        f'{label} Analysis'
                    )
                    
                    visualizations[label] = {
                        'figure': json.loads(json.dumps(fig.to_dict(), cls=plotly.utils.PlotlyJSONEncoder)),
                        'data': chart_data.to_dict(orient='records'),
                        'x_col': 'Month',
                        'y_col': 'Value',
                        'ppt_type': visual_type,
                        'color_override': None,
                        'metrics': metrics
                    }

        # Prepare display data
        display_df = formatted_df.copy()
        if is_branch_analysis or is_product_analysis:
            first_col = display_df.columns[0]
            visual_cols = [col for col in display_df.columns if column_filter(col, selected_month, selected_year)]
            display_df = display_df[[first_col] + visual_cols] if visual_cols else display_df[[first_col]]
        
        display_data = make_jsonly_serializable(display_df).to_dict(orient='records')

        return jsonify({
            'visualizations': visualizations,
            'display_data': display_data,
            'selected_filter': selected_filter,
            'status': 'success'
        })

    except Exception as e:
        logging.error(f"Error in visualizations endpoint: {str(e)}", exc_info=True)
        return jsonify({
            'error': 'Internal server error',
            'details': str(e),
            'status': 'error'
        }), 500

@main_bp.route('/download-csv', methods=['POST'])
def download_csv():
    try:
        if not request.is_json:
            return jsonify({'error': 'Request must be JSON'}), 400

        data = request.get_json()
        
        # Get the raw table data (exactly as shown in preview)
        table_data = data.get('table_data')
        columns = data.get('columns', [])
        
        if not table_data or not columns:
            return jsonify({'error': 'Missing table data or columns'}), 400

        # Create DataFrame while preserving original formatting
        df = pd.DataFrame(table_data)
        
        # Ensure columns are in correct order
        if len(columns) == len(df.columns):
            df.columns = columns
        
        # Create CSV with proper formatting
        output = io.StringIO()
        
        # Write header
        output.write(','.join(f'"{col}"' for col in df.columns) + '\n')
        
        # Write rows with proper formatting
        for _, row in df.iterrows():
            csv_row = []
            for col in df.columns:
                value = row[col]
                if pd.isna(value) or value in ['', None]:
                    csv_row.append('')
                else:
                    # Preserve numbers with commas and strings with quotes
                    if isinstance(value, (int, float)) and not pd.isna(value):
                        csv_row.append(f'{value:,.2f}'.rstrip('0').rstrip('.') if '.' in f'{value}' else f'{value:,}')
                    else:
                        csv_row.append(f'"{str(value)}"')
            output.write(','.join(csv_row) + '\n')
        
        csv_content = output.getvalue()
        output.close()
        
        # Create response
        response = Response(
            csv_content,
            mimetype='text/csv',
            headers={
                'Content-Disposition': f'attachment; filename={secure_filename(data.get("filename", "export.csv"))}'
            }
        )
        
        return response

    except Exception as e:
        return jsonify({'error': f'CSV generation failed: {str(e)}'}), 500
    
@main_bp.route('/download-ppt', methods=['POST'])
def download_ppt():
    try:
        current_app.logger.info("PPT download request received")
        
        if not request.is_json:
            current_app.logger.error("Request must be JSON")
            return jsonify({'error': 'Request must be JSON'}), 400
            
        data = request.get_json()
        current_app.logger.info(f"Request data keys: {data.keys()}")
        
        # Special handling for Budget vs Actual charts
        if data.get('is_budget_vs_actual'):
            current_app.logger.info("Processing Budget vs Actual chart")
            
            # Validate required data
            if not all(key in data for key in ['vis_label', 'table_name', 'chart_data']):
                return jsonify({'error': 'Missing required data for Budget vs Actual'}), 400

            # Create presentation
            ppt = Presentation()
            slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Title and content
            
            # Set title
            title = f"{data['vis_label']} - {data['table_name']}"
            if data.get('selected_filter') and data['selected_filter'] != "Select All":
                title += f" ({data['selected_filter']})"
            slide.shapes.title.text = title
            
            # Handle pie chart case
            if data.get('is_pie_chart'):
                current_app.logger.info("Creating Budget vs Actual pie chart")
                
                # Extract data
                budget_data = next((d for d in data['chart_data'] if d.get('name') == 'Budget'), None)
                actual_data = next((d for d in data['chart_data'] if d.get('name') == 'Actual'), None)
                
                if not budget_data or not actual_data:
                    raise ValueError("Missing Budget or Actual data for pie chart")
                
                # Create pie chart
                plt.figure(figsize=(10, 6))
                sizes = [budget_data['y'], actual_data['y']]
                labels = ['Budget', 'Actual']
                
                # Ensure positive values
                if sizes[0] <= 0 or sizes[1] <= 0:
                    raise ValueError("Pie chart requires positive values for Budget and Actual")
                
                plt.pie(
                    sizes,
                    labels=labels,
                    autopct='%1.1f%%',
                    startangle=90,
                    colors=['#2E86AB', '#FF8C00'],
                    textprops={'fontsize': 12}
                )
                plt.axis('equal')
                plt.title(title, pad=20)
                
            else:
                # Handle bar/line chart case
                current_app.logger.info("Creating Budget vs Actual bar/line chart")
                
                # Extract data
                budget_trace = next((d for d in data['chart_data'] if d.get('name') == 'Budget'), None)
                actual_trace = next((d for d in data['chart_data'] if d.get('name') == 'Actual'), None)
                
                if not budget_trace or not actual_trace:
                    raise ValueError("Missing Budget or Actual data traces")
                
                # Create plot
                plt.figure(figsize=(10, 6))
                bar_width = 0.35
                months = budget_trace['x']
                x_pos = np.arange(len(months))
                
                if data['ppt_type'] == 'line':
                    plt.plot(months, budget_trace['y'], 
                            marker='o', 
                            color='#2E86AB', 
                            label='Budget',
                            linewidth=2)
                    plt.plot(months, actual_trace['y'], 
                            marker='o', 
                            color='#FF8C00', 
                            label='Actual',
                            linewidth=2)
                else:  # Default to bar chart
                    plt.bar(x_pos - bar_width/2, budget_trace['y'], 
                           bar_width, 
                           color='#2E86AB',
                           label='Budget')
                    plt.bar(x_pos + bar_width/2, actual_trace['y'], 
                           bar_width, 
                           color='#FF8C00',
                           label='Actual')
                
                plt.xlabel('Month')
                plt.ylabel('Value')
                plt.xticks(x_pos, months, rotation=45, ha='right')
                plt.legend()
                plt.grid(True, linestyle='--', alpha=0.3)
            
            plt.tight_layout()
            
            # Save plot to buffer
            img_buffer = BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            # Add image to slide
            slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
            
            # Save PPT to buffer
            ppt_buffer = BytesIO()
            ppt.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            return send_file(
                ppt_buffer,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name=f"budget_vs_actual_{data['ppt_type']}.pptx"
            )
        
        # Validate required fields
        required_fields = ['vis_label', 'table_name', 'chart_data', 'x_col', 'y_col', 'ppt_type']
        missing_fields = [field for field in required_fields if field not in data]
        if missing_fields:
            current_app.logger.error(f"Missing required fields: {missing_fields}")
            return jsonify({'error': f'Missing required fields: {missing_fields}'}), 400

        # Process chart data - handle both array and object formats
        try:
            chart_data = data['chart_data']
            
            # If it's Plotly figure data, extract the actual data points
            if isinstance(chart_data, dict) and 'data' in chart_data:
                chart_data = chart_data['data']
                
            # Convert to DataFrame
            if isinstance(chart_data, list):
                # Handle array of traces
                if len(chart_data) > 0 and 'x' in chart_data[0] and 'y' in chart_data[0]:
                    # This is Plotly trace data - combine all traces
                    combined_data = []
                    for trace in chart_data:
                        if 'name' in trace:  # For grouped data
                            for i, x_val in enumerate(trace['x']):
                                combined_data.append({
                                    data['x_col']: x_val,
                                    data['y_col']: trace['y'][i],
                                    'Metric': trace['name']
                                })
                        else:  # For single trace
                            for i, x_val in enumerate(trace['x']):
                                combined_data.append({
                                    data['x_col']: x_val,
                                    data['y_col']: trace['y'][i]
                                })
                    chart_df = pd.DataFrame(combined_data)
                else:
                    # Regular array of objects
                    chart_df = pd.DataFrame(chart_data)
            else:
                current_app.logger.error("Invalid chart_data format")
                return jsonify({'error': 'chart_data must be an array or Plotly figure data'}), 400
                
            current_app.logger.info(f"Processed chart data shape: {chart_df.shape}")
            
            if chart_df.empty:
                current_app.logger.error("Empty chart data after processing")
                return jsonify({'error': 'No valid chart data available'}), 400
                
            # Ensure required columns exist
            if data['x_col'] not in chart_df.columns:
                current_app.logger.error(f"x_col '{data['x_col']}' not found in data")
                return jsonify({'error': f"Column {data['x_col']} not found in chart data"}), 400
                
            if data['y_col'] not in chart_df.columns:
                current_app.logger.error(f"y_col '{data['y_col']}' not found in data")
                return jsonify({'error': f"Column {data['y_col']} not found in chart data"}), 400

            # Convert numeric data
            chart_df[data['y_col']] = pd.to_numeric(
                chart_df[data['y_col']].astype(str).str.replace(r'[^\d.-]', '', regex=True),
                errors='coerce'
            )
            chart_df = chart_df.dropna(subset=[data['y_col']])
            
            if chart_df.empty:
                current_app.logger.error("No numeric data after conversion")
                return jsonify({'error': 'No valid numeric data available'}), 400
                
        except Exception as e:
            current_app.logger.error(f"Data processing failed: {str(e)}", exc_info=True)
            return jsonify({'error': f'Invalid chart data format: {str(e)}'}), 400

        # Generate PPT
        try:
            ppt = Presentation()
            slide_layout = ppt.slide_layouts[5]  # Title and content
            slide = ppt.slides.add_slide(slide_layout)
            
            # Set title
            title = f"{data['vis_label']} - {data['table_name']}"
            if data.get('selected_filter') and data['selected_filter'] != "Select All":
                title += f" ({data['selected_filter']})"
            slide.shapes.title.text = title
            
            # Create chart
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # Determine chart type
            chart_type = data['ppt_type'].lower()
            
            if chart_type == 'pie':
                # Filter out non-positive values for pie charts
                chart_df = chart_df[chart_df[data['y_col']] > 0]
                if chart_df.empty:
                    raise ValueError("No positive values for pie chart")
                    
                wedges, texts, autotexts = ax.pie(
                    chart_df[data['y_col']],
                    labels=chart_df[data['x_col']],
                    autopct='%1.1f%%',
                    startangle=90,
                    colors=['#2E86AB', '#FF8C00', '#A23B72', '#F18F01', '#C73E1D']
                )
                ax.axis('equal')  # Equal aspect ratio ensures pie is drawn as circle
                
            elif chart_type == 'bar' or chart_type == 'column':
                if 'Metric' in chart_df.columns:  # Grouped bar chart
                    groups = chart_df.groupby([data['x_col'], 'Metric'])[data['y_col']].sum().unstack()
                    groups.plot(kind='bar', ax=ax, color=['#2E86AB', '#FF8C00'])
                else:
                    chart_df.plot(
                        x=data['x_col'],
                        y=data['y_col'],
                        kind='bar',
                        ax=ax,
                        color=data.get('color_override', '#2E86AB')
                    )
                ax.set_ylabel(data['y_col'])
                
            elif chart_type == 'line':
                if 'Metric' in chart_df.columns:  # Multiple lines
                    for name, group in chart_df.groupby('Metric'):
                        group.plot(
                            x=data['x_col'],
                            y=data['y_col'],
                            ax=ax,
                            label=name,
                            marker='o',
                            color='#FF8C00' if 'Actual' in str(name) else '#2E86AB'
                        )
                else:
                    chart_df.plot(
                        x=data['x_col'],
                        y=data['y_col'],
                        kind='line',
                        ax=ax,
                        marker='o',
                        color=data.get('color_override', '#2E86AB')
                    )
                ax.set_ylabel(data['y_col'])
                
            ax.set_xlabel(data['x_col'])
            ax.set_title(title)
            ax.grid(True, linestyle='--', alpha=0.7)
            
            # Rotate x-axis labels for better readability
            plt.xticks(rotation=45, ha='right')
            
            # Save plot to buffer
            img_buffer = BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            # Add image to slide
            slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
            
            # Save PPT to buffer
            ppt_buffer = BytesIO()
            ppt.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            return send_file(
                ppt_buffer,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name=f"{secure_filename(title.lower().replace(' ', '_'))}.pptx"
            )
            
        except Exception as e:
            current_app.logger.error(f"PPT generation failed: {str(e)}", exc_info=True)
            return jsonify({'error': f'PPT generation failed: {str(e)}'}), 500

    except Exception as e:
        current_app.logger.error(f"Unexpected error in PPT endpoint: {str(e)}", exc_info=True)
        return jsonify({'error': f'Internal server error: {str(e)}'}), 500
    
@main_bp.route('/download-budget-actual-ppt', methods=['POST'])
def download_budget_actual_ppt():
    try:
        data = request.get_json()
        
        # Validate required data
        if not all(key in data for key in ['months', 'budget', 'actual', 'title']):
            return jsonify({'error': 'Missing required data'}), 400

        # Create PowerPoint
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title and content
        
        # Set title
        title = data['title']
        if data.get('selected_filter') and data['selected_filter'] != "Select All":
            title += f" ({data['selected_filter']})"
        slide.shapes.title.text = title
        
        # Create plot with consistent month ordering
        plt.figure(figsize=(10, 6))
        bar_width = 0.35
        index = np.arange(len(data['months']))
        
        # Plot bars with consistent colors
        plt.bar(index - bar_width/2, data['budget'], bar_width, 
                label='Budget', color='#2E86AB')
        plt.bar(index + bar_width/2, data['actual'], bar_width,
                label='Actual', color='#FF8C00')
        
        # Customize plot to match frontend
        plt.xlabel('Month')
        plt.ylabel('Value')
        plt.title(title)
        plt.xticks(index, data['months'], rotation=45, ha='right')
        plt.legend()
        plt.grid(True, linestyle='--', alpha=0.3)
        plt.tight_layout()
        
        # Add value labels
        for i, (budget, actual) in enumerate(zip(data['budget'], data['actual'])):
            plt.text(i - bar_width/2, budget + 0.05 * max(data['budget']), 
                    f'{budget:,.0f}', ha='center', va='bottom')
            plt.text(i + bar_width/2, actual + 0.05 * max(data['actual']), 
                    f'{actual:,.0f}', ha='center', va='bottom')
        
        # Save plot to buffer
        img_buffer = BytesIO()
        plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
        img_buffer.seek(0)
        plt.close()
        
        # Add image to slide
        slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
        
        # Save PPT to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return send_file(
            ppt_buffer,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='budget_vs_actual.pptx'
        )
        
    except Exception as e:
        current_app.error(f"Error creating Budget vs Actual PPT: {str(e)}")
        return jsonify({'error': str(e)}), 500

@main_bp.route('/download-master-ppt', methods=['POST'])
def download_master_ppt():
    try:
        data = request.get_json()
        current_app.logger.info("Master PPT download request received")
        
        # Validate required fields
        required_fields = ['all_data', 'table_name', 'selected_sheet', 'visual_type']
        missing_fields = [field for field in required_fields if field not in data]
        if missing_fields:
            current_app.logger.error(f"Missing required fields: {missing_fields}")
            return jsonify({'error': f'Missing required fields: {missing_fields}'}), 400

        # Extract data
        all_data = data['all_data']
        table_name = data['table_name']
        selected_sheet = data['selected_sheet']
        visual_type = data['visual_type'].lower().replace(" chart", "")
        selected_filter = data.get('selected_filter')

        ppt = Presentation()

        # Title slide
        title_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = f"Complete Analysis Report - {table_name}"
        if selected_filter and selected_filter != "Select All":
            title += f" - {selected_filter}"
        title_slide.shapes.title.text = title
        if len(title_slide.placeholders) > 1:
            subtitle = title_slide.placeholders[1]
            subtitle.text = f"Sheet: {selected_sheet}\nGenerated on: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

        # Loop through all data items
        for item in all_data:
            try:
                label = item['label']
                chart_data = item['data']

                if not chart_data:
                    current_app.logger.debug(f"Skipping {label}: No data provided")
                    continue

                # Convert to DataFrame
                if isinstance(chart_data, list):
                    chart_df = pd.DataFrame(chart_data)
                elif isinstance(chart_data, dict):
                    chart_df = pd.DataFrame([chart_data])
                else:
                    current_app.logger.error(f"Invalid chart_data type for {label}")
                    continue

                chart_df = make_jsonly_serializable(chart_df)
                if chart_df.empty:
                    current_app.logger.warning(f"Skipping {label}: Data is empty after cleaning")
                    continue

                # Determine x_col and y_col
                x_col = item.get('x_col')
                y_col = item.get('y_col')
                if not x_col or not y_col:
                    if len(chart_df.columns) >= 2:
                        x_col = chart_df.columns[0]
                        y_col = chart_df.columns[1]
                    else:
                        current_app.logger.warning(f"Skipping {label}: Insufficient columns")
                        continue

                # Special unified handling for Budget vs Actual
                if label == "Budget vs Actual":
                    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                    slide_title = f"{label} - {table_name}"
                    if selected_filter and selected_filter != "Select All":
                        slide_title += f" ({selected_filter})"
                    slide.shapes.title.text = slide_title

                    budget_data = chart_df[chart_df['Metric'] == 'Budget']
                    actual_data = chart_df[chart_df['Metric'] == 'Actual']

                    plt.figure(figsize=(10, 6))

                    if visual_type == 'line':
                        months = budget_data[x_col].tolist()
                        plt.plot(months, budget_data[y_col], marker='o',
                                 label='Budget', color='#2E86AB')
                        plt.plot(months, actual_data[y_col], marker='o',
                                 label='Actual', color='#FF8C00')
                        plt.legend()
                        plt.grid(True, linestyle='--', alpha=0.3)

                    elif visual_type == 'pie':
                        budget_value = budget_data[y_col].sum()
                        actual_value = actual_data[y_col].sum()
                        if budget_value > 0 and actual_value > 0:
                            plt.pie([budget_value, actual_value],
                                    labels=['Budget', 'Actual'],
                                    autopct='%1.1f%%',
                                    startangle=90,
                                    colors=['#2E86AB', '#FF8C00'])
                            plt.axis('equal')

                    else:  # Default bar
                        bar_width = 0.35
                        months = budget_data[x_col].tolist()
                        x_pos = np.arange(len(months))
                        plt.bar(x_pos - bar_width/2, budget_data[y_col], bar_width,
                                label='Budget', color='#2E86AB')
                        plt.bar(x_pos + bar_width/2, actual_data[y_col], bar_width,
                                label='Actual', color='#FF8C00')
                        plt.xticks(x_pos, months, rotation=45, ha='right')
                        plt.legend()
                        plt.grid(True, linestyle='--', alpha=0.3)

                    plt.title(slide_title)
                    plt.tight_layout()
                    img_buffer = BytesIO()
                    plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
                    img_buffer.seek(0)
                    plt.close()

                    slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
                    continue  # done with this slide

                # Ensure numeric data for other charts
                try:
                    chart_df[y_col] = pd.to_numeric(
                        chart_df[y_col].astype(str).str.replace(r'[^\d.-]', '', regex=True),
                        errors='coerce'
                    )
                    chart_df = chart_df.dropna(subset=[y_col])
                except Exception as e:
                    current_app.logger.warning(f"Could not convert {y_col} to numeric for {label}: {str(e)}")
                    continue
                if chart_df.empty:
                    current_app.logger.warning(f"Skipping {label}: No numeric data after conversion")
                    continue

                # Filter for pie
                if visual_type == 'pie':
                    chart_df = chart_df[chart_df[y_col] > 0].copy()
                    if chart_df.empty:
                        current_app.logger.warning(f"Skipping {label}: No positive values for pie chart")
                        continue

                # Generic chart slide
                slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                slide_title = f"{label} - {table_name}"
                if selected_filter and selected_filter != "Select All":
                    slide_title += f" ({selected_filter})"
                slide.shapes.title.text = slide_title

                plt.figure(figsize=(10, 6))

                if visual_type == 'pie':
                    plt.pie(chart_df[y_col], labels=chart_df[x_col],
                            autopct='%1.1f%%', startangle=90,
                            colors=sns.color_palette('Set2'))
                    plt.axis('equal')

                elif visual_type == 'bar':
                    if 'Metric' in chart_df.columns:
                        pivot_df = chart_df.pivot(index=x_col, columns='Metric', values=y_col)
                        pivot_df.plot(kind='bar', ax=plt.gca(), color=['#2E86AB', '#FF8C00'])
                    else:
                        plt.bar(chart_df[x_col], chart_df[y_col], color='#2E86AB')
                    plt.xticks(rotation=45, ha='right')

                elif visual_type == 'line':
                    if 'Metric' in chart_df.columns:
                        for name, group in chart_df.groupby('Metric'):
                            plt.plot(group[x_col], group[y_col], marker='o', label=name,
                                     color='#FF8C00' if 'Actual' in str(name) else '#2E86AB')
                        plt.legend()
                    else:
                        plt.plot(chart_df[x_col], chart_df[y_col], marker='o', color='#2E86AB')
                    plt.xticks(rotation=45, ha='right')

                plt.title(slide_title)
                plt.tight_layout()
                img_buffer = BytesIO()
                plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
                img_buffer.seek(0)
                plt.close()
                slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))

            except Exception as e:
                current_app.logger.error(f"Error creating chart for {label}: {str(e)}", exc_info=True)
                continue

        # Save final PPT
        ppt_buffer = BytesIO()
        ppt.save(ppt_buffer)
        ppt_buffer.seek(0)
        return send_file(
            ppt_buffer,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f"{secure_filename(table_name.lower().replace(' ', '_'))}_report.pptx"
        )

    except Exception as e:
        current_app.logger.error(f"Error generating master PPT: {str(e)}", exc_info=True)
        return jsonify({'error': f'Failed to generate PPT: {str(e)}'}), 500



