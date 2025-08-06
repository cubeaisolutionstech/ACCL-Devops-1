from flask import Blueprint, request, jsonify, send_file
import traceback
import uuid
from datetime import datetime
import logging
import pandas as pd
import os
from io import BytesIO

from utils.executive_budget_vs_billed import (
    calculate_executive_budget_vs_billed, 
    auto_map_executive_columns, 
    load_data,
    get_executives_and_branches,
    get_available_months
)
from utils.executive_odc import (
    load_data as load_od_data,
    auto_map_od_columns,
    get_od_executives_and_branches,
    get_available_months_od,
    calculate_od_values
)
from utils.executive_product_growth import (
    load_data as load_product_data,
    auto_map_product_growth_columns,
    get_product_growth_options,
    get_product_growth_months,
    calculate_product_growth
)

# NEW IMPORTS for Customer and OD Target Analysis (both from same file)
from utils.executive_nbc import (
    auto_map_customer_columns,
    get_customer_options,
    create_customer_table,
    auto_map_od_columns as auto_map_od_target_columns,
    get_od_options as get_od_target_options,
    filter_os_qty,
    load_data,
)

from utils.executive_ppt_generator import (
    create_executive_budget_ppt,
    create_executive_od_ppt,
    create_product_growth_ppt,
    create_consolidated_ppt,
    validate_ppt_data,
    validate_product_growth_ppt_data,
    create_nbc_individual_ppt,
    create_od_individual_ppt
)

executive_bp = Blueprint('executive', __name__, url_prefix='/api/executive')
logger = logging.getLogger(__name__)

@executive_bp.route('/customer_auto_map_columns', methods=['POST'])
def customer_auto_map_columns():
    """Auto-map columns for customer analysis"""
    try:
        data = request.get_json()
        sales_file_path = data.get('sales_file_path')
        
        if not sales_file_path:
            return jsonify({'error': 'Sales file path is required'}), 400
        
        # Load data
        sales_df = load_data(sales_file_path)
        
        # Auto-map columns
        mapping_result = auto_map_customer_columns(sales_df)
        
        return jsonify(mapping_result)
        
    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in customer_auto_map_columns route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/customer_get_options', methods=['POST'])
def customer_get_options():
    """Get available options for customer analysis"""
    try:
        data = request.get_json()
        sales_file_path = data.get('sales_file_path')
        
        if not sales_file_path:
            return jsonify({'error': 'Sales file path is required'}), 400
        
        # Load data
        sales_df = load_data(sales_file_path)
        
        # Auto-map to find columns
        mapping_result = auto_map_customer_columns(sales_df)
        
        if not mapping_result.get('success'):
            return jsonify({'error': 'Failed to map columns'}), 400
        
        date_col = mapping_result['mapping'].get('date')
        branch_col = mapping_result['mapping'].get('branch')
        executive_col = mapping_result['mapping'].get('executive')
        
        # Get options
        options_result = get_customer_options(sales_df, date_col, branch_col, executive_col)
        
        return jsonify({
            **options_result,
            'column_mapping': mapping_result
        })
        
    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in customer_get_options route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/calculate_customer_analysis', methods=['POST'])
def calculate_customer_analysis():
    """Calculate customer analysis"""
    try:
        data = request.get_json()
        logger.info(f"Received Customer Analysis request data keys: {list(data.keys())}")
        
        # Get file path
        sales_file_path = data.get('sales_file_path')
        
        if not sales_file_path:
            return jsonify({'error': 'Sales file path is required'}), 400
        
        # Get column mappings
        date_col = data.get('date_column')
        branch_col = data.get('branch_column')
        customer_id_col = data.get('customer_id_column')
        executive_col = data.get('executive_column')
        
        # Get filters
        selected_months = data.get('selected_months', [])
        selected_branches = data.get('selected_branches', [])
        selected_executives = data.get('selected_executives', [])
        
        logger.info(f"Columns: date={date_col}, branch={branch_col}, customer_id={customer_id_col}, executive={executive_col}")
        logger.info(f"Filters: months={len(selected_months)}, branches={len(selected_branches)}, executives={len(selected_executives)}")
        
        # Validate required mappings
        required_cols = [date_col, branch_col, customer_id_col, executive_col]
        if not all(required_cols):
            missing = [name for name, val in zip(['date', 'branch', 'customer_id', 'executive'], required_cols) if not val]
            return jsonify({'error': f'Missing required column mappings: {missing}'}), 400
        
        # Load data
        sales_df = load_data(sales_file_path)
        
        logger.info(f"Sales data loaded: {len(sales_df)} rows")
        
        # Call calculation function
        result = create_customer_table(
            sales_df=sales_df,
            date_col=date_col,
            branch_col=branch_col,
            customer_id_col=customer_id_col,
            executive_col=executive_col,
            selected_months=selected_months if selected_months else None,
            selected_branches=selected_branches if selected_branches else None,
            selected_executives=selected_executives if selected_executives else None
        )
        
        if result.get('success'):
            logger.info(f"Customer analysis calculation successful!")
            return jsonify(result)
        else:
            logger.error(f"Customer analysis calculation failed: {result.get('error', 'Unknown error')}")
            return jsonify({'error': result.get('error', 'Unknown error')}), 500
            
    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in calculate_customer_analysis route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/generate_customer_ppt', methods=['POST'])
def generate_customer_ppt():
    """Generate PowerPoint presentation for Customer Analysis results"""
    try:
        data = request.get_json()
        logger.info(f"Received Customer PPT generation request with keys: {list(data.keys())}")

        # Get results data
        results_data = data.get('results_data')
        if not results_data:
            return jsonify({'error': 'Results data is required for PPT generation'}), 400

        # Get optional parameters
        title = data.get('title', 'Customer Analysis Report')

        # 🔥 Load static logo
        import os
        logo_file = None
        static_logo_path = os.path.join("static", "logo.jpeg")
        if os.path.exists(static_logo_path):
            logo_file = open(static_logo_path, "rb")
            logger.info(f"✅ Static logo loaded for Customer PPT from: {static_logo_path}")
        else:
            logger.warning("⚠️ Static logo not found, continuing without logo.")

        logger.info(f"Generating Customer PPT: {title}")

        # Validate structure
        if not isinstance(results_data, dict) or 'results' not in results_data:
            return jsonify({'error': 'Invalid results data format'}), 400

        ppt_buffers = []

        # Generate PPT for each financial year
        for fin_year, fy_data in results_data['results'].items():
            customer_df = pd.DataFrame(fy_data['data'])
            sorted_months = fy_data['sorted_months']

            # Generate individual PPT for this financial year
            ppt_buffer = create_nbc_individual_ppt(
                customer_df=customer_df,
                customer_title="Number of Billed Customers",
                sorted_months=sorted_months,
                financial_year=fin_year,
                logo_file=logo_file
            )

            if ppt_buffer:
                ppt_buffers.append((fin_year, ppt_buffer))

        if logo_file:
            logo_file.close()  # 🧼 Clean up

        if not ppt_buffers:
            return jsonify({'error': 'Failed to generate PowerPoint presentation'}), 500

        # Use the first generated PPT (could be enhanced to zip/combine later)
        final_buffer = ppt_buffers[0][1]

        # Generate unique filename
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"Customer_Analysis_{timestamp}.pptx"

        logger.info(f"✅ Successfully generated Customer PPT: {filename}")

        return send_file(
            final_buffer,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        logger.error(f"Error in generate_customer_ppt route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

# ============= OD TARGET ANALYSIS ROUTES =============

@executive_bp.route('/od_target_auto_map_columns', methods=['POST'])
def od_target_auto_map_columns():
    """Auto-map columns for OD target analysis"""
    try:
        data = request.get_json()
        os_file_path = data.get('os_file_path')
        
        if not os_file_path:
            return jsonify({'error': 'OS file path is required'}), 400
        
        # Load data
        os_df = load_data(os_file_path)
        
        # Auto-map columns
        mapping_result = auto_map_od_target_columns(os_df)
        
        return jsonify(mapping_result)
        
    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in od_target_auto_map_columns route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/od_target_get_options', methods=['POST'])
def od_target_get_options():
    """Get available options for OD target analysis"""
    try:
        data = request.get_json()
        os_file_path = data.get('os_file_path')
        
        if not os_file_path:
            return jsonify({'error': 'OS file path is required'}), 400
        
        # Load data
        os_df = load_data(os_file_path)
        
        # Auto-map to find columns
        mapping_result = auto_map_od_target_columns(os_df)
        
        if not mapping_result.get('success'):
            return jsonify({'error': 'Failed to map columns'}), 400
        
        due_date_col = mapping_result['mapping'].get('due_date')
        area_col = mapping_result['mapping'].get('area')
        executive_col = mapping_result['mapping'].get('executive')
        
        # Get options
        options_result = get_od_target_options(os_df, due_date_col, area_col, executive_col)
        
        return jsonify({
            **options_result,
            'column_mapping': mapping_result
        })
        
    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in od_target_get_options route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/calculate_od_target', methods=['POST'])
def calculate_od_target():
    """Calculate OD target analysis"""
    try:
        data = request.get_json()
        logger.info(f"Received OD Target request data keys: {list(data.keys())}")
        
        # Get file path
        os_file_path = data.get('os_file_path')
        
        if not os_file_path:
            return jsonify({'error': 'OS file path is required'}), 400
        
        # Get column mappings
        area_col = data.get('area_column')
        net_value_col = data.get('net_value_column')
        due_date_col = data.get('due_date_column')
        executive_col = data.get('executive_column')
        
        # Get filters
        selected_branches = data.get('selected_branches', [])
        selected_years = data.get('selected_years', [])
        till_month = data.get('till_month')
        selected_executives = data.get('selected_executives', [])
        
        logger.info(f"Columns: area={area_col}, net_value={net_value_col}, due_date={due_date_col}, executive={executive_col}")
        logger.info(f"Filters: branches={len(selected_branches)}, years={selected_years}, till_month={till_month}, executives={len(selected_executives)}")
        
        # Validate required mappings
        required_cols = [area_col, net_value_col, due_date_col, executive_col]
        if not all(required_cols):
            missing = [name for name, val in zip(['area', 'net_value', 'due_date', 'executive'], required_cols) if not val]
            return jsonify({'error': f'Missing required column mappings: {missing}'}), 400
        
        # Load data
        os_df = load_data(os_file_path)
        
        logger.info(f"OS data loaded: {len(os_df)} rows")
        
        # Call calculation function
        result = filter_os_qty(
            os_df=os_df,
            os_area_col=area_col,
            os_qty_col=net_value_col,  # Using net_value as quantity for OD target
            os_due_date_col=due_date_col,
            os_exec_col=executive_col,
            selected_branches=selected_branches if selected_branches else None,
            selected_years=selected_years if selected_years else None,
            till_month=till_month,
            selected_executives=selected_executives if selected_executives else None
        )
        
        if result.get('success'):
            logger.info(f"OD target calculation successful!")
            return jsonify(result)
        else:
            logger.error(f"OD target calculation failed: {result.get('error', 'Unknown error')}")
            return jsonify({'error': result.get('error', 'Unknown error')}), 500
            
    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in calculate_od_target route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/generate_od_target_ppt', methods=['POST'])
def generate_od_target_ppt():
    """Generate PowerPoint presentation for OD Target results"""
    try:
        data = request.get_json()
        logger.info(f"Received OD Target PPT generation request with keys: {list(data.keys())}")

        # Get results data
        results_data = data.get('results_data')
        if not results_data:
            return jsonify({'error': 'Results data is required for PPT generation'}), 400

        # Get optional parameters
        title = data.get('title', 'OD Target Report')

        # 🔥 Load static logo
        import os
        logo_file = None
        static_logo_path = os.path.join("static", "logo.jpeg")
        if os.path.exists(static_logo_path):
            logo_file = open(static_logo_path, "rb")
            logger.info(f"✅ Static logo loaded for OD Target PPT from: {static_logo_path}")
        else:
            logger.warning("⚠️ Static logo not found. Proceeding without logo.")

        logger.info(f"Generating OD Target PPT: {title}")

        # Validate structure
        if not isinstance(results_data, dict) or 'data' not in results_data:
            return jsonify({'error': 'Invalid results data format'}), 400

        # Convert to DataFrame
        od_target_df = pd.DataFrame(results_data['data'])

        # Add date range to title if present
        period_title = title
        if results_data.get('start_date') and results_data.get('end_date'):
            period_title = f"{title} ({results_data['start_date']} to {results_data['end_date']})"

        # Generate PPT
        ppt_buffer = create_od_individual_ppt(
            od_target_df=od_target_df,
            od_title=period_title,
            logo_file=logo_file
        )

        if logo_file:
            logo_file.close()  # 🧼 Always clean up

        if not ppt_buffer:
            return jsonify({'error': 'Failed to generate PowerPoint presentation'}), 500

        # Create filename
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"OD_Target_Analysis_{timestamp}.pptx"

        logger.info(f"✅ Successfully generated OD Target PPT: {filename}")

        return send_file(
            ppt_buffer,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        logger.error(f"Error in generate_od_target_ppt route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500


# ============= PRODUCT GROWTH ROUTES (EXISTING) =============

@executive_bp.route('/product_auto_map_columns', methods=['POST'])
def product_auto_map_columns():
    """Auto-map columns for product growth analysis"""
    try:
        data = request.get_json()
        ly_file_path = data.get('ly_file_path')
        cy_file_path = data.get('cy_file_path')
        budget_file_path = data.get('budget_file_path')
        
        if not ly_file_path or not cy_file_path or not budget_file_path:
            return jsonify({'error': 'All three file paths are required'}), 400
        
        # Load data
        ly_df = load_product_data(ly_file_path)
        cy_df = load_product_data(cy_file_path)
        budget_df = load_product_data(budget_file_path)
        
        # Auto-map columns
        mapping_result = auto_map_product_growth_columns(ly_df, cy_df, budget_df)
        
        return jsonify(mapping_result)
        
    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in product_auto_map_columns route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/product_get_options', methods=['POST'])
def product_get_options():
    """Get available options for product growth analysis"""
    try:
        data = request.get_json()
        ly_file_path = data.get('ly_file_path')
        cy_file_path = data.get('cy_file_path')
        budget_file_path = data.get('budget_file_path')
        
        if not ly_file_path or not cy_file_path or not budget_file_path:
            return jsonify({'error': 'All three file paths are required'}), 400
        
        # Load data
        ly_df = load_product_data(ly_file_path)
        cy_df = load_product_data(cy_file_path)
        budget_df = load_product_data(budget_file_path)
        
        # Auto-map to find columns
        mapping_result = auto_map_product_growth_columns(ly_df, cy_df, budget_df)
        
        ly_exec_col = mapping_result['ly_mapping'].get('executive')
        cy_exec_col = mapping_result['cy_mapping'].get('executive')
        budget_exec_col = mapping_result['budget_mapping'].get('executive')
        ly_company_group_col = mapping_result['ly_mapping'].get('company_group')
        cy_company_group_col = mapping_result['cy_mapping'].get('company_group')
        budget_company_group_col = mapping_result['budget_mapping'].get('company_group')
        ly_product_group_col = mapping_result['ly_mapping'].get('product_group')
        cy_product_group_col = mapping_result['cy_mapping'].get('product_group')
        budget_product_group_col = mapping_result['budget_mapping'].get('product_group')
        ly_sl_code_col = mapping_result['ly_mapping'].get('sl_code')
        cy_sl_code_col = mapping_result['cy_mapping'].get('sl_code')
        budget_sl_code_col = mapping_result['budget_mapping'].get('sl_code')
        
        # Get options
        options_data = get_product_growth_options(
            ly_df, cy_df, budget_df,
            ly_exec_col, cy_exec_col, budget_exec_col,
            ly_company_group_col, cy_company_group_col, budget_company_group_col,
            ly_product_group_col, cy_product_group_col, budget_product_group_col,
            ly_sl_code_col, cy_sl_code_col, budget_sl_code_col
        )
        
        # Get available months
        ly_date_col = mapping_result['ly_mapping'].get('date')
        cy_date_col = mapping_result['cy_mapping'].get('date')
        months_data = get_product_growth_months(ly_df, cy_df, ly_date_col, cy_date_col)
        
        return jsonify({
            **options_data,
            **months_data,
            'column_mapping': mapping_result
        })
        
    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in product_get_options route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/calculate_product_growth', methods=['POST'])
def calculate_product_growth_route():
    """Calculate product growth analysis"""
    try:
        data = request.get_json()
        logger.info(f"Received Product Growth request data keys: {list(data.keys())}")
        
        # Get file paths
        ly_file_path = data.get('ly_file_path')
        cy_file_path = data.get('cy_file_path')
        budget_file_path = data.get('budget_file_path')
        
        if not ly_file_path or not cy_file_path or not budget_file_path:
            return jsonify({'error': 'All three file paths are required'}), 400
        
        # Get LY column mappings
        ly_date = data.get('ly_date')
        ly_quantity = data.get('ly_quantity')
        ly_value = data.get('ly_value')
        ly_company_group = data.get('ly_company_group')
        ly_product_group = data.get('ly_product_group')
        ly_executive = data.get('ly_executive')
        ly_sl_code = data.get('ly_sl_code')
        
        # Get CY column mappings
        cy_date = data.get('cy_date')
        cy_quantity = data.get('cy_quantity')
        cy_value = data.get('cy_value')
        cy_company_group = data.get('cy_company_group')
        cy_product_group = data.get('cy_product_group')
        cy_executive = data.get('cy_executive')
        cy_sl_code = data.get('cy_sl_code')
        
        # Get Budget column mappings
        budget_quantity = data.get('budget_quantity')
        budget_value = data.get('budget_value')
        budget_company_group = data.get('budget_company_group')
        budget_product_group = data.get('budget_product_group')
        budget_executive = data.get('budget_executive')
        budget_sl_code = data.get('budget_sl_code')
        
        # Get filters
        selected_executives = data.get('selected_executives', [])
        selected_company_groups = data.get('selected_company_groups', [])
        selected_product_groups = data.get('selected_product_groups', [])
        ly_month = data.get('ly_month')
        cy_month = data.get('cy_month')
        
        logger.info(
            f"LY columns: date={ly_date}, qty={ly_quantity}, value={ly_value}, "
            f"company={ly_company_group}, product={ly_product_group}"
        )
        logger.info(
            f"CY columns: date={cy_date}, qty={cy_quantity}, value={cy_value}, "
            f"company={cy_company_group}, product={cy_product_group}"
        )
        logger.info(
            f"Budget columns: qty={budget_quantity}, value={budget_value}, "
            f"company={budget_company_group}, product={budget_product_group}"
        )
        logger.info(
            f"Selected filters: executives={len(selected_executives)}, "
            f"companies={len(selected_company_groups)}, products={len(selected_product_groups)}"
        )
        
        # Validate required mappings
        required_ly_cols = [ly_date, ly_quantity, ly_value, ly_company_group, ly_product_group, ly_executive]
        required_cy_cols = [cy_date, cy_quantity, cy_value, cy_company_group, cy_product_group, cy_executive]
        required_budget_cols = [budget_quantity, budget_value, budget_company_group, budget_product_group, budget_executive]
        
        if not all(required_ly_cols):
            missing = [name for name, val in zip(['date', 'quantity', 'value', 'company_group', 'product_group', 'executive'], required_ly_cols) if not val]
            return jsonify({'error': f'Missing required LY column mappings: {missing}'}), 400
            
        if not all(required_cy_cols):
            missing = [name for name, val in zip(['date', 'quantity', 'value', 'company_group', 'product_group', 'executive'], required_cy_cols) if not val]
            return jsonify({'error': f'Missing required CY column mappings: {missing}'}), 400
            
        if not all(required_budget_cols):
            missing = [name for name, val in zip(['quantity', 'value', 'company_group', 'product_group', 'executive'], required_budget_cols) if not val]
            return jsonify({'error': f'Missing required Budget column mappings: {missing}'}), 400
        
        # Load data
        ly_df = load_product_data(ly_file_path)
        cy_df = load_product_data(cy_file_path)
        budget_df = load_product_data(budget_file_path)
        
        logger.info(f"LY data loaded: {len(ly_df)} rows")
        logger.info(f"CY data loaded: {len(cy_df)} rows")
        logger.info(f"Budget data loaded: {len(budget_df)} rows")
        
        # Call calculation function - This now returns data in React-compatible format
        result = calculate_product_growth(
            ly_df, cy_df, budget_df, ly_month, cy_month,
            ly_date, cy_date, ly_quantity, cy_quantity, ly_value, cy_value,
            budget_quantity, budget_value, ly_company_group, cy_company_group, budget_company_group,
            ly_product_group, cy_product_group, budget_product_group,
            ly_sl_code, cy_sl_code, budget_sl_code,
            ly_executive, cy_executive, budget_executive,
            selected_executives, selected_company_groups, selected_product_groups
        )
        
        if result.get('success'):
            logger.info(f"Product Growth calculation successful! Returning streamlit_result with {len(result.get('streamlit_result', {}))} companies")
            return jsonify(result)
        else:
            logger.error(f"Product Growth calculation failed: {result.get('error', 'Unknown error')}")
            return jsonify({'error': result.get('error', 'Unknown error')}), 500
            
    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in calculate_product_growth route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

# Update this function in your executive_routes.py

@executive_bp.route('/generate_product_growth_ppt', methods=['POST'])
def generate_product_growth_ppt():
    """Generate PowerPoint presentation for Product Growth results with proper months - FIXED"""
    try:
        data = request.get_json()
        logger.info(f"Received Product Growth PPT generation request with keys: {list(data.keys())}")
        logger.info(f"🔍 Full request data: {data}")

        # Handle different data formats
        streamlit_result = data.get('streamlit_result')
        results_data = data.get('results_data')

        # Use streamlit_result if available, otherwise try results_data
        if streamlit_result:
            group_results = streamlit_result
            logger.info("Using streamlit_result data for PPT generation")
        elif results_data and isinstance(results_data, dict) and 'streamlit_result' in results_data:
            group_results = results_data['streamlit_result']
            logger.info("Using results_data.streamlit_result for PPT generation")
        elif results_data:
            group_results = transform_results_to_streamlit_format(results_data)
            logger.info("Transformed results_data to streamlit format")
        else:
            logger.error("No valid results data found in request")
            return jsonify({'error': 'Results data (streamlit_result) is required for PPT generation'}), 400

        if not group_results:
            logger.error("Group results is empty after processing")
            return jsonify({'error': 'No valid data found for PPT generation'}), 400

        # 🔒 Extract month info
        ly_month = data.get('ly_month')
        cy_month = data.get('cy_month')
        month_title = data.get('month_title')

        if not ly_month and results_data:
            ly_month = results_data.get('ly_month')
        if not cy_month and results_data:
            cy_month = results_data.get('cy_month')
        if not month_title and results_data:
            month_title = results_data.get('month_title')

        # Check nested keys
        for key in ['filters', 'request_data', 'params']:
            if key in data:
                nested_data = data[key]
                if isinstance(nested_data, dict):
                    ly_month = ly_month or nested_data.get('ly_month')
                    cy_month = cy_month or nested_data.get('cy_month')
                    month_title = month_title or nested_data.get('month_title')

        logger.info(f"🔍 Extracted ly_month: '{ly_month}'")
        logger.info(f"🔍 Extracted cy_month: '{cy_month}'")
        logger.info(f"🔍 Extracted month_title: '{month_title}'")
        print(f"🔍 DEBUG MONTHS: ly_month={ly_month}, cy_month={cy_month}, month_title={month_title}")

        # Final title logic
        if ly_month and cy_month:
            final_month_title = f"LY: {ly_month} vs CY: {cy_month}"
            logger.info(f"✅ Using extracted months for title: {final_month_title}")
        elif month_title and month_title != "Product Growth Analysis":
            final_month_title = month_title
            logger.info(f"✅ Using provided month_title: {final_month_title}")
        else:
            final_month_title = "Product Growth Analysis"
            logger.info(f"✅ Using default title: {final_month_title}")

        # 🔥 Load static logo
        import os
        logo_file = None
        static_logo_path = os.path.join("static", "logo.jpeg")
        if os.path.exists(static_logo_path):
            logo_file = open(static_logo_path, "rb")
            logger.info(f"✅ Static logo loaded for Product Growth PPT from: {static_logo_path}")
        else:
            logger.warning("⚠️ Static logo not found, continuing without logo.")

        logger.info(f"Generating Product Growth PPT for: {final_month_title}")
        logger.info(f"Found data for {len(group_results)} companies: {list(group_results.keys())}")

        # Generate the PPT
        ppt_buffer = create_product_growth_ppt(
            group_results, 
            final_month_title, 
            logo_file, 
            ly_month=ly_month, 
            cy_month=cy_month
        )

        if logo_file:
            logo_file.close()

        if not ppt_buffer:
            logger.error("Failed to create PPT buffer")
            return jsonify({'error': 'Failed to generate PowerPoint presentation'}), 500

        # Create unique filename
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        if ly_month and cy_month:
            ly_clean = str(ly_month).replace(' ', '_').replace('/', '_').replace('\\', '_')
            cy_clean = str(cy_month).replace(' ', '_').replace('/', '_').replace('\\', '_')
            filename = f"Executive_Product_Growth_LY_{ly_clean}_vs_CY_{cy_clean}_{timestamp}.pptx"
        else:
            filename = f"Executive_Product_Growth_{final_month_title.replace(' ', '_').replace(':', '').replace('/', '_')}_{timestamp}.pptx"

        logger.info(f"Successfully generated PPT: {filename}")

        return send_file(
            ppt_buffer,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        logger.error(f"Error in generate_product_growth_ppt route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

def transform_results_to_streamlit_format(results_data):
    """Transform various result formats to streamlit format - ENHANCED VERSION"""
    try:
        # Check if it's already in streamlit format
        if isinstance(results_data, dict):
            # Check if all values have 'qty_df' and 'value_df'
            if all(isinstance(v, dict) and 'qty_df' in v and 'value_df' in v for v in results_data.values()):
                logger.info("Data is already in streamlit format")
                return results_data
            
            # Check if there's a nested streamlit_result
            if 'streamlit_result' in results_data:
                return results_data['streamlit_result']
        
        logger.warning("Could not transform results_data to streamlit format")
        return {}
        
    except Exception as e:
        logger.error(f"Error transforming results to streamlit format: {str(e)}")
        return {}

def transform_results_to_streamlit_format(results_data):
   """Transform various result formats to streamlit format"""
   try:
       # Check if it's already in streamlit format
       if isinstance(results_data, dict):
           # Check if all values have 'qty_df' and 'value_df'
           if all(isinstance(v, dict) and 'qty_df' in v and 'value_df' in v for v in results_data.values()):
               logger.info("Data is already in streamlit format")
               return results_data
           
           # Check if there's a nested streamlit_result
           if 'streamlit_result' in results_data:
               return results_data['streamlit_result']
       
       logger.warning("Could not transform results_data to streamlit format")
       return {}
       
   except Exception as e:
       logger.error(f"Error transforming results to streamlit format: {str(e)}")
       return {}

def create_title_slide(prs, title, logo_file=None):
   """Create title slide for PPT"""
   try:
       from pptx.util import Inches, Pt
       from pptx.enum.text import PP_ALIGN
       from pptx.dml.color import RGBColor
       
       blank_slide_layout = prs.slide_layouts[6]
       title_slide = prs.slides.add_slide(blank_slide_layout)
       
       # Company name
       company_name = title_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
       company_frame = company_name.text_frame
       company_frame.text = "Asia Crystal Commodity LLP"
       p = company_frame.paragraphs[0]
       p.alignment = PP_ALIGN.CENTER
       p.font.name = "Times New Roman"
       p.font.size = Pt(36)
       p.font.bold = True
       p.font.color.rgb = RGBColor(0, 112, 192)
       
       # Add logo if provided
       if logo_file is not None:
           try:
               import base64
               logo_data = base64.b64decode(logo_file)
               logo_buffer = BytesIO(logo_data)
               logo = title_slide.shapes.add_picture(logo_buffer, Inches(5.665), Inches(1.5), width=Inches(2), height=Inches(2))
           except Exception as e:
               logger.error(f"Error adding logo to slide: {e}")
       
       # Title
       title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(1))
       title_frame = title_box.text_frame
       title_frame.text = title
       p = title_frame.paragraphs[0]
       p.alignment = PP_ALIGN.CENTER
       p.font.name = "Times New Roman"
       p.font.size = Pt(32)
       p.font.bold = True
       p.font.color.rgb = RGBColor(0, 128, 0)
       
       # Subtitle
       subtitle = title_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.33), Inches(1))
       subtitle_frame = subtitle.text_frame
       subtitle_frame.text = "ACCLLP"
       p = subtitle_frame.paragraphs[0]
       p.alignment = PP_ALIGN.CENTER
       p.font.name = "Times New Roman"
       p.font.size = Pt(28)
       p.font.bold = True
       p.font.color.rgb = RGBColor(0, 112, 192)
       
       return title_slide
       
   except Exception as e:
       logger.error(f"Error creating title slide: {e}")
       return None

def add_table_slide(prs, df, title, percent_cols=None):
   """Add table slide to PPT"""
   try:
       from pptx.util import Inches, Pt
       from pptx.enum.text import PP_ALIGN
       from pptx.dml.color import RGBColor
       import pandas as pd
       
       if percent_cols is None:
           percent_cols = []
       
       if df.empty:
           logger.warning(f"Empty DataFrame for slide: {title}")
           return
       
       slide_layout = prs.slide_layouts[6]
       slide = prs.slides.add_slide(slide_layout)
       
       # Add title
       title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
       title_frame = title_shape.text_frame
       title_frame.text = title
       p = title_frame.paragraphs[0]
       p.font.size = Pt(28)
       p.font.bold = True
       p.font.color.rgb = RGBColor(0, 112, 192)
       p.alignment = PP_ALIGN.CENTER
       
       # Create table
       columns = df.columns.tolist()
       num_rows = len(df) + 1
       num_cols = len(columns)
       
       if num_rows <= 1 or num_cols <= 0:
           logger.warning(f"Invalid table dimensions for slide: {title}")
           return
       
       table_height = min(Inches(0.4 * len(df) + 0.4), Inches(5.5))
       
       table = slide.shapes.add_table(
           num_rows, num_cols, 
           Inches(0.5), Inches(1.5), 
           Inches(12), table_height
       ).table
       
       # Set column widths
       if num_cols > 0:
           table.columns[0].width = Inches(3.0)
           remaining_width = 9.0
           if num_cols > 1:
               col_width = remaining_width / (num_cols - 1)
               for i in range(1, num_cols):
                   table.columns[i].width = Inches(col_width)
       
       # Add headers
       for i, col_name in enumerate(columns):
           cell = table.cell(0, i)
           cell.text = str(col_name)
           cell.fill.solid()
           cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
           cell.text_frame.paragraphs[0].font.size = Pt(14)
           cell.text_frame.paragraphs[0].font.bold = True
           cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
           cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
       
       # Add data rows
       for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
           is_total_row = 'TOTAL' in str(row.iloc[0]).upper()
           
           for col_idx, col_name in enumerate(columns):
               cell = table.cell(row_idx, col_idx)
               value = row[col_name]
               
               # Format percentage columns
               if col_idx in percent_cols and isinstance(value, (int, float)) and not pd.isna(value):
                   cell.text = f"{value}%"
               elif isinstance(value, (int, float)) and not pd.isna(value):
                   if value == int(value):
                       cell.text = str(int(value))
                   else:
                       cell.text = f"{value:.2f}"
               else:
                   cell.text = str(value) if not pd.isna(value) else ""
               
               cell.text_frame.paragraphs[0].font.size = Pt(12)
               cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
               cell.fill.solid()
               
               if is_total_row:
                   cell.fill.fore_color.rgb = RGBColor(211, 211, 211)
                   cell.text_frame.paragraphs[0].font.bold = True
               else:
                   if row_idx % 2 == 0:
                       cell.fill.fore_color.rgb = RGBColor(221, 235, 247)
                   else:
                       cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
       
       logger.info(f"Successfully added table slide: {title}")
       
   except Exception as e:
       logger.error(f"Error adding table slide '{title}': {e}", exc_info=True)

# ============= BUDGET VS BILLED ROUTES (EXISTING) =============

@executive_bp.route('/calculate_budget_vs_billed', methods=['POST'])
def calculate_budget_vs_billed():
   """Calculate executive budget vs billed analysis"""
   try:
       data = request.get_json()
       logger.info(f"Received request data keys: {list(data.keys())}")
       
       # Get file paths
       sales_file_path = data.get('sales_file_path')
       budget_file_path = data.get('budget_file_path')
       
       if not sales_file_path or not budget_file_path:
           return jsonify({'error': 'Sales and budget file paths are required'}), 400
       
       # Get sales column mappings
       sales_date = data.get('sales_date')
       sales_value = data.get('sales_value')
       sales_quantity = data.get('sales_quantity')
       sales_product_group = data.get('sales_product_group')
       sales_sl_code = data.get('sales_sl_code')
       sales_executive = data.get('sales_executive')
       sales_area = data.get('sales_area')
       
       # Get budget column mappings
       budget_value = data.get('budget_value')
       budget_quantity = data.get('budget_quantity')
       budget_product_group = data.get('budget_product_group')
       budget_sl_code = data.get('budget_sl_code')
       budget_executive = data.get('budget_executive')
       budget_area = data.get('budget_area')
       
       # Get filters
       selected_executives = data.get('selected_executives', [])
       selected_months = data.get('selected_months', [])
       selected_branches = data.get('selected_branches', [])
       
       logger.info(
           f"Sales columns: date={sales_date}, value={sales_value}, qty={sales_quantity}, "
           f"exec={sales_executive}, area={sales_area}"
       )
       logger.info(
           f"Budget columns: value={budget_value}, qty={budget_quantity}, "
           f"exec={budget_executive}, area={budget_area}"
       )
       logger.info(f"Selected executives: {selected_executives}")
       logger.info(f"Selected months: {selected_months}")
       logger.info(f"Selected branches: {selected_branches}")
       
       # Validate required mappings
       required_sales_cols = [sales_date, sales_value, sales_quantity, sales_executive]
       required_budget_cols = [budget_value, budget_quantity, budget_executive]
       
       if not all(required_sales_cols):
           missing = [name for name, val in zip(['date', 'value', 'quantity', 'executive'], required_sales_cols) if not val]
           return jsonify({'error': f'Missing required sales column mappings: {missing}'}), 400
           
       if not all(required_budget_cols):
           missing = [name for name, val in zip(['value', 'quantity', 'executive'], required_budget_cols) if not val]
           return jsonify({'error': f'Missing required budget column mappings: {missing}'}), 400
       
       # Call calculation function
       result = calculate_executive_budget_vs_billed(
           sales_file_path=sales_file_path,
           budget_file_path=budget_file_path,
           sales_date=sales_date,
           sales_value=sales_value,
           sales_quantity=sales_quantity,
           sales_product_group=sales_product_group,
           sales_sl_code=sales_sl_code,
           sales_executive=sales_executive,
           sales_area=sales_area,
           budget_value=budget_value,
           budget_quantity=budget_quantity,
           budget_product_group=budget_product_group,
           budget_sl_code=budget_sl_code,
           budget_executive=budget_executive,
           budget_area=budget_area,
           selected_executives=selected_executives,
           selected_months=selected_months,
           selected_branches=selected_branches
       )
       
       if result.get('success'):
           return jsonify(result)
       else:
           return jsonify({'error': result.get('error', 'Unknown error')}), 500
           
   except (ValueError, KeyError, FileNotFoundError) as e:
       logger.error(f"Error in calculate_budget_vs_billed route: {str(e)}", exc_info=True)
       return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/auto_map_columns', methods=['POST'])
def auto_map_columns():
   """Auto-map columns for executive analysis"""
   try:
       data = request.get_json()
       sales_file_path = data.get('sales_file_path')
       budget_file_path = data.get('budget_file_path')
       
       if not sales_file_path or not budget_file_path:
           return jsonify({'error': 'Sales and budget file paths are required'}), 400
       
       # Load data
       sales_df = load_data(sales_file_path)
       budget_df = load_data(budget_file_path)
       
       # Auto-map columns
       mapping_result = auto_map_executive_columns(sales_df, budget_df)
       
       return jsonify(mapping_result)
       
   except (ValueError, KeyError, FileNotFoundError) as e:
       logger.error(f"Error in auto_map_columns route: {str(e)}", exc_info=True)
       return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/get_exec_branch_options', methods=['POST'])
def get_exec_branch_options():
   """Get available executives and branches from uploaded files"""
   try:
       data = request.get_json()
       sales_file_path = data.get('sales_file_path')
       budget_file_path = data.get('budget_file_path')
       
       if not sales_file_path or not budget_file_path:
           return jsonify({'error': 'Sales and budget file paths are required'}), 400
       
       # Load data
       sales_df = load_data(sales_file_path)
       budget_df = load_data(budget_file_path)
       
       # Auto-map to find executive and branch columns
       mapping_result = auto_map_executive_columns(sales_df, budget_df)
       
       sales_exec_col = mapping_result['sales_mapping'].get('exec')
       budget_exec_col = mapping_result['budget_mapping'].get('exec')
       sales_area_col = mapping_result['sales_mapping'].get('area')
       budget_area_col = mapping_result['budget_mapping'].get('area')
       
       # Get executives and branches
       exec_branch_data = get_executives_and_branches(
           sales_df, budget_df, 
           sales_exec_col, budget_exec_col,
           sales_area_col, budget_area_col
       )
       
       # Get available months from sales data
       sales_date_col = mapping_result['sales_mapping'].get('date')
       available_months = get_available_months(sales_df, sales_date_col)
       
       return jsonify({
           **exec_branch_data,
           'available_months': available_months,
           'column_mapping': mapping_result
       })
       
   except (ValueError, KeyError, FileNotFoundError) as e:
       logger.error(f"Error in get_exec_branch_options route: {str(e)}", exc_info=True)
       return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/get_available_months', methods=['POST'])
def get_available_months_route():
   """Get available months from sales data"""
   try:
       data = request.get_json()
       sales_file_path = data.get('sales_file_path')
       
       if not sales_file_path:
           return jsonify({'error': 'Sales file path is required'}), 400
       
       # Load sales data
       sales_df = load_data(sales_file_path)
       
       # Auto-map to find date column
       mapping_result = auto_map_executive_columns(sales_df, sales_df)  # Use same df for both
       sales_date_col = mapping_result['sales_mapping'].get('date')
       
       available_months = get_available_months(sales_df, sales_date_col)
       
       return jsonify({
           'available_months': available_months,
           'date_column': sales_date_col
       })
       
   except (ValueError, KeyError, FileNotFoundError) as e:
       logger.error(f"Error in get_available_months route: {str(e)}", exc_info=True)
       return jsonify({'error': f'Server error: {str(e)}'}), 500

# ============= OD VS COLLECTION ROUTES (EXISTING) =============

@executive_bp.route('/od_auto_map_columns', methods=['POST'])
def od_auto_map_columns():
   """Auto-map columns for OD analysis"""
   try:
       data = request.get_json()
       os_jan_file_path = data.get('os_jan_file_path')
       os_feb_file_path = data.get('os_feb_file_path')
       sales_file_path = data.get('sales_file_path')
       
       if not os_jan_file_path or not os_feb_file_path or not sales_file_path:
           return jsonify({'error': 'All three file paths are required'}), 400
       
       # Load data
       os_jan_df = load_od_data(os_jan_file_path)
       os_feb_df = load_od_data(os_feb_file_path)
       sales_df = load_od_data(sales_file_path)
       
       # Auto-map columns
       mapping_result = auto_map_od_columns(os_jan_df, os_feb_df, sales_df)
       
       return jsonify(mapping_result)
       
   except (ValueError, KeyError, FileNotFoundError) as e:
       logger.error(f"Error in od_auto_map_columns route: {str(e)}", exc_info=True)
       return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/od_get_exec_branch_options', methods=['POST'])
def od_get_exec_branch_options():
   """Get available executives and branches from OD files"""
   try:
       data = request.get_json()
       os_jan_file_path = data.get('os_jan_file_path')
       os_feb_file_path = data.get('os_feb_file_path')
       sales_file_path = data.get('sales_file_path')
       
       if not os_jan_file_path or not os_feb_file_path or not sales_file_path:
           return jsonify({'error': 'All three file paths are required'}), 400
       
       # Load data
       os_jan_df = load_od_data(os_jan_file_path)
       os_feb_df = load_od_data(os_feb_file_path)
       sales_df = load_od_data(sales_file_path)
       
       # Auto-map to find executive and branch columns
       mapping_result = auto_map_od_columns(os_jan_df, os_feb_df, sales_df)
       
       os_jan_exec_col = mapping_result['os_jan_mapping'].get('executive')
       os_feb_exec_col = mapping_result['os_feb_mapping'].get('executive')
       sales_exec_col = mapping_result['sales_mapping'].get('executive')
       os_jan_area_col = mapping_result['os_jan_mapping'].get('area')
       os_feb_area_col = mapping_result['os_feb_mapping'].get('area')
       sales_area_col = mapping_result['sales_mapping'].get('area')
       
       # Get executives and branches
       exec_branch_data = get_od_executives_and_branches(
           os_jan_df, os_feb_df, sales_df,
           os_jan_exec_col, os_feb_exec_col, sales_exec_col,
           os_jan_area_col, os_feb_area_col, sales_area_col
       )
       
       return jsonify({
           **exec_branch_data,
           'column_mapping': mapping_result
       })
       
   except (ValueError, KeyError, FileNotFoundError) as e:
       logger.error(f"Error in od_get_exec_branch_options route: {str(e)}", exc_info=True)
       return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/od_get_available_months', methods=['POST'])
def od_get_available_months():
   """Get available months from OD data"""
   try:
       data = request.get_json()
       os_jan_file_path = data.get('os_jan_file_path')
       os_feb_file_path = data.get('os_feb_file_path')
       sales_file_path = data.get('sales_file_path')
       
       if not os_jan_file_path or not os_feb_file_path or not sales_file_path:
           return jsonify({'error': 'All three file paths are required'}), 400
       
       # Load data
       os_jan_df = load_od_data(os_jan_file_path)
       os_feb_df = load_od_data(os_feb_file_path)
       sales_df = load_od_data(sales_file_path)
       
       # Auto-map to find date columns
       mapping_result = auto_map_od_columns(os_jan_df, os_feb_df, sales_df)
       
       os_jan_due_date_col = mapping_result['os_jan_mapping'].get('due_date')
       os_jan_ref_date_col = mapping_result['os_jan_mapping'].get('ref_date')
       os_feb_due_date_col = mapping_result['os_feb_mapping'].get('due_date')
       os_feb_ref_date_col = mapping_result['os_feb_mapping'].get('ref_date')
       sales_bill_date_col = mapping_result['sales_mapping'].get('bill_date')
       sales_due_date_col = mapping_result['sales_mapping'].get('due_date')
       
       available_months = get_available_months_od(
           os_jan_df, os_feb_df, sales_df,
           os_jan_due_date_col, os_jan_ref_date_col,
           os_feb_due_date_col, os_feb_ref_date_col,
           sales_bill_date_col, sales_due_date_col
       )
       
       return jsonify({
           'available_months': available_months,
           'column_mapping': mapping_result
       })
       
   except (ValueError, KeyError, FileNotFoundError) as e:
       logger.error(f"Error in od_get_available_months route: {str(e)}", exc_info=True)
       return jsonify({'error': f'Server error: {str(e)}'}), 500

@executive_bp.route('/calculate_od_vs_collection', methods=['POST'])
def calculate_od_vs_collection():
   """Calculate OD Target vs Collection analysis"""
   try:
       data = request.get_json()
       logger.info(f"Received OD vs Collection request data keys: {list(data.keys())}")
       
       # Get file paths
       os_jan_file_path = data.get('os_jan_file_path')
       os_feb_file_path = data.get('os_feb_file_path')
       sales_file_path = data.get('sales_file_path')
       
       if not os_jan_file_path or not os_feb_file_path or not sales_file_path:
           return jsonify({'error': 'All three file paths are required'}), 400
       
       # Get OS Jan column mappings
       os_jan_due_date = data.get('os_jan_due_date')
       os_jan_ref_date = data.get('os_jan_ref_date')
       os_jan_net_value = data.get('os_jan_net_value')
       os_jan_executive = data.get('os_jan_executive')
       os_jan_sl_code = data.get('os_jan_sl_code')
       os_jan_area = data.get('os_jan_area')
       
       # Get OS Feb column mappings
       os_feb_due_date = data.get('os_feb_due_date')
       os_feb_ref_date = data.get('os_feb_ref_date')
       os_feb_net_value = data.get('os_feb_net_value')
       os_feb_executive = data.get('os_feb_executive')
       os_feb_sl_code = data.get('os_feb_sl_code')
       os_feb_area = data.get('os_feb_area')
       
       # Get Sales column mappings
       sales_bill_date = data.get('sales_bill_date')
       sales_due_date = data.get('sales_due_date')
       sales_value = data.get('sales_value')
       sales_executive = data.get('sales_executive')
       sales_sl_code = data.get('sales_sl_code')
       sales_area = data.get('sales_area')
       
       # Get filters
       selected_executives = data.get('selected_executives', [])
       selected_month = data.get('selected_month')
       selected_branches = data.get('selected_branches', [])
       
       logger.info(
           f"OS Jan columns: due_date={os_jan_due_date}, ref_date={os_jan_ref_date}, "
           f"net_value={os_jan_net_value}, exec={os_jan_executive}, area={os_jan_area}"
       )
       logger.info(
           f"OS Feb columns: due_date={os_feb_due_date}, ref_date={os_feb_ref_date}, "
           f"net_value={os_feb_net_value}, exec={os_feb_executive}, area={os_feb_area}"
       )
       logger.info(
           f"Sales columns: bill_date={sales_bill_date}, due_date={sales_due_date}, "
           f"value={sales_value}, exec={sales_executive}, area={sales_area}"
       )
       logger.info(f"Selected executives: {selected_executives}")
       logger.info(f"Selected month: {selected_month}")
       logger.info(f"Selected branches: {selected_branches}")
       
       # Validate required mappings
       required_os_jan_cols = [os_jan_due_date, os_jan_ref_date, os_jan_net_value, os_jan_executive, os_jan_sl_code, os_jan_area]
       required_os_feb_cols = [os_feb_due_date, os_feb_ref_date, os_feb_net_value, os_feb_executive, os_feb_sl_code, os_feb_area]
       required_sales_cols = [sales_bill_date, sales_due_date, sales_value, sales_executive, sales_sl_code, sales_area]
       
       if not all(required_os_jan_cols):
           missing = [name for name, val in zip(['due_date', 'ref_date', 'net_value', 'executive', 'sl_code', 'area'], required_os_jan_cols) if not val]
           return jsonify({'error': f'Missing required OS Jan column mappings: {missing}'}), 400
           
       if not all(required_os_feb_cols):
           missing = [name for name, val in zip(['due_date', 'ref_date', 'net_value', 'executive', 'sl_code', 'area'], required_os_feb_cols) if not val]
           return jsonify({'error': f'Missing required OS Feb column mappings: {missing}'}), 400
           
       if not all(required_sales_cols):
           missing = [name for name, val in zip(['bill_date', 'due_date', 'value', 'executive', 'sl_code', 'area'], required_sales_cols) if not val]
           return jsonify({'error': f'Missing required Sales column mappings: {missing}'}), 400
       
       if not selected_month:
           return jsonify({'error': 'Selected month is required'}), 400
       
       # Load data
       os_jan_df = load_od_data(os_jan_file_path)
       os_feb_df = load_od_data(os_feb_file_path)
       sales_df = load_od_data(sales_file_path)
       
       logger.info(f"OS Jan data loaded: {len(os_jan_df)} rows")
       logger.info(f"OS Feb data loaded: {len(os_feb_df)} rows")
       logger.info(f"Sales data loaded: {len(sales_df)} rows")
       
       # Call calculation function
       result = calculate_od_values(
           os_jan_df, os_feb_df, sales_df, selected_month,
           os_jan_due_date, os_jan_ref_date, os_jan_net_value, os_jan_executive, os_jan_sl_code, os_jan_area,
           os_feb_due_date, os_feb_ref_date, os_feb_net_value, os_feb_executive, os_feb_sl_code, os_feb_area,
           sales_bill_date, sales_due_date, sales_value, sales_executive, sales_sl_code, sales_area,
           selected_executives, selected_branches
       )
       
       if result.get('success'):
           return jsonify(result)
       else:
           return jsonify({'error': result.get('error', 'Unknown error')}), 500
           
   except (ValueError, KeyError, FileNotFoundError) as e:
       logger.error(f"Error in calculate_od_vs_collection route: {str(e)}", exc_info=True)
       return jsonify({'error': f'Server error: {str(e)}'}), 500

# ============= PPT GENERATION ROUTES (EXISTING) =============

@executive_bp.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    """Generate PowerPoint presentation for executive budget vs billed results"""
    try:
        data = request.get_json()
        logger.info(f"Received PPT generation request with keys: {list(data.keys())}")

        # Get results data
        results_data = data.get('results_data')
        if not results_data:
            return jsonify({'error': 'Results data is required for PPT generation'}), 400

        # Get optional parameters
        month_title = data.get('month_title', 'Executive Analysis')

        # 🔥 Load static logo
        import os
        logo_file = None
        static_logo_path = os.path.join("static", "logo.jpeg")
        if os.path.exists(static_logo_path):
            logo_file = open(static_logo_path, "rb")
            logger.info(f"✅ Static logo loaded for Executive PPT from: {static_logo_path}")
        else:
            logger.warning("⚠️ Static logo not found, continuing without logo.")

        logger.info(f"Generating PPT for month: {month_title}")

        # Validate data structure
        try:
            validate_ppt_data(results_data)
        except ValueError as e:
            return jsonify({'error': f'Invalid data structure: {str(e)}'}), 400

        # Generate PPT
        ppt_buffer = create_executive_budget_ppt(
            results_data=results_data,
            month_title=month_title,
            logo_file=logo_file
        )

        if logo_file:
            logo_file.close()  # 🧼 Clean up

        if not ppt_buffer:
            return jsonify({'error': 'Failed to generate PowerPoint presentation'}), 500

        # Generate unique filename
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"Executive_Budget_vs_Billed_{month_title.replace(' ', '_')}_{timestamp}.pptx"

        return send_file(
            ppt_buffer,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in generate_ppt route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500


@executive_bp.route('/generate_od_ppt', methods=['POST'])
def generate_od_ppt():
    """Generate PowerPoint presentation for OD vs Collection results"""
    try:
        data = request.get_json()
        logger.info(f"Received OD PPT generation request with keys: {list(data.keys())}")

        # Get results data
        results_data = data.get('results_data')
        if not results_data:
            return jsonify({'error': 'Results data is required for PPT generation'}), 400

        # Get optional parameters
        month_title = data.get('month_title', 'OD Analysis')

        # 🔥 Load static logo
        import os
        logo_file = None
        static_logo_path = os.path.join("static", "logo.jpeg")
        if os.path.exists(static_logo_path):
            logo_file = open(static_logo_path, "rb")
            logger.info(f"✅ Static logo loaded for OD PPT from: {static_logo_path}")
        else:
            logger.warning("⚠️ Static logo not found, proceeding without logo.")

        logger.info(f"Generating OD PPT for month: {month_title}")

        # Validate data structure
        try:
            validate_ppt_data(results_data)
        except ValueError as e:
            return jsonify({'error': f'Invalid data structure: {str(e)}'}), 400

        # Generate PPT
        ppt_buffer = create_executive_od_ppt(
            results_data=results_data,
            month_title=month_title,
            logo_file=logo_file
        )

        if logo_file:
            logo_file.close()  # 🧼 Clean up file handle

        if not ppt_buffer:
            return jsonify({'error': 'Failed to generate PowerPoint presentation'}), 500

        # Generate unique filename
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"Executive_OD_Target_vs_Collection_{month_title.replace(' ', '_')}_{timestamp}.pptx"

        return send_file(
            ppt_buffer,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in generate_od_ppt route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500


@executive_bp.route('/generate_consolidated_ppt', methods=['POST'])
def generate_consolidated_ppt():
    """Generate consolidated PowerPoint presentation with multiple reports"""
    try:
        data = request.get_json()
        logger.info(f"Received consolidated PPT generation request")

        # 1. Validate reports
        reports_data = data.get('reports_data', [])
        if not reports_data:
            return jsonify({'error': 'Reports data is required for consolidated PPT generation'}), 400

        for report in reports_data:
            df_data = report.get("df", [])
            columns = report.get("columns", [])
            df = pd.DataFrame(df_data)

            if columns:
                df.attrs["columns"] = columns
                logger.info(f"✅ Embedded column order for: {report.get('title')} → {columns}")
            else:
                logger.warning(f"⚠️ No columns passed for: {report.get('title')}")

            report["df"] = df
            report.pop("columns", None)

        # 2. Load default logo from static folder (if it exists)
        logo_file = None
        static_logo_path = os.path.join("static", "logo.jpeg")  # or logo.jpg/jpeg
        if os.path.exists(static_logo_path):
            logo_file = open(static_logo_path, "rb")
            logger.info(f"✅ Static logo loaded from: {static_logo_path}")
        else:
            logger.warning("⚠️ Static logo not found. Proceeding without logo.")

        # 3. Generate PPT
        title = data.get('title', 'Consolidated Executive Report')
        ppt_buffer = create_consolidated_ppt(
            dfs_info=reports_data,
            logo_file=logo_file,
            title=title
        )

        if logo_file:
            logo_file.close()  # 🔐 Always close manually opened files

        if not ppt_buffer:
            return jsonify({'error': 'Failed to generate consolidated PowerPoint presentation'}), 500

        filename = f"Consolidated_Executive_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

        return send_file(
            ppt_buffer,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except (ValueError, KeyError, FileNotFoundError) as e:
        logger.error(f"Error in generate_consolidated_ppt route: {str(e)}", exc_info=True)
        return jsonify({'error': f'Server error: {str(e)}'}), 500

# ============= HEALTH CHECK ROUTE =============

@executive_bp.route('/health', methods=['GET'])
def health_check():
   """Health check endpoint"""
   return jsonify({
       'status': 'healthy', 
       'timestamp': datetime.now().isoformat(),
       'routes': [
           # Customer Analysis Routes
           'customer_auto_map_columns',
           'customer_get_options', 
           'calculate_customer_analysis',
           'generate_customer_ppt',
           # OD Target Analysis Routes
           'od_target_auto_map_columns',
           'od_target_get_options',
           'calculate_od_target',
           'generate_od_target_ppt',
           # Product Growth Routes
           'product_auto_map_columns',
           'product_get_options', 
           'calculate_product_growth',
           'generate_product_growth_ppt',
           # Budget vs Billed Routes
           'calculate_budget_vs_billed',
           'auto_map_columns',
           'get_exec_branch_options',
           'get_available_months',
           # OD vs Collection Routes
           'od_auto_map_columns',
           'od_get_exec_branch_options', 
           'od_get_available_months',
           'calculate_od_vs_collection',
           # PPT Generation Routes
           'generate_ppt',
           'generate_od_ppt',
           'generate_consolidated_ppt'
       ]
   })