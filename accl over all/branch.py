#good
import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
from matplotlib.table import Table
from datetime import datetime
from dateutil.relativedelta import relativedelta
import logging
import uuid
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)
st.set_page_config(page_title="ACCLLP Integrated Dashboard", page_icon="📊", layout="wide")
if 'init' not in st.session_state:
    st.session_state.init = True
    st.session_state.sales_file = None
    st.session_state.budget_file = None
    st.session_state.last_year_sales_file = None 
    st.session_state.os_jan_file = None
    st.session_state.os_feb_file = None
    st.session_state.logo_file = None
    st.session_state.budget_results = None
    st.session_state.customers_results = None
    st.session_state.od_results = None
    st.session_state.product_results = None
branch_mapping = {
    'PONDY': 'PUDUCHERRY', 'PDY': 'PUDUCHERRY', 'Puducherry - PDY': 'PUDUCHERRY',
    'COVAI': 'COIMBATORE', 'CBE': 'COIMBATORE', 'Coimbatore - CBE': 'COIMBATORE',
    'ERD': 'ERODE', 'Erode - ERD': 'ERODE', 'ERD002': 'ERODE', 'ERD001': 'ERODE',
    'ERDTD1': 'ERODE', 'ERD003': 'ERODE', 'ERD004': 'ERODE', 'ERD005': 'ERODE',
    'ERD007': 'ERODE',
    'KRR': 'KARUR',
    'Chennai - CHN': 'CHENNAI', 'CHN': 'CHENNAI',
    'Tirupur - TPR': 'TIRUPUR', 'TPR': 'TIRUPUR',
    'Madurai - MDU': 'MADURAI', 'MDU': 'MADURAI',
    'POULTRY': 'POULTRY', 'Poultry - PLT': 'POULTRY',
    'SALEM': 'SALEM', 'Salem - SLM': 'SALEM',
    'HO': 'HO',
    'SLM002': 'SALEM', 'SLMTD1': 'SALEM',
    'BHV1': 'BHAVANI',
    'CBU': 'COIMBATORE',
    'VLR': 'VELLORE', 
    'TRZ': 'TRICHY', 
    'TVL': 'TIRUNELVELI',
    'NGS': 'NAGERCOIL', 
    'PONDICHERRY': 'PUDUCHERRY',
    'BLR': 'BANGALORE', 'BANGALORE': 'BANGALORE', 'BGLR': 'BANGALORE'
}
nbc_branch_mapping = {
    'PONDY': 'PONDY', 'PDY': 'PONDY',
    'COVAI': 'ERODE_CBE_KRR',
    'ERODE': 'ERODE_CBE_KRR',
    'KARUR': 'ERODE_CBE_KRR',
    'SLM002': 'SALEM', 'SLMTD1': 'SALEM',
    'BHV1': 'BHAVANI',
    'MDU': 'MADURAI',
    'CHENNAI': 'CHENNAI', 'CHN': 'CHENNAI',
    'TIRUPUR': 'TIRUPUR', 'TPR': 'TIRUPUR',
    'POULTRY': 'POULTRY'
}
def find_column_by_names(columns, target_names):
    """Find a column by checking multiple possible names (case-insensitive)."""
    columns_upper = [col.upper() for col in columns]
    target_names_upper = [name.upper() for name in target_names]
    
    for target in target_names_upper:
        if target in columns_upper:
            return columns[columns_upper.index(target)]
    return None

def auto_map_budget_columns(sales_columns, budget_columns):
    """Auto-map Budget vs Billed columns based on specifications."""
    # Sales column mappings
    sales_mapping = {
        'date': find_column_by_names(sales_columns, ['Date', 'Bill Date', 'Invoice Date']),
        'value': find_column_by_names(sales_columns, ['Value', 'Invoice Value', 'Amount']),
        'product_group': find_column_by_names(sales_columns, ['Type (Make)', 'Product Group', 'Type', 'Make']),
        'area': find_column_by_names(sales_columns, ['Branch', 'Area', 'Location']),
        'quantity': find_column_by_names(sales_columns, ['Actual Quantity', 'Quantity', 'Qty']),
        'sl_code': find_column_by_names(sales_columns, ['Customer Code', 'SL Code', 'Customer ID']),
        'executive': find_column_by_names(sales_columns, ['Executive', 'Sales Executive', 'Executive Name'])
    }
    
    # Budget column mappings
    budget_mapping = {
        'area': find_column_by_names(budget_columns, ['Branch', 'Area', 'Location']),
        'quantity': find_column_by_names(budget_columns, ["Qty – Apr'25", 'Quantity', 'Qty', 'Budget Qty']),
        'sl_code': find_column_by_names(budget_columns, ['SL Code', 'Customer Code', 'Customer ID']),
        'value': find_column_by_names(budget_columns, ["Value – Apr'25", 'Value', 'Budget Value', 'Amount']),
        'product_group': find_column_by_names(budget_columns, ['Product Group', 'Type(Make)', 'Product', 'Type']),
        'executive': find_column_by_names(budget_columns, ['Executive Name', 'Executive', 'Sales Executive'])
    }
    
    return sales_mapping, budget_mapping
def auto_map_od_columns(os_first_columns, os_second_columns, sales_columns):
    """Auto-map OD Target vs Collection columns."""
    os_first_mapping = {
        'due_date': find_column_by_names(os_first_columns, ['Due Date', 'Due_Date', 'DueDate']),
        'branch': find_column_by_names(os_first_columns, ['Branch', 'Unit', 'Area', 'Location']),
        'ref_date': find_column_by_names(os_first_columns, ['Ref. Date', 'Ref Date', 'Reference Date']),
        'net_value': find_column_by_names(os_first_columns, ['Net Value', 'NetValue', 'Amount', 'Value']),
        'executive': find_column_by_names(os_first_columns, ['Executive Name', 'Executive', 'Sales Executive']),
        'region': find_column_by_names(os_first_columns, ['Region', 'Area Region', 'Zone'])
    }
    os_second_mapping = {
        'due_date': find_column_by_names(os_second_columns, ['Due Date', 'Due_Date', 'DueDate']),
        'branch': find_column_by_names(os_second_columns, ['Branch', 'Unit', 'Area', 'Location']),
        'ref_date': find_column_by_names(os_second_columns, ['Ref. Date', 'Ref Date', 'Reference Date']),
        'net_value': find_column_by_names(os_second_columns, ['Net Value', 'NetValue', 'Amount', 'Value']),
        'executive': find_column_by_names(os_second_columns, ['Executive Name', 'Executive', 'Sales Executive']),
        'region': find_column_by_names(os_second_columns, ['Region', 'Area Region', 'Zone'])
    }
    sales_mapping = {
        'bill_date': find_column_by_names(sales_columns, ['Date', 'Bill Date', 'Invoice Date']),
        'branch': find_column_by_names(sales_columns, ['Branch', 'Unit', 'Area', 'Location']),
        'due_date': find_column_by_names(sales_columns, ['Due Date', 'Due_Date', 'DueDate']),
        'value': find_column_by_names(sales_columns, ['Invoice Value', 'Value', 'Amount']),
        'executive': find_column_by_names(sales_columns, ['Executive', 'Sales Executive', 'Executive Name']),
        'region': find_column_by_names(sales_columns, ['Region', 'Area Region', 'Zone'])
    }
    
    return os_first_mapping, os_second_mapping, sales_mapping
def auto_map_product_growth_columns(ly_columns, cy_columns, budget_columns):
    """Auto-map Product Growth columns."""
    ly_mapping = {
        'date': find_column_by_names(ly_columns, ['Date', 'Bill Date', 'Invoice Date']),
        'product_group': find_column_by_names(ly_columns, ['Type (Make)', 'Product Group', 'Type', 'Make']),
        'quantity': find_column_by_names(ly_columns, ['Actual Quantity', 'Quantity', 'Qty']),
        'company_group': find_column_by_names(ly_columns, ['Company Group', 'Company', 'Group']),
        'value': find_column_by_names(ly_columns, ['Value', 'Invoice Value', 'Amount']),
        'executive': find_column_by_names(ly_columns, ['Executive', 'Sales Executive', 'Executive Name'])
    }
    cy_mapping = {
        'date': find_column_by_names(cy_columns, ['Month Format', 'Date', 'Bill Date', 'Invoice Date']),
        'product_group': find_column_by_names(cy_columns, ['Type (Make)', 'Product Group', 'Type', 'Make']),
        'quantity': find_column_by_names(cy_columns, ['Actual Quantity', 'Quantity', 'Qty']),
        'company_group': find_column_by_names(cy_columns, ['Company Group', 'Company', 'Group']),
        'value': find_column_by_names(cy_columns, ['Value', 'Invoice Value', 'Amount']),
        'executive': find_column_by_names(cy_columns, ['Executive', 'Sales Executive', 'Executive Name'])
    }
    budget_mapping = {
        'quantity': find_column_by_names(budget_columns, ["Qty – Apr'25", 'Quantity', 'Qty', 'Budget Qty']),
        'company_group': find_column_by_names(budget_columns, ['Company Group', 'Company', 'Group']),
        'value': find_column_by_names(budget_columns, ["Value – Apr'25", 'Value', 'Budget Value', 'Amount']),
        'executive': find_column_by_names(budget_columns, ['Executive Name', 'Executive', 'Sales Executive']),
        'product_group': find_column_by_names(budget_columns, ['Product Group', 'Type(Make)', 'Product', 'Type'])
    }
    
    return ly_mapping, cy_mapping, budget_mapping
def auto_map_nbc_columns(sales_columns):
    """Auto-map Number of Billed Customers columns."""
    return {
        'date': find_column_by_names(sales_columns, ['Date', 'Bill Date', 'Invoice Date']),
        'customer_id': find_column_by_names(sales_columns, ['Customer Code', 'SL Code', 'Customer ID']),
        'branch': find_column_by_names(sales_columns, ['Branch', 'Area', 'Location']),
        'executive': find_column_by_names(sales_columns, ['Executive', 'Sales Executive', 'Executive Name'])
    }
def auto_map_od_target_columns(os_columns):
    """Auto-map OD Target columns."""
    return {
        'area': find_column_by_names(os_columns, ['Branch', 'Area', 'Location', 'Unit']),
        'due_date': find_column_by_names(os_columns, ['Due Date', 'Due_Date', 'DueDate']),
        'net_value': find_column_by_names(os_columns, ['Net Value', 'NetValue', 'Amount', 'Value']),
        'executive': find_column_by_names(os_columns, ['Executive Name', 'Executive', 'Sales Executive'])
    }
def extract_region_from_branch(branch):
    """Extract region from branch name based on location groupings."""
    if pd.isna(branch):
        return 'Unknown'
    
    branch = str(branch).upper().strip()
    region_map = {
        'SOUTH': ['MADURAI', 'TIRUNELVELI', 'NAGERCOIL', 'TRICHY'],
        'WEST': ['COIMBATORE', 'ERODE', 'TIRUPUR', 'KARUR', 'SALEM', 'BHAVANI'],
        'EAST': ['CHENNAI', 'VELLORE'],
        'CENTRAL': ['PUDUCHERRY', 'PONDICHERRY', 'PONDY'],
        'EXTERNAL': ['BANGALORE', 'BLR', 'BGLR'],
        'SPECIAL': ['POULTRY', 'HO']
    }
    normalized_branch = map_branch(branch, case='upper')
    for region, branches in region_map.items():
        if normalized_branch in branches:
            return region
    
    return 'Others'
def get_excel_sheets(file):
    """Get sheet names from an Excel file."""
    try:
        xl = pd.ExcelFile(file)
        return xl.sheet_names
    except Exception as e:
        logger.error(f"Error reading Excel sheets: {e}")
        st.error(f"Error reading Excel sheets: {e}")
        return []
def map_branch(branch_name, case='upper'):
    """Map branch codes/names to standardized names."""
    if pd.isna(branch_name):
        return 'Unknown'
    branch_str = str(branch_name).strip().upper()
    if ' - ' in branch_str:
        branch_str = branch_str.split(' - ')[-1].strip()
    mapped_branch = branch_mapping.get(branch_str, branch_str)
    if case == 'title':
        return mapped_branch.title()
    else:  # default: upper
        return mapped_branch
def create_table_image(df, title, percent_cols=None):
    """Create a table image for display in Streamlit."""
    if df is None or df.empty:
        return None    
    fig, ax = plt.subplots(figsize=(10, len(df) * 0.5))
    ax.axis('off')
    table = Table(ax, bbox=[0, 0, 1, 1])
    nrows, ncols = df.shape
    width, height = 1.0 / ncols, 1.0 / nrows
    for col_idx, col_name in enumerate(df.columns):
        table.add_cell(0, col_idx, width, height, text=str(col_name), loc='center', facecolor='#0070C0')
        table[0, col_idx].set_text_props(weight='bold', color='white')
    for row_idx in range(nrows):
        for col_idx in range(ncols):
            value = df.iloc[row_idx, col_idx]
            if percent_cols and col_idx in percent_cols:
                text = f"{value}%"
            else:
                text = str(value)            
            facecolor = '#f0f0f0' if row_idx % 2 == 0 else 'white'
            is_total = row_idx == nrows - 1 and str(df.iloc[row_idx, 0]).upper() in ['TOTAL', 'GRAND TOTAL']
            if is_total:
                facecolor = '#D3D3D3'
                table.add_cell(row_idx + 1, col_idx, width, height, text=text, loc='center', facecolor=facecolor).set_text_props(weight='bold')
            else:
                table.add_cell(row_idx + 1, col_idx, width, height, text=text, loc='center', facecolor=facecolor)
    
    table.auto_set_font_size(False)
    table.set_fontsize(12)
    ax.add_table(table)
    plt.title(title, fontsize=14 if len(title) > 50 else 16, weight='bold', color='#0070C0', pad=20)
    plt.tight_layout()
    img_buffer = BytesIO()
    plt.savefig(img_buffer, format='png', bbox_inches='tight', dpi=100)
    img_buffer.seek(0)
    plt.close()
    return img_buffer
def create_title_slide(prs, title, logo_file=None):
    """Create a standard title slide for PPT."""
    blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout
    title_slide = prs.slides.add_slide(blank_slide_layout)
    company_name = title_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    company_frame = company_name.text_frame
    company_frame.text = "Asia Crystal Commodity LLP"
    p = company_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    if logo_file is not None:
        try:
            logo_buffer = BytesIO(logo_file.read())
            logo = title_slide.shapes.add_picture(logo_buffer, Inches(5.665), Inches(1.5), width=Inches(2), height=Inches(2))
            # Reset file pointer for future use
            logo_file.seek(0)
        except Exception as e:
            logger.error(f"Error adding logo to slide: {e}")
    title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    subtitle = title_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.33), Inches(1))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "ACCLLP"
    p = subtitle_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)    
    return title_slide
def add_table_slide(prs, df, title, percent_cols=None):
    """Add a slide with a table to a PPT presentation."""
    if df is None or df.empty:
        return None
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.75))
    title_frame = title_shape.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5)).table
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    for row_idx in range(df.shape[0]):
        is_total_row = row_idx == rows - 1 and str(df.iloc[row_idx, 0]).upper() in ['TOTAL', 'GRAND TOTAL']
        for col_idx in range(cols):
            cell = table.cell(row_idx + 1, col_idx)
            value = df.iloc[row_idx, col_idx]
            if percent_cols and col_idx in percent_cols:
                cell.text = f"{value}%"
            else:
                cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            if row_idx % 2 == 0 and not is_total_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            if is_total_row:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(211, 211, 211)    
    return slide
def create_consolidated_ppt(all_dfs_with_titles, logo_file=None, title="ACCLLP Consolidated Report"):
    """Create a consolidated PPT with all report data."""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        create_title_slide(prs, title, logo_file)
        for df_info in all_dfs_with_titles:
            if df_info and 'df' in df_info and 'title' in df_info:
                add_table_slide(
                    prs, 
                    df_info['df'], 
                    df_info['title'],
                    df_info.get('percent_cols')
                )
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating consolidated PPT: {e}")
        st.error(f"Error creating consolidated PPT: {e}")
        return None
def calculate_values(sales_df, budget_df, selected_month, sales_executives, budget_executives,
                     sales_date_col, sales_area_col, sales_value_col, sales_qty_col, sales_product_group_col, sales_sl_code_col, sales_exec_col,
                     budget_area_col, budget_value_col, budget_qty_col, budget_product_group_col, budget_sl_code_col, budget_exec_col, selected_branches=None):
    """
    Calculate budget vs billed values with improved optimized logic.
    Returns four DataFrames for the different report sections.
    """
    try:
        # Create copies to avoid modifying original DataFrames
        sales_df = sales_df.copy()
        budget_df = budget_df.copy()
        
        # Validate column existence
        required_sales_cols = [sales_date_col, sales_area_col, sales_value_col, sales_qty_col, sales_exec_col,
                              sales_product_group_col, sales_sl_code_col]
        required_budget_cols = [budget_area_col, budget_value_col, budget_qty_col, budget_exec_col,
                               budget_product_group_col, budget_sl_code_col]
        
        missing_sales_cols = [col for col in required_sales_cols if col not in sales_df.columns]
        missing_budget_cols = [col for col in required_budget_cols if col not in budget_df.columns]
        
        if missing_sales_cols:
            st.error(f"Missing columns in sales data: {missing_sales_cols}")
            return None, None, None, None
        if missing_budget_cols:
            st.error(f"Missing columns in budget data: {missing_budget_cols}")
            return None, None, None, None

        # Get branch information for UI
        raw_sales_branches = sales_df[sales_area_col].dropna().astype(str).str.upper().str.split(' - ').str[-1].unique().tolist()
        raw_budget_branches = budget_df[budget_area_col].dropna().astype(str).str.upper().str.split(' - ').str[-1].str.replace('AAAA - ', '', regex=False).unique().tolist()
        all_branches = sorted(set([branch_mapping.get(branch, branch) for branch in raw_sales_branches + raw_budget_branches]))
        
        # Filter by executives
        if sales_executives:
            sales_df = sales_df[sales_df[sales_exec_col].isin(sales_executives)].copy()
        if budget_executives:
            budget_df = budget_df[budget_df[budget_exec_col].isin(budget_executives)].copy()

        if sales_df.empty or budget_df.empty:
            st.warning("No data found for selected executives in one or both files.")
            return None, None, None, None
            
        # Convert date and numeric columns
        sales_df[sales_date_col] = pd.to_datetime(sales_df[sales_date_col], dayfirst=True, errors='coerce')
        sales_df[sales_value_col] = pd.to_numeric(sales_df[sales_value_col], errors='coerce').fillna(0)
        sales_df[sales_qty_col] = pd.to_numeric(sales_df[sales_qty_col], errors='coerce').fillna(0)
        budget_df[budget_value_col] = pd.to_numeric(budget_df[budget_value_col], errors='coerce').fillna(0)
        budget_df[budget_qty_col] = pd.to_numeric(budget_df[budget_qty_col], errors='coerce').fillna(0)
        
        # Filter sales data for the selected month
        filtered_sales_df = sales_df[sales_df[sales_date_col].dt.strftime('%b %y') == selected_month].copy()
        if filtered_sales_df.empty:
            st.warning(f"No sales data found for {selected_month} in '{sales_date_col}'. Check date format.")
            return None, None, None, None

        # Standardize string columns
        filtered_sales_df[sales_area_col] = filtered_sales_df[sales_area_col].astype(str).str.strip()
        budget_df[budget_area_col] = budget_df[budget_area_col].astype(str).str.strip()
        filtered_sales_df[sales_product_group_col] = filtered_sales_df[sales_product_group_col].astype(str).str.strip()
        filtered_sales_df[sales_sl_code_col] = filtered_sales_df[sales_sl_code_col].astype(str).str.strip().str.replace('\.0$', '', regex=True)
        budget_df[budget_product_group_col] = budget_df[budget_product_group_col].astype(str).str.strip()
        budget_df[budget_sl_code_col] = budget_df[budget_sl_code_col].astype(str).str.strip().str.replace('\.0$', '', regex=True)

        # Apply branch mapping and standardization
        budget_df[budget_area_col] = budget_df[budget_area_col].str.split(' - ').str[-1].str.upper()
        budget_df[budget_area_col] = budget_df[budget_area_col].str.replace('AAAA - ', '', regex=False).str.upper()
        budget_df[budget_area_col] = budget_df[budget_area_col].replace(branch_mapping)
        filtered_sales_df[sales_area_col] = filtered_sales_df[sales_area_col].str.upper().replace(branch_mapping)
        
        # Apply branch filtering if specified
        if selected_branches:
            filtered_sales_df = filtered_sales_df[filtered_sales_df[sales_area_col].isin(selected_branches)].copy()
            budget_df = budget_df[budget_df[budget_area_col].isin(selected_branches)].copy()
            default_branches = selected_branches
        else:
            default_branches = all_branches

        # Standardize product groups and SL codes
        filtered_sales_df[sales_product_group_col] = filtered_sales_df[sales_product_group_col].str.upper()
        filtered_sales_df[sales_sl_code_col] = filtered_sales_df[sales_sl_code_col].str.upper()
        budget_df[budget_product_group_col] = budget_df[budget_product_group_col].str.upper()
        budget_df[budget_sl_code_col] = budget_df[budget_sl_code_col].str.upper()

        # STEP 1: Process Budget Data with improved logic
        print("Processing Budget Data...")
        
        # Group by Branch + SL Code + Product Group and sum quantities/values
        budget_grouped = budget_df.groupby([
            budget_area_col,
            budget_sl_code_col, 
            budget_product_group_col
        ]).agg({
            budget_qty_col: 'sum',
            budget_value_col: 'sum'
        }).reset_index()
        
        # Filter: Only include rows where BOTH qty > 0 AND value > 0
        budget_valid = budget_grouped[
            (budget_grouped[budget_qty_col] > 0) & 
            (budget_grouped[budget_value_col] > 0)
        ].copy()
        
        if budget_valid.empty:
            st.error("No valid budget data found (with qty > 0 and value > 0).")
            return None, None, None, None
        
        print(f"Valid budget records: {len(budget_valid)}")
        
        # STEP 2: Process Sales Data - Match exactly with budget combinations
        print("Processing Sales Data...")
        
        # Initialize results storage
        final_results = []
        
        # For each valid budget record, find matching sales
        for _, budget_row in budget_valid.iterrows():
            branch = budget_row[budget_area_col]
            sl_code = budget_row[budget_sl_code_col]
            product = budget_row[budget_product_group_col]
            budget_qty = budget_row[budget_qty_col]
            budget_value = budget_row[budget_value_col]
            
            # Find matching sales records (same branch + sl_code + product)
            matching_sales = filtered_sales_df[
                (filtered_sales_df[sales_area_col] == branch) &
                (filtered_sales_df[sales_sl_code_col] == sl_code) &
                (filtered_sales_df[sales_product_group_col] == product)
            ]
            
            # Sum all matching sales records for this combination
            if not matching_sales.empty:
                sales_qty_total = matching_sales[sales_qty_col].sum()
                sales_value_total = matching_sales[sales_value_col].sum()
            else:
                sales_qty_total = 0
                sales_value_total = 0
            
            # Apply the comparison logic for budget vs billed reports
            # If sales > budget, use budget; else use sales
            final_qty = budget_qty if sales_qty_total > budget_qty else sales_qty_total
            final_value = budget_value if sales_value_total > budget_value else sales_value_total
            
            # Store result
            final_results.append({
                'Branch': branch,
                'SL_Code': sl_code,
                'Product': product,
                'Budget_Qty': budget_qty,
                'Sales_Qty': sales_qty_total,
                'Final_Qty': final_qty,
                'Budget_Value': budget_value,
                'Sales_Value': sales_value_total,
                'Final_Value': final_value
            })
        
        # Convert to DataFrame for easier manipulation
        results_df = pd.DataFrame(final_results)
        
        print(f"Processed {len(results_df)} budget-sales combinations")
        
        # STEP 3: Aggregate by Branch for Budget vs Billed reports
        
        # Budget vs Billed Quantity DataFrame
        branch_qty_summary = results_df.groupby('Branch').agg({
            'Budget_Qty': 'sum',
            'Final_Qty': 'sum'
        }).reset_index()
        branch_qty_summary.columns = ['Area', 'Budget Qty', 'Billed Qty']
        
        # Calculate achievement percentage for quantity
        branch_qty_summary['%'] = branch_qty_summary.apply(
            lambda row: int((row['Billed Qty'] / row['Budget Qty'] * 100)) if row['Budget Qty'] > 0 else 0,
            axis=1
        )
        
        # Ensure all selected branches are included (even if no data)
        budget_vs_billed_qty_df = pd.DataFrame({'Area': default_branches})
        budget_vs_billed_qty_df = pd.merge(
            budget_vs_billed_qty_df,
            branch_qty_summary,
            on='Area',
            how='left'
        ).fillna({'Budget Qty': 0, 'Billed Qty': 0, '%': 0})
        
        # Budget vs Billed Value DataFrame
        branch_value_summary = results_df.groupby('Branch').agg({
            'Budget_Value': 'sum',
            'Final_Value': 'sum'
        }).reset_index()
        branch_value_summary.columns = ['Area', 'Budget Value', 'Billed Value']
        
        # Calculate achievement percentage for value
        branch_value_summary['%'] = branch_value_summary.apply(
            lambda row: int((row['Billed Value'] / row['Budget Value'] * 100)) if row['Budget Value'] > 0 else 0,
            axis=1
        )
        
        # Ensure all selected branches are included (even if no data)
        budget_vs_billed_value_df = pd.DataFrame({'Area': default_branches})
        budget_vs_billed_value_df = pd.merge(
            budget_vs_billed_value_df,
            branch_value_summary,
            on='Area',
            how='left'
        ).fillna({'Budget Value': 0, 'Billed Value': 0, '%': 0})
        
        # STEP 4: Create Overall Sales DataFrames (use all sales data, not just matched)
        
        # Overall Sales - use complete sales data for selected branches
        overall_sales_data = filtered_sales_df.groupby(sales_area_col).agg({
            sales_qty_col: 'sum',
            sales_value_col: 'sum'
        }).reset_index()
        overall_sales_data.columns = ['Area', 'Overall_Sales_Qty', 'Overall_Sales_Value']
        
        # For overall sales, use the same budget totals for consistency
        budget_totals = results_df.groupby('Branch').agg({
            'Budget_Qty': 'sum',
            'Budget_Value': 'sum'
        }).reset_index()
        budget_totals.columns = ['Area', 'Budget_Qty', 'Budget_Value']
        
        # Overall Sales Quantity DataFrame
        overall_sales_qty_df = pd.DataFrame({'Area': default_branches})
        overall_sales_qty_df = pd.merge(
            overall_sales_qty_df,
            budget_totals[['Area', 'Budget_Qty']].rename(columns={'Budget_Qty': 'Budget Qty'}),
            on='Area',
            how='left'
        ).fillna({'Budget Qty': 0})
        
        overall_sales_qty_df = pd.merge(
            overall_sales_qty_df,
            overall_sales_data[['Area', 'Overall_Sales_Qty']].rename(columns={'Overall_Sales_Qty': 'Billed Qty'}),
            on='Area',
            how='left'
        ).fillna({'Billed Qty': 0})
        
        # Overall Sales Value DataFrame
        overall_sales_value_df = pd.DataFrame({'Area': default_branches})
        overall_sales_value_df = pd.merge(
            overall_sales_value_df,
            budget_totals[['Area', 'Budget_Value']].rename(columns={'Budget_Value': 'Budget Value'}),
            on='Area',
            how='left'
        ).fillna({'Budget Value': 0})
        
        overall_sales_value_df = pd.merge(
            overall_sales_value_df,
            overall_sales_data[['Area', 'Overall_Sales_Value']].rename(columns={'Overall_Sales_Value': 'Billed Value'}),
            on='Area',
            how='left'
        ).fillna({'Billed Value': 0})
        
        # STEP 5: Add Total Rows
        
        # Add total row for budget vs billed quantity
        total_budget_qty = budget_vs_billed_qty_df['Budget Qty'].sum()
        total_billed_qty = budget_vs_billed_qty_df['Billed Qty'].sum()
        total_percentage_qty = int((total_billed_qty / total_budget_qty * 100)) if total_budget_qty > 0 else 0
        
        total_row_qty = pd.DataFrame({
            'Area': ['TOTAL'],
            'Budget Qty': [total_budget_qty],
            'Billed Qty': [total_billed_qty],
            '%': [total_percentage_qty]
        })
        budget_vs_billed_qty_df = pd.concat([budget_vs_billed_qty_df, total_row_qty], ignore_index=True)
        
        # Add total row for budget vs billed value
        total_budget_value = budget_vs_billed_value_df['Budget Value'].sum()
        total_billed_value = budget_vs_billed_value_df['Billed Value'].sum()
        total_percentage_value = int((total_billed_value / total_budget_value * 100)) if total_budget_value > 0 else 0
        
        total_row_value = pd.DataFrame({
            'Area': ['TOTAL'],
            'Budget Value': [total_budget_value],
            'Billed Value': [total_billed_value],
            '%': [total_percentage_value]
        })
        budget_vs_billed_value_df = pd.concat([budget_vs_billed_value_df, total_row_value], ignore_index=True)
        
        # Add total rows for overall sales
        total_row_overall_qty = pd.DataFrame({
            'Area': ['TOTAL'],
            'Budget Qty': [overall_sales_qty_df['Budget Qty'].sum()],
            'Billed Qty': [overall_sales_qty_df['Billed Qty'].sum()]
        })
        overall_sales_qty_df = pd.concat([overall_sales_qty_df, total_row_overall_qty], ignore_index=True)
        
        total_row_overall_value = pd.DataFrame({
            'Area': ['TOTAL'],
            'Budget Value': [overall_sales_value_df['Budget Value'].sum()],
            'Billed Value': [overall_sales_value_df['Billed Value'].sum()]
        })
        overall_sales_value_df = pd.concat([overall_sales_value_df, total_row_overall_value], ignore_index=True)
        
        # Convert to integers for final display
        budget_vs_billed_value_df['Budget Value'] = budget_vs_billed_value_df['Budget Value'].round(0).astype(int)
        budget_vs_billed_value_df['Billed Value'] = budget_vs_billed_value_df['Billed Value'].round(0).astype(int)
        budget_vs_billed_qty_df['Budget Qty'] = budget_vs_billed_qty_df['Budget Qty'].round(0).astype(int)
        budget_vs_billed_qty_df['Billed Qty'] = budget_vs_billed_qty_df['Billed Qty'].round(0).astype(int)
        overall_sales_qty_df['Budget Qty'] = overall_sales_qty_df['Budget Qty'].round(0).astype(int)
        overall_sales_qty_df['Billed Qty'] = overall_sales_qty_df['Billed Qty'].round(0).astype(int)
        overall_sales_value_df['Budget Value'] = overall_sales_value_df['Budget Value'].round(0).astype(int)
        overall_sales_value_df['Billed Value'] = overall_sales_value_df['Billed Value'].round(0).astype(int)
        
        print("Budget vs Billed calculation completed successfully!")
        
        return budget_vs_billed_value_df, budget_vs_billed_qty_df, overall_sales_qty_df, overall_sales_value_df

    except Exception as e:
        logger.error(f"Error in calculate_values: {str(e)}")
        st.error(f"Error calculating budget values: {str(e)}")
        return None, None, None, None
def create_budget_ppt(budget_vs_billed_value_df, budget_vs_billed_qty_df, overall_sales_qty_df, overall_sales_value_df, month_title=None, logo_file=None):
    """Create PPT presentation with budget vs billed data"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout
        title_slide = prs.slides.add_slide(blank_slide_layout)
        company_name = title_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
        company_frame = company_name.text_frame
        company_frame.text = "Asia Crystal Commodity LLP"
        p = company_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        if logo_file is not None:
            try:
                logo_buffer = BytesIO(logo_file.read())
                logo = title_slide.shapes.add_picture(logo_buffer, Inches(5.665), Inches(1.5), width=Inches(2), height=Inches(2))
                logo_file.seek(0)
            except Exception as e:
                logger.error(f"Error adding logo: {e}")
        title = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(1))
        title_frame = title.text_frame
        title_frame.text = f"Monthly Review Meeting – {month_title}"
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)
        subtitle = title_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.33), Inches(1))
        subtitle_frame = subtitle.text_frame
        subtitle_frame.text = "ACCLLP"
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        def add_table_slide(title_text, df):
            slide = prs.slides.add_slide(blank_slide_layout)
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.75))
            title_text_frame = title_shape.text_frame
            title_text_frame.text = title_text
            p = title_text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 112, 192)
            table_width = Inches(9)
            table_height = Inches(3)
            left = Inches(2.165)  # (13.33 - 9) / 2 to center the table
            top = Inches(1.5)
            rows, cols = df.shape[0] + 1, df.shape[1]
            table = slide.shapes.add_table(rows, cols, left, top, table_width, table_height).table
            for col_idx, col_name in enumerate(df.columns):
                cell = table.cell(0, col_idx)
                cell.text = str(col_name)
                cell.text_frame.paragraphs[0].font.size = Pt(14)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 112, 192)
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
            total_row_idx = df.shape[0] - 1
            for row_idx in range(df.shape[0]):
                is_total_row = row_idx == total_row_idx
                for col_idx in range(df.shape[1]):
                    cell = table.cell(row_idx + 1, col_idx)
                    value = df.iloc[row_idx, col_idx]
                    cell.text = f"{value}%" if col_idx == 3 else str(value)
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    if is_total_row:
                        cell.text_frame.paragraphs[0].font.bold = True
                        fill = cell.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(211, 211, 211)
        add_table_slide("BUDGET AGAINST BILLED (Qty in Mt)", budget_vs_billed_qty_df)
        add_table_slide("BUDGET AGAINST BILLED (Value in Lakhs)", budget_vs_billed_value_df)
        add_table_slide("OVERALL SALES (Qty in Mt)", overall_sales_qty_df)
        add_table_slide("OVERALL SALES (Value in Lakhs)", overall_sales_value_df)
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating PPT: {e}")
        st.error(f"Error creating PPT: {e}")
        return None
def create_budget_table_image(df, title):
    """Create a table image for the budget vs billed report"""
    fig, ax = plt.subplots(figsize=(10, len(df) * 0.5))
    ax.axis('off')
    table = Table(ax, bbox=[0, 0, 1, 1])
    nrows, ncols = df.shape
    width, height = 1.0 / ncols, 1.0 / nrows
    for col_idx, col_name in enumerate(df.columns):
        table.add_cell(0, col_idx, width, height, text=str(col_name), loc='center', facecolor='#0070C0')
        table[0, col_idx].set_text_props(weight='bold', color='white')
    for row_idx in range(nrows):
        for col_idx in range(ncols):
            value = df.iloc[row_idx, col_idx]
            text = f"{value}%" if col_idx == 3 else str(value)
            facecolor = '#f0f0f0' if row_idx % 2 == 0 else 'white'
            if row_idx == nrows - 1:
                facecolor = '#D3D3D3'
                table.add_cell(row_idx + 1, col_idx, width, height, text=text, loc='center', facecolor=facecolor).set_text_props(weight='bold')
            else:
                table.add_cell(row_idx + 1, col_idx, width, height, text=text, loc='center', facecolor=facecolor)
    table.auto_set_font_size(False)
    table.set_fontsize(12)
    ax.add_table(table)
    plt.title(title, fontsize=16, weight='bold', color='#0070C0', pad=20)
    plt.tight_layout()
    img_buffer = BytesIO()
    plt.savefig(img_buffer, format='png', bbox_inches='tight', dpi=100)
    img_buffer.seek(0)
    plt.close()
    return img_buffer
def tab_budget_vs_billed():
    """Implements the Budget vs Billed Report Tab with AUTO-MAPPING."""
    st.header("Budget vs Billed Report") 
    dfs_info = []  
    sales_file = st.session_state.sales_file
    budget_file = st.session_state.budget_file  
    if not sales_file or not budget_file:
        st.warning("Please upload both Sales and Budget files in the sidebar")
        return dfs_info  
    try:
        sales_xl = pd.ExcelFile(sales_file)
        budget_xl = pd.ExcelFile(budget_file)
        sales_sheets = sales_xl.sheet_names
        budget_sheets = budget_xl.sheet_names
        st.subheader("Sheet Selection")
        col1, col2 = st.columns(2)
        with col1:
            sales_sheet = st.selectbox("Sales Sheet Name", sales_sheets, 
                                    index=sales_sheets.index('Sheet1') if 'Sheet1' in sales_sheets else 0, 
                                    key='budget_sales_sheet')
        with col2:
            budget_sheet = st.selectbox("Budget Sheet Name", budget_sheets, 
                                        index=budget_sheets.index('Sheet1') if 'Sheet1' in budget_sheets else 0, 
                                        key='budget_budget_sheet')
        st.subheader("Header Row Selection")
        col1, col2 = st.columns(2)
        with col1:
            sales_header_row = st.number_input("Sales Header Row (1-based)", min_value=1, value=1, step=1, key='budget_sales_header') - 1
        with col2:
            budget_header_row = st.number_input("Budget Header Row (1-based)", min_value=1, value=1, step=1, key='budget_budget_header') - 1
        sales_df = pd.read_excel(sales_file, sheet_name=sales_sheet, header=sales_header_row, dtype={'SL Code': str})
        budget_df = pd.read_excel(budget_file, sheet_name=budget_sheet, header=budget_header_row, dtype={'SL Code': str})
        sales_columns = sales_df.columns.tolist()
        budget_columns = budget_df.columns.tolist()
        sales_mapping, budget_mapping = auto_map_budget_columns(sales_columns, budget_columns)
        st.subheader("Sales Column Mapping")
        col1, col2, col3 = st.columns(3)
        with col1:
            sales_date_col = st.selectbox("Sales Date Column", sales_columns, 
                                        index=sales_columns.index(sales_mapping['date']) if sales_mapping['date'] else 0, 
                                        key='budget_sales_date')
            sales_area_col = st.selectbox("Sales Area Column", sales_columns, 
                                        index=sales_columns.index(sales_mapping['area']) if sales_mapping['area'] else 0, 
                                        key='budget_sales_area')
        with col2:
            sales_value_col = st.selectbox("Sales Value Column", sales_columns, 
                                        index=sales_columns.index(sales_mapping['value']) if sales_mapping['value'] else 0, 
                                        key='budget_sales_value')
            sales_qty_col = st.selectbox("Sales Quantity Column", sales_columns, 
                                        index=sales_columns.index(sales_mapping['quantity']) if sales_mapping['quantity'] else 0, 
                                        key='budget_sales_qty')
        with col3:
            sales_product_group_col = st.selectbox("Sales Product Group Column", sales_columns, 
                                                index=sales_columns.index(sales_mapping['product_group']) if sales_mapping['product_group'] else 0, 
                                                key='budget_sales_product')
            sales_sl_code_col = st.selectbox("Sales SL Code Column", sales_columns, 
                                            index=sales_columns.index(sales_mapping['sl_code']) if sales_mapping['sl_code'] else 0, 
                                            key='budget_sales_sl_code')
        
        sales_exec_col = st.selectbox("Sales Executive Column", sales_columns, 
                                    index=sales_columns.index(sales_mapping['executive']) if sales_mapping['executive'] else 0, 
                                    key='budget_sales_exec')
        st.subheader("Budget Column Mapping")
        col1, col2, col3 = st.columns(3)
        with col1:
            budget_area_col = st.selectbox("Budget Area Column", budget_columns, 
                                        index=budget_columns.index(budget_mapping['area']) if budget_mapping['area'] else 0, 
                                        key='budget_budget_area')
            budget_value_col = st.selectbox("Budget Value Column", budget_columns, 
                                            index=budget_columns.index(budget_mapping['value']) if budget_mapping['value'] else 0, 
                                            key='budget_budget_value')
        with col2:
            budget_qty_col = st.selectbox("Budget Quantity Column", budget_columns, 
                                        index=budget_columns.index(budget_mapping['quantity']) if budget_mapping['quantity'] else 0, 
                                        key='budget_budget_qty')
            budget_product_group_col = st.selectbox("Budget Product Group Column", budget_columns, 
                                                    index=budget_columns.index(budget_mapping['product_group']) if budget_mapping['product_group'] else 0, 
                                                    key='budget_budget_product')
        with col3:
            budget_sl_code_col = st.selectbox("Budget SL Code Column", budget_columns, 
                                            index=budget_columns.index(budget_mapping['sl_code']) if budget_mapping['sl_code'] else 0, 
                                            key='budget_budget_sl_code')
            budget_exec_col = st.selectbox("Budget Executive Column", budget_columns, 
                                        index=budget_columns.index(budget_mapping['executive']) if budget_mapping['executive'] else 0, 
                                        key='budget_budget_exec')
        if 'budget_sales_exec_all' not in st.session_state:
            st.session_state.budget_sales_exec_all = True
        if 'budget_budget_exec_all' not in st.session_state:
            st.session_state.budget_budget_exec_all = True
        if 'budget_branch_all' not in st.session_state:
            st.session_state.budget_branch_all = True
        st.subheader("Executive Selection")
        sales_executives = sorted(sales_df[sales_exec_col].dropna().unique().tolist())
        budget_executives = sorted(budget_df[budget_exec_col].dropna().unique().tolist())
        
        if not sales_executives or not budget_executives:
            st.error("No executives found in selected Executive columns.")
            return dfs_info
        col1, col2 = st.columns(2)
        with col1:
            st.write("**Sales Executives**")
            sales_select_all = st.checkbox("Select All Sales Executives", value=st.session_state.budget_sales_exec_all, key='budget_sales_exec_all')
            selected_sales_execs = []
            if sales_select_all != st.session_state.budget_sales_exec_all:
                st.session_state.budget_sales_exec_all = sales_select_all
                for exec in sales_executives:
                    st.session_state[f'budget_sales_exec_{exec}'] = sales_select_all
            for exec in sales_executives:
                if sales_select_all:
                    st.session_state[f'budget_sales_exec_{exec}'] = True
                if st.checkbox(exec, value=st.session_state.get(f'budget_sales_exec_{exec}', True), key=f'budget_sales_exec_{exec}'):
                    selected_sales_execs.append(exec)
                if not st.session_state.get(f'budget_sales_exec_{exec}', True) and st.session_state.budget_sales_exec_all:
                    st.session_state.budget_sales_exec_all = False      
        with col2:
            st.write("**Budget Executives**")
            budget_select_all = st.checkbox("Select All Budget Executives", value=st.session_state.budget_budget_exec_all, key='budget_budget_exec_all')
            selected_budget_execs = []
            if budget_select_all != st.session_state.budget_budget_exec_all:
                st.session_state.budget_budget_exec_all = budget_select_all
                for exec in budget_executives:
                    st.session_state[f'budget_budget_exec_{exec}'] = budget_select_all
            for exec in budget_executives:
                if budget_select_all:
                    st.session_state[f'budget_budget_exec_{exec}'] = True
                if st.checkbox(exec, value=st.session_state.get(f'budget_budget_exec_{exec}', True), key=f'budget_budget_exec_{exec}'):
                    selected_budget_execs.append(exec)
                if not st.session_state.get(f'budget_budget_exec_{exec}', True) and st.session_state.budget_budget_exec_all:
                    st.session_state.budget_budget_exec_all = False
        all_branches = pd.concat([sales_df[sales_area_col], budget_df[budget_area_col]]).dropna().str.upper().str.split(' - ').str[-1].unique().tolist()
        all_branches = sorted(set([branch_mapping.get(branch, branch) for branch in all_branches]))      
        st.subheader("Branch Selection")
        branch_select_all = st.checkbox("Select All Branches", value=True, key='od_filter_branch_all')
        selected_branches = []
        if branch_select_all != st.session_state.budget_branch_all:
            st.session_state.budget_branch_all = branch_select_all
            for branch in all_branches:
                st.session_state[f'budget_branch_{branch}'] = branch_select_all
        branch_cols = st.columns(4)  # Adjust the number of columns based on how many branches there are      
        for i, branch in enumerate(all_branches):
            col_idx = i % 4  # Distribute branches across columns
            with branch_cols[col_idx]:
                if branch_select_all:
                    st.session_state[f'budget_branch_{branch}'] = True
                if st.checkbox(branch, value=st.session_state.get(f'budget_branch_{branch}', True), key=f'budget_branch_{branch}'):
                    selected_branches.append(branch)
                if not st.session_state.get(f'budget_branch_{branch}', True) and st.session_state.budget_branch_all:
                    st.session_state.budget_branch_all = False
        months = pd.to_datetime(sales_df[sales_date_col], dayfirst=True, errors='coerce').dt.strftime('%b %y').dropna().unique()
        if len(months) == 0:
            st.error(f"No valid months found in '{sales_date_col}'. Check date format.")
            return dfs_info          
        selected_month = st.selectbox("Select Sales Month", months, index=0, key='budget_month')
        if st.button("Calculate Budget vs Billed", key='budget_calculate'):
            # Use the fixed calculation function
            budget_vs_billed_value_df, budget_vs_billed_qty_df, overall_sales_qty_df, overall_sales_value_df = calculate_values(
                sales_df, budget_df, selected_month, selected_sales_execs, selected_budget_execs,
                sales_date_col, sales_area_col, sales_value_col, sales_qty_col, sales_product_group_col, sales_sl_code_col, sales_exec_col,
                budget_area_col, budget_value_col, budget_qty_col, budget_product_group_col, budget_sl_code_col, budget_exec_col,
                selected_branches
            )
            
            if all(df is not None for df in [budget_vs_billed_value_df, budget_vs_billed_qty_df, overall_sales_qty_df, overall_sales_value_df]):
                st.subheader("Results")              
                result_tabs = st.tabs(["Budget vs Billed Quantity", "Budget vs Billed Value", "Overall Sales (Qty)", "Overall Sales (Value)"])              
                with result_tabs[0]:
                    st.write("**Budget vs Billed Quantity**")
                    st.dataframe(budget_vs_billed_qty_df)
                    img_buffer = create_budget_table_image(budget_vs_billed_qty_df, f"BUDGET AGAINST BILLED (Qty in Mt) - {selected_month}")
                    st.image(img_buffer, use_column_width=True)                
                with result_tabs[1]:
                    st.write("**Budget vs Billed Value**")
                    st.dataframe(budget_vs_billed_value_df)
                    img_buffer = create_budget_table_image(budget_vs_billed_value_df, f"BUDGET AGAINST BILLED (Value in Lakhs) - {selected_month}")
                    st.image(img_buffer, use_column_width=True)                              
                with result_tabs[2]:
                    st.write("**Overall Sales Quantity**")
                    st.dataframe(overall_sales_qty_df)
                    img_buffer = create_budget_table_image(overall_sales_qty_df, f"OVERALL SALES (Qty In Mt) - {selected_month}")
                    st.image(img_buffer, use_column_width=True)              
                with result_tabs[3]:
                    st.write("**Overall Sales Value**")
                    st.dataframe(overall_sales_value_df)
                    img_buffer = create_budget_table_image(overall_sales_value_df, f"OVERALL SALES (Value in Lakhs) - {selected_month}")
                    st.image(img_buffer, use_column_width=True)
                logo_file = st.session_state.logo_file
                ppt_buffer = create_budget_ppt(    
                    budget_vs_billed_qty_df,
                    budget_vs_billed_value_df, 
                    overall_sales_qty_df, 
                    overall_sales_value_df, 
                    month_title=selected_month, 
                    logo_file=logo_file
                )
                if ppt_buffer:
                    st.download_button(
                        label="Download Budget vs Billed PPT",
                        data=ppt_buffer,
                        file_name=f"Budget_vs_Billed_{selected_month}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key='budget_download'
                    )
                dfs_info = [
                    {'df': budget_vs_billed_qty_df, 'title': f"BUDGET AGAINST BILLED QUANTITY (Qty in Mt) - {selected_month}", 'percent_cols': [3]},
                    {'df': budget_vs_billed_value_df, 'title': f"BUDGET AGAINST BILLED (Value in Lakhs) - {selected_month}", 'percent_cols': [3]},
                    {'df': overall_sales_qty_df, 'title': f"OVERALL SALES (Qty In Mt) - {selected_month}", 'percent_cols': [3]},
                    {'df': overall_sales_value_df, 'title': f"OVERALL SALES (Value in Lakhs) - {selected_month}", 'percent_cols': [3]}
                ]
                st.session_state.budget_results = dfs_info
            else:
                st.error("Failed to calculate values. Check your data and selections.")
    except Exception as e:
        st.error(f"Error in Budget vs Billed tab: {e}")
        logger.error(f"Error in Budget vs Billed tab: {e}")  
    # Ensure dfs_info is populated or remove the return statement if not needed
    if dfs_info:
        # If dfs_info is not used, remove the return statement
        pass
    else:
        st.warning("No data to return. Please check your inputs and configurations.")
        return None
def determine_financial_year(date):
    """Determine the financial year for a given date (April to March)."""
    year = date.year
    month = date.month
    if month >= 4:
        return f"{year % 100}-{year % 100 + 1}"  # e.g., April 2024 -> 24-25
    else:
        return f"{year % 100 - 1}-{year % 100}"  # e.g., February 2025 -> 24-25
def create_customer_table(sales_df, date_col, branch_col, customer_id_col, executive_col, selected_branches=None, selected_executives=None):
    """Create a table of unique customer counts per branch and month for available data."""
    sales_df = sales_df.copy()
    for col in [date_col, branch_col, customer_id_col, executive_col]:
        if col not in sales_df.columns:
            st.error(f"Column '{col}' not found in sales data.")
            return None
    try:
        sales_df[date_col] = pd.to_datetime(sales_df[date_col], errors='coerce')
    except Exception as e:
        st.error(f"Error converting '{date_col}' to datetime: {e}. Ensure dates are in a valid format.")
        return None
    valid_dates = sales_df[date_col].notna()
    if not valid_dates.any():
        st.error(f"Column '{date_col}' contains no valid dates.")
        return None
    sales_df['Financial_Year'] = sales_df[date_col].apply(determine_financial_year)
    available_financial_years = sorted(sales_df['Financial_Year'].unique())
    sales_df = sales_df[sales_df[date_col].notna()].copy()
    if sales_df.empty:
        st.error("No valid dates found in data after filtering.")
        return None
    sales_df['Month'] = sales_df[date_col].dt.month
    sales_df['Year'] = sales_df[date_col].dt.year
    sales_df['Month_Name'] = sales_df[date_col].dt.strftime('%b-%Y').str.upper()
    try:
        sales_df['Raw_Branch'] = sales_df[branch_col].astype(str).str.upper()
    except Exception as e:
        st.error(f"Error processing branch column '{branch_col}': {e}.")
        return None
    sales_df['Mapped_Branch'] = sales_df['Raw_Branch'].replace(nbc_branch_mapping)
    if selected_branches:
        sales_df = sales_df[sales_df['Mapped_Branch'].isin(selected_branches)]
        if sales_df.empty:
            st.error("No data matches the selected branches.")
            return None
    if selected_executives:
        sales_df = sales_df[sales_df[executive_col].isin(selected_executives)]
        if sales_df.empty:
            st.error("No data matches the selected executives.")
            return None
    result_dict = {}
    for fy in available_financial_years:
        fy_df = sales_df[sales_df['Financial_Year'] == fy].copy()
        
        if fy_df.empty:
            continue
        fy_df['Date_Sort'] = fy_df[date_col]  # For sorting months chronologically
        available_months = fy_df.sort_values('Date_Sort')['Month_Name'].unique()
        if selected_branches:
            branches_to_display = selected_branches
        else:
            branches_to_display = sorted(fy_df['Mapped_Branch'].dropna().unique())
        grouped_df = fy_df.groupby(['Raw_Branch', 'Month_Name'])[customer_id_col].nunique().reset_index(name='Count')
        grouped_df['Mapped_Branch'] = grouped_df['Raw_Branch'].replace(nbc_branch_mapping)
        pivot_df = grouped_df.pivot_table(
            values='Count',
            index='Mapped_Branch',
            columns='Month_Name',
            aggfunc='sum',
            fill_value=0
        ).reset_index()
        result_df = pd.DataFrame({'Branch Name': branches_to_display})
        result_df = pd.merge(result_df, pivot_df, left_on='Branch Name', right_on='Mapped_Branch', how='left').fillna(0)
        result_df = result_df.drop(columns=['Mapped_Branch'] if 'Mapped_Branch' in result_df.columns else [])
        for month in available_months:
            if month not in result_df.columns:
                result_df[month] = 0
        result_df = result_df[['Branch Name'] + [month for month in available_months if month in result_df.columns]]
        result_df.insert(0, 'S.No', range(1, len(result_df) + 1))
        total_row = {'S.No': '', 'Branch Name': 'GRAND TOTAL'}
        for month in available_months:
            if month in result_df.columns:
                total_row[month] = result_df[month].sum()
        result_df = pd.concat([result_df, pd.DataFrame([total_row])], ignore_index=True)
        for col in available_months:
            if col in result_df.columns:
                result_df[col] = result_df[col].round(0).astype(int)
        result_dict[fy] = (result_df, list(available_months))    
    return result_dict
def create_customer_table_image(df, title, sorted_months, financial_year):
    """Create a table image from the DataFrame."""
    fig, ax = plt.subplots(figsize=(14, len(df) * 0.6))
    ax.axis('off')
    columns = ['S.No', 'Branch Name'] + sorted_months
    rows = len(df)
    ncols = len(columns)
    width = 1.0 / ncols
    height = 1.0 / rows
    table = Table(ax, bbox=[0, 0, 1, 1])
    for col_idx, col_name in enumerate(columns):
        table.add_cell(0, col_idx, width, height, text=col_name, loc='center', facecolor='#0070C0')
        table[0, col_idx].set_text_props(weight='bold', color='white', fontsize=10)
    for row_idx in range(rows):
        for col_idx in range(ncols):
            value = df.iloc[row_idx, col_idx]
            text = str(value)
            facecolor = '#DDEBF7' if row_idx % 2 == 0 else 'white'
            if row_idx == rows - 1:  # GRAND TOTAL
                facecolor = '#D3D3D3'
                table.add_cell(row_idx + 1, col_idx, width, height, text=text, loc='center', facecolor=facecolor).set_text_props(weight='bold', fontsize=10)
            else:
                table.add_cell(row_idx + 1, col_idx, width, height, text=text, loc='center', facecolor=facecolor).set_text_props(fontsize=10)
    table[(0, 0)].width = 0.05
    table[(0, 1)].width = 0.15
    for col_idx in range(2, ncols):
        table[(0, col_idx)].width = 0.08
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    ax.add_table(table)
    plt.suptitle(f"{title} - FY {financial_year}", fontsize=14, weight='bold', color='#0070C0', y=1.02)
    img_buffer = BytesIO()
    plt.savefig(img_buffer, format='png', bbox_inches='tight', dpi=150)
    plt.close()
    return img_buffer
def create_customer_ppt_slide(slide, df, title, sorted_months, financial_year):
    """Add a slide with the billed customers table to the presentation."""
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = f"{title} - FY {financial_year}"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER
    columns = ['S.No', 'Branch Name'] + sorted_months
    rows = len(df) + 1
    ncols = len(columns)
    table_width = Inches(12.0)
    table_height = Inches(0.3 * len(df) + 0.4)
    left = Inches(0.65)
    top = Inches(1.2)
    table = slide.shapes.add_table(rows, ncols, left, top, table_width, table_height).table
    col_widths = [Inches(0.5), Inches(2.0)] + [Inches(0.75)] * len(sorted_months)
    for col_idx in range(ncols):
        table.columns[col_idx].width = col_widths[col_idx]
    for col_idx, col_name in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    total_row_idx = len(df) - 1
    for row_idx in range(len(df)):
        is_total_row = row_idx == total_row_idx
        for col_idx in range(ncols):
            cell = table.cell(row_idx + 1, col_idx)
            value = df.iloc[row_idx, col_idx]
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            if is_total_row:
                cell.fill.fore_color.rgb = RGBColor(211, 211, 211)
                cell.text_frame.paragraphs[0].font.bold = True
            else:
                if row_idx % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(221, 235, 247)
                else:
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255) 
def extract_area_name(area):
    """Extract and standardize branch names."""
    if pd.isna(area) or str(area).strip() == '':
        return None  # Return None for empty/null values  
    area = str(area).strip().upper()
    if area == 'HO' or area.endswith('-HO'):
        return None
    branch_variations = {
        'PUDUCHERRY': ['PUDUCHERRY', 'PONDY', 'PUDUCHERRY - PONDY', 'PONDICHERRY', 'PUDUCHERI', 'aaaa - PUDUCHERRY', 'AAAA - PUDUCHERRY'],
        'COIMBATORE': ['COIMBATORE', 'CBE', 'COIMBATORE - CBE', 'COIMBATURE', 'aaaa - COIMBATORE', 'AAAA - COIMBATORE'],
        'KARUR': ['KARUR', 'KRR', 'KARUR - KRR', 'aaaa - KARUR', 'AAAA - KARUR'],
        'MADURAI': ['MADURAI', 'MDU', 'MADURAI - MDU', 'MADURA', 'aaaa - MADURAI', 'AAAA - MADURAI'],
        'CHENNAI': ['CHENNAI', 'CHN', 'aaaa - CHENNAI', 'AAAA - CHENNAI'],
    }
    
    for standard_name, variations in branch_variations.items():
        for variation in variations:
            if variation in area:
                return standard_name
    prefixes = ['AAAA - ', 'aaaa - ', 'BBB - ', 'bbb - ', 'ASIA CRYSTAL COMMODITY LLP - ']
    for prefix in prefixes:
        if area.startswith(prefix.upper()):
            return area[len(prefix):].strip()
    separators = [' - ', '-', ':']
    for sep in separators:
        if sep in area:
            return area.split(sep)[-1].strip()    
    return area

def extract_executive_name(executive):
    """Normalize executive names, treating null/empty as 'BLANK'."""
    if pd.isna(executive) or str(executive).strip() == '':
        return 'BLANK'  # Treat null/empty as 'BLANK'
    return str(executive).strip().upper()  # Standardize to uppercase
def filter_os_qty(os_df, os_area_col, os_qty_col, os_due_date_col, os_exec_col, 
               selected_branches=None, selected_years=None, till_month=None, selected_executives=None):
    """Filter by due date and aggregate net values by area, applying branch/executive filters."""
    required_columns = [os_area_col, os_qty_col, os_due_date_col, os_exec_col]
    for col in required_columns:
        if col not in os_df.columns:
            st.error(f"Column '{col}' not found in OS data.")
            return None, None, None
    os_df = os_df.copy()
    os_df[os_area_col] = os_df[os_area_col].apply(extract_area_name)
    os_df[os_exec_col] = os_df[os_exec_col].apply(extract_executive_name)
    try:
        os_df[os_due_date_col] = pd.to_datetime(os_df[os_due_date_col], errors='coerce')
    except Exception as e:
        st.error(f"Error converting '{os_due_date_col}' to datetime: {e}. Ensure dates are in 'YYYY-MM-DD' format.")
        return None, None, None
    start_date, end_date = None, None
    if selected_years and till_month:
        month_map = {
            'January': 1, 'February': 2, 'March': 3, 'April': 4, 'May': 5, 'June': 6,
            'July': 7, 'August': 8, 'September': 9, 'October': 10, 'November': 11, 'December': 12
        }
        till_month_num = month_map.get(till_month)
        if not till_month_num:
            st.error("Invalid month selected.")
            return None, None, None
        selected_years = [int(year) for year in selected_years]
        earliest_year = min(selected_years)
        latest_year = max(selected_years)
        start_date = datetime(earliest_year, 1, 1)
        end_date = (datetime(latest_year, till_month_num, 1) + relativedelta(months=1) - relativedelta(days=1))
        os_df = os_df[
            (os_df[os_due_date_col].notna()) &
            (os_df[os_due_date_col] >= start_date) &
            (os_df[os_due_date_col] <= end_date)
        ]
        if os_df.empty:
            st.error(f"No data matches the period from Jan {earliest_year} to {end_date.strftime('%b %Y')}.")
            return None, None, None
    all_branches = sorted(os_df[os_area_col].dropna().unique())
    if not selected_branches:
        selected_branches = all_branches  # Default to all branches if none selected
    if sorted(selected_branches) != all_branches:
        os_df = os_df[os_df[os_area_col].isin(selected_branches)]
        if os_df.empty:
            st.error("No data matches the selected branches.")
            return None, None, None
    all_executives = sorted(os_df[os_exec_col].dropna().unique())
    if selected_executives and sorted(selected_executives) != sorted(all_executives):
        os_df = os_df[os_df[os_exec_col].isin(selected_executives)]
        if os_df.empty:
            st.error("No data matches the selected executives.")
            return None, None, None
    os_df[os_qty_col] = pd.to_numeric(os_df[os_qty_col], errors='coerce').fillna(0)
    if os_df[os_qty_col].isna().any():
        st.warning(f"Non-numeric values in '{os_qty_col}' replaced with 0.")
    os_df_positive = os_df[os_df[os_qty_col] > 0].copy()
    if os_df_positive.empty:
        st.warning("No positive net values found in the filtered data.")
    os_grouped_qty = (os_df_positive.groupby(os_area_col)
                    .agg({os_qty_col: 'sum'})
                    .reset_index()
                    .rename(columns={os_area_col: 'Area', os_qty_col: 'TARGET'}))
    os_grouped_qty['TARGET'] = os_grouped_qty['TARGET'] / 100000  # Convert to lakhs
    branches_to_display = selected_branches if selected_branches else all_branches
    result_df = pd.DataFrame({'Area': branches_to_display})
    result_df = pd.merge(result_df, os_grouped_qty, on='Area', how='left').fillna({'TARGET': 0})
    total_row = pd.DataFrame([{'Area': 'TOTAL', 'TARGET': result_df['TARGET'].sum()}])
    result_df = pd.concat([result_df, total_row], ignore_index=True)
    result_df['TARGET'] = result_df['TARGET'].round(2)
    return result_df, start_date, end_date
def create_od_table_image(df, title, columns_to_show=None):
    """Create a table image from the OD Target DataFrame."""
    if columns_to_show is None:
        columns_to_show = ['Area', 'TARGET (Lakhs)']    
    fig, ax = plt.subplots(figsize=(10, len(df) * 0.5))
    ax.axis('off')
    nrows, ncols = len(df), len(columns_to_show)
    table = Table(ax, bbox=[0, 0, 1, 1])
    for col_idx, col_name in enumerate(columns_to_show):
        table.add_cell(0, col_idx, 1.0/ncols, 1.0/nrows, text=col_name, loc='center', facecolor='#F2F2F2')
        table[0, col_idx].set_text_props(weight='bold', color='black', fontsize=12)

    for row_idx in range(len(df)):
        for col_idx, col_name in enumerate(columns_to_show):
            value = df.iloc[row_idx]['Area'] if col_name == 'Area' else df.iloc[row_idx]['TARGET']
            text = str(value) if col_name == 'Area' else f"{float(value):.2f}"
            facecolor = '#DDEBF7' if row_idx % 2 == 0 else 'white'
            if row_idx == len(df) - 1 and df.iloc[row_idx, 0] == 'TOTAL':
                facecolor = '#D3D3D3'
                table.add_cell(row_idx + 1, col_idx, 1.0/ncols, 1.0/nrows, text=text, loc='center', facecolor=facecolor).set_text_props(weight='bold', fontsize=12)
            else:
                table.add_cell(row_idx + 1, col_idx, 1.0/ncols, 1.0/nrows, text=text, loc='center', facecolor=facecolor).set_text_props(fontsize=10)
    table[(0, 0)].width = 0.6
    table[(0, 1)].width = 0.4
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    ax.add_table(table)
    plt.suptitle(title, fontsize=16, weight='bold', color='black', y=1.05)
    img_buffer = BytesIO()
    plt.savefig(img_buffer, format='png', bbox_inches='tight', dpi=150)
    plt.close()
    return img_buffer
def create_od_ppt_slide(slide, df, title):
    """Add a slide with the OD Target table to the presentation."""
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    columns_to_show = ['Area', 'TARGET (Lakhs)']
    rows, cols = df.shape[0] + 1, len(columns_to_show)
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(5)).table
    for col_idx, col_name in enumerate(columns_to_show):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for row_idx in range(len(df)):
        is_total_row = row_idx == len(df) - 1
        for col_idx, col_name in enumerate(columns_to_show):
            cell = table.cell(row_idx + 1, col_idx)
            value = df.iloc[row_idx]['Area'] if col_name == 'Area' else df.iloc[row_idx]['TARGET']
            cell.text = str(value) if col_name == 'Area' else f"{float(value):.2f}"
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 204) if is_total_row else (RGBColor(221, 235, 247) if row_idx % 2 == 0 else RGBColor(255, 255, 255))
            if is_total_row:
                cell.text_frame.paragraphs[0].font.bold = True
def create_combined_nbc_od_ppt(customer_df, customer_title, sorted_months, od_target_df, od_title, logo_file=None):
    """Create a single PPT with two slides: one for billed customers, one for OD Target."""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        create_title_slide(prs, "Number of Billed Customers & OD Target Report", logo_file)
        slide_layout = prs.slide_layouts[6]  # Blank slide
        slide1 = prs.slides.add_slide(slide_layout)
        create_customer_ppt_slide(slide1, customer_df, customer_title, sorted_months)
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide2 = prs.slides.add_slide(slide_layout)
        create_od_ppt_slide(slide2, od_target_df, od_title)
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating Combined NBC & OD Target PPT: {e}")
        st.error(f"Error creating Combined NBC & OD Target PPT: {e}")
        return None
def tab_billed_customers():
    """Number of Billed Customers Report and OD Target Report Tab with AUTO-MAPPING."""
    st.header("Number of Billed Customers & OD Target Report")   
    dfs_info = []
    sales_file = st.session_state.sales_file
    os_jan_file = st.session_state.os_jan_file
    os_feb_file = st.session_state.os_feb_file    
    if not sales_file:
        st.warning("Please upload the Sales file in the sidebar")
        return dfs_info
    nbc_tab, od_tab = st.tabs(["Number of Billed Customers", "OD Target"])
    with nbc_tab:
        st.subheader("Number of Billed Customers Setup")        
        try:
            sales_sheets = get_excel_sheets(sales_file)
            
            if not sales_sheets:
                st.error("No sheets found in Sales file.")
                return dfs_info
            sales_sheet = st.selectbox("Select Sales Sheet", sales_sheets, 
                                    index=sales_sheets.index('Sheet1') if 'Sheet1' in sales_sheets else 0, 
                                    key='nbc_sales_sheet')
            sales_df = pd.read_excel(sales_file, sheet_name=sales_sheet)
            columns = sales_df.columns.tolist()
            nbc_mapping = auto_map_nbc_columns(columns)
            st.subheader("Column Mapping")
            col1, col2 = st.columns(2)
            with col1:
                date_col = st.selectbox(
                    "Date Column",
                    columns,
                    index=columns.index(nbc_mapping['date']) if nbc_mapping['date'] else 0,
                    help="This column should contain dates (e.g., '2024-04-01').",
                    key='nbc_date_col'
                )
                branch_col = st.selectbox(
                    "Branch Column",
                    columns,
                    index=columns.index(nbc_mapping['branch']) if nbc_mapping['branch'] else 0,
                    help="This column should contain branch names.",
                    key='nbc_branch_col'
                )
            with col2:
                customer_id_col = st.selectbox(
                    "Customer ID Column",
                    columns,
                    index=columns.index(nbc_mapping['customer_id']) if nbc_mapping['customer_id'] else 0,
                    help="This column should contain unique customer identifiers.",
                    key='nbc_customer_id_col'
                )
                executive_col = st.selectbox(
                    "Executive Column",
                    columns,
                    index=columns.index(nbc_mapping['executive']) if nbc_mapping['executive'] else 0,
                    help="This column should contain executive names for filtering.",
                    key='nbc_executive_col'
                )
            st.subheader("Filter Options")
            filter_tab1, filter_tab2 = st.tabs(["Branches", "Executives"])            
            with filter_tab1:
                raw_branches = sales_df[branch_col].dropna().astype(str).str.upper().unique().tolist()
                all_nbc_branches = sorted(set([nbc_branch_mapping.get(branch.split(' - ')[-1], branch.split(' - ')[-1]) for branch in raw_branches]))
                
                branch_select_all = st.checkbox("Select All Branches", value=True, key='nbc_branch_all')
                if branch_select_all:
                    selected_branches = all_nbc_branches
                else:
                    selected_branches = st.multiselect("Select Branches", all_nbc_branches, key='nbc_branches')
            with filter_tab2:
                all_executives = sorted(sales_df[executive_col].dropna().unique().tolist())
                
                exec_select_all = st.checkbox("Select All Executives", value=True, key='nbc_exec_all')
                if exec_select_all:
                    selected_executives = all_executives
                else:
                    selected_executives = st.multiselect("Select Executives", all_executives, key='nbc_executives')
            if st.button("Generate Billed Customers Report", key='nbc_generate'):
                results = create_customer_table(
                    sales_df, date_col, branch_col, customer_id_col, executive_col,
                    selected_branches=selected_branches,
                    selected_executives=selected_executives
                )                
                if results:
                    st.subheader("Results")
                    for fy, (result_df, sorted_months) in results.items():
                        st.write(f"**Financial Year: {fy}**")
                        st.dataframe(result_df)
                        title = "NUMBER OF BILLED CUSTOMERS"
                        img_buffer = create_customer_table_image(result_df, title, sorted_months, fy)
                        if img_buffer:
                            st.image(img_buffer, use_column_width=True)
                        customers_dfs = [{'df': result_df, 'title': f"NUMBER OF BILLED CUSTOMERS - FY {fy}"}]
                        st.session_state.customers_results = customers_dfs
                        dfs_info.append({'df': result_df, 'title': f"NUMBER OF BILLED CUSTOMERS - FY {fy}"})
                else:
                    st.error("Failed to generate Number of Billed Customers report. Check your data and selections.")
        except Exception as e:
            st.error(f"Error in Number of Billed Customers tab: {e}")
            logger.error(f"Error in Number of Billed Customers tab: {e}")
    with od_tab:
        st.subheader("OD Target Setup")
        if not os_jan_file or not os_feb_file:
            missing_files = []
            if not os_jan_file:
                missing_files.append("OS-Previous Month")
            if not os_feb_file:
                missing_files.append("OS-Current Month")
            st.error(f"Please upload {' and '.join(missing_files)} files in the sidebar for OD Target calculation.")
        else:
            try:
                st.subheader("File and Sheet Selection")
                os_file_choice = st.radio(
                    "Choose OS file for OD Target calculation:", 
                    ["OS-Previous Month", "OS-Current Month"],
                    key="od_file_choice"
                )
                os_file = os_jan_file if os_file_choice == "OS-Previous Month" else os_feb_file
                os_sheets = get_excel_sheets(os_file)                
                if not os_sheets:
                    st.error(f"No sheets found in {os_file_choice} file.")
                else:
                    os_sheet = st.selectbox(
                        f"Select {os_file_choice} Sheet", 
                        os_sheets, 
                        index=os_sheets.index('Sheet1') if 'Sheet1' in os_sheets else 0,
                        key='od_sheet'
                    )
                    header_row = st.number_input(
                        "Header Row (1-based)", 
                        min_value=1, 
                        max_value=11, 
                        value=1, 
                        key='od_header_row'
                    ) - 1
                    os_df = pd.read_excel(os_file, sheet_name=os_sheet, header=header_row)
                    if st.checkbox("Preview Raw OS Data", key='od_preview'):
                        st.write(f"Raw {os_file_choice} Data (first 20 rows):")
                        st.dataframe(os_df.head(20))
                    columns = os_df.columns.tolist()
                    od_mapping = auto_map_od_target_columns(columns)                    
                    st.subheader("OS Column Mapping")
                    col1, col2 = st.columns(2)
                    with col1:
                        os_area_col = st.selectbox(
                            "Area Column",
                            columns,
                            index=columns.index(od_mapping['area']) if od_mapping['area'] else 0,
                            help="Contains branch names (e.g., COMPANY).",
                            key='od_area_col'
                        )
                        os_qty_col = st.selectbox(
                            "Net Value Column",
                            columns,
                            index=columns.index(od_mapping['net_value']) if od_mapping['net_value'] else 0,
                            help="Contains net values in INR.",
                            key='od_qty_col'
                        )
                    with col2:
                        os_due_date_col = st.selectbox(
                            "Due Date Column",
                            columns,
                            index=columns.index(od_mapping['due_date']) if od_mapping['due_date'] else 0,
                            help="Contains due dates (e.g., '2018-02-15').",
                            key='od_due_date_col'
                        )
                        os_exec_col = st.selectbox(
                            "Executive Column",
                            columns,
                            index=columns.index(od_mapping['executive']) if od_mapping['executive'] else 0,
                            help="Contains executive names (BLANK for null/empty).",
                            key='od_exec_col'
                        )
                    st.subheader("Due Date Filter")
                    try:
                        os_df[os_due_date_col] = pd.to_datetime(os_df[os_due_date_col], errors='coerce')
                        years = sorted(os_df[os_due_date_col].dt.year.dropna().astype(int).unique())
                    except Exception as e:
                        st.error(f"Error processing due dates: {e}. Ensure valid date format.")
                        years = []

                    if not years:
                        st.warning(f"No valid years found in {os_file_choice}'s due date column.")
                    else:
                        selected_years = st.multiselect(
                            "Select years for filtering",
                            options=[str(year) for year in years],
                            default=[str(year) for year in years],
                            key='od_year_multiselect'
                        )

                        if not selected_years:
                            st.error("Please select at least one year.")
                        else:
                            month_options = ['January', 'February', 'March', 'April', 'May', 'June',
                                            'July', 'August', 'September', 'October', 'November', 'December']
                            till_month = st.selectbox("Select end month", month_options, key='od_till_month')
                    st.subheader("Branch Selection")
                    os_branches = sorted(set([b for b in os_df[os_area_col].dropna().apply(extract_area_name) if b]))
                    if not os_branches:
                        st.error(f"No valid branches found in {os_file_choice} data. Check area column.")
                    else:
                        os_branch_select_all = st.checkbox("Select All OS Branches", value=True, key='od_os_branch_all')
                        if os_branch_select_all:
                            selected_os_branches = os_branches
                        else:
                            selected_os_branches = st.multiselect("Select OS Branches", os_branches, key='od_os_branches')
                    st.subheader("Executive Selection")
                    os_executives = sorted(set([e for e in os_df[os_exec_col].apply(extract_executive_name) if e]))
                    
                    os_exec_select_all = st.checkbox("Select All OS Executives", value=True, key='od_os_exec_all')
                    if os_exec_select_all:
                        selected_os_executives = os_executives
                    else:
                        selected_os_executives = st.multiselect("Select OS Executives", os_executives, key='od_executives')
                    if st.button("Generate OD Target Report", key='nbc_od_generate'):
                        if selected_years and till_month:
                            od_target_df, start_date, end_date = filter_os_qty(
                                os_df, os_area_col, os_qty_col, os_due_date_col, os_exec_col,
                                selected_branches=selected_os_branches,
                                selected_years=selected_years,
                                till_month=till_month,
                                selected_executives=selected_os_executives
                            )
                            
                            if od_target_df is not None:
                                start_str = start_date.strftime('%b %Y') if start_date else 'All Periods'
                                end_str = end_date.strftime('%b %Y') if end_date else 'All Periods'                                
                                od_title = f"OD Target-{end_str}(Value in Lakhs)"
                                st.subheader(od_title)
                                st.dataframe(od_target_df)                                
                                img_buffer = create_od_table_image(od_target_df, od_title)
                                if img_buffer:
                                    st.image(img_buffer, use_column_width=True)
                                if 'customers_results' not in st.session_state or st.session_state.customers_results is None:
                                    st.session_state.customers_results = []
                                st.session_state.customers_results.append({'df': od_target_df, 'title': od_title})
                                dfs_info.append({'df': od_target_df, 'title': od_title})
                        else:
                            st.error("Please select at least one year and a month.")
            except Exception as e:
                st.error(f"Error in OD Target tab: {e}")
                logger.error(f"Error in OD Target tab: {e}")
    st.divider()
    st.subheader("Combined Report")    
    if hasattr(st.session_state, 'nbc_results') and hasattr(st.session_state, 'od_results'):
        try:
            nbc_data = st.session_state.nbc_results
            od_data = st.session_state.od_results
            
            ppt_buffer = create_combined_nbc_od_ppt(
                nbc_data['df'],
                nbc_data['title'],
                nbc_data['sorted_months'],
                od_data['df'],
                od_data['title'],
                st.session_state.logo_file
            )            
            if ppt_buffer:
                st.download_button(
                    label="Download Combined NBC & OD Target PPT",
                    data=ppt_buffer,
                    file_name="NBC_OD_Target_Combined.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key='nbc_od_combined_download'
                )
        except Exception as e:
            st.error(f"Error creating combined report: {e}")
            logger.error(f"Error creating combined report: {e}")
        if hasattr(st.session_state, 'nbc_results') and isinstance(st.session_state.nbc_results, dict) and 'df' in st.session_state.nbc_results:
            if 'customers_results' not in st.session_state or not st.session_state.customers_results:
                st.session_state.customers_results = []
            nbc_data = {'df': st.session_state.nbc_results['df'], 'title': st.session_state.nbc_results.get('title', 'Number of Billed Customers')}
            if not any(d.get('title') == nbc_data['title'] for d in st.session_state.customers_results):
                st.session_state.customers_results.append(nbc_data)
        if hasattr(st.session_state, 'od_results') and isinstance(st.session_state.od_results, dict) and 'df' in st.session_state.od_results:
            if 'customers_results' not in st.session_state or not st.session_state.customers_results:
                st.session_state.customers_results = []
            od_data = {'df': st.session_state.od_results['df'], 'title': st.session_state.od_results.get('title', 'OD Target')}
            if not any(d.get('title') == od_data['title'] for d in st.session_state.customers_results):
                st.session_state.customers_results.append(od_data)                
    elif hasattr(st.session_state, 'nbc_results'):
        st.info("Only Number of Billed Customers report is available. Generate OD Target report to enable combined download.")
    elif hasattr(st.session_state, 'od_results'):
        st.info("Only OD Target report is available. Generate Number of Billed Customers report to enable combined download.")
    else:
        st.info("Generate both Number of Billed Customers and OD Target reports to enable combined download.")    
    return dfs_info
def validate_numeric_column(df, col_name, file_name):
    """Validate that a column contains numeric data."""
    try:
        df[col_name] = pd.to_numeric(df[col_name], errors='coerce')
        if df[col_name].isna().all():
            return False, f"Column '{col_name}' in {file_name} contains no valid numeric data."
        return True, None
    except Exception as e:
        return False, f"Column '{col_name}' in {file_name} contains non-numeric data: {e}"
def get_available_months_from_sales(total_sale, sale_bill_date_col):
    """Get unique months from sales bill date column only."""
    months = set()    
    if sale_bill_date_col in total_sale.columns:
        total_sale[sale_bill_date_col] = pd.to_datetime(total_sale[sale_bill_date_col], errors='coerce')
        valid_dates = total_sale[sale_bill_date_col].dropna()
        month_years = valid_dates.dt.strftime('%b-%y').unique()
        months.update(month_years)    
    months = sorted(list(months), key=lambda x: pd.to_datetime("01-" + x, format="%d-%b-%y"))
    return months
def get_cumulative_branches(os_first, os_second, total_sale, 
                         os_first_unit_col, os_second_unit_col, sale_branch_col):
    """Get cumulative branches from all files."""
    all_branches = set()
    if os_first_unit_col in os_first.columns:
        os_first_branches = os_first[os_first_unit_col].dropna().apply(lambda x: map_branch(x, 'title')).unique()
        all_branches.update(os_first_branches)
    if os_second_unit_col in os_second.columns:
        os_second_branches = os_second[os_second_unit_col].dropna().apply(lambda x: map_branch(x, 'title')).unique()
        all_branches.update(os_second_branches)
    if sale_branch_col in total_sale.columns:
        sale_branches = total_sale[sale_branch_col].dropna().apply(lambda x: map_branch(x, 'title')).unique()
        all_branches.update(sale_branches)
    all_branches = sorted([b for b in all_branches if b and b != 'Unknown'])
    return all_branches
def get_cumulative_regions(os_first, os_second, total_sale, 
                         os_first_region_col, os_second_region_col, sale_region_col):
   """Get cumulative regions from all files."""
   all_regions = set()
   if os_first_region_col and os_first_region_col in os_first.columns:
       os_first_regions = os_first[os_first_region_col].dropna().astype(str).str.strip().unique()
       all_regions.update(os_first_regions)
   if os_second_region_col and os_second_region_col in os_second.columns:
       os_second_regions = os_second[os_second_region_col].dropna().astype(str).str.strip().unique()
       all_regions.update(os_second_regions)
   if sale_region_col and sale_region_col in total_sale.columns:
       sale_regions = total_sale[sale_region_col].dropna().astype(str).str.strip().unique()
       all_regions.update(sale_regions)
   all_regions = sorted([r for r in all_regions if r and r.strip() != ''])
   return all_regions
def create_region_branch_mapping(os_first, os_second, total_sale,
                               os_first_unit_col, os_first_region_col,
                               os_second_unit_col, os_second_region_col,
                               sale_branch_col, sale_region_col):
    """Create dynamic region to branch mapping from files."""
    region_branch_mapping = {}
    mappings_data = []
    if (os_first_region_col and os_first_region_col in os_first.columns and 
        os_first_unit_col in os_first.columns):
        os_first_mapping = os_first[[os_first_unit_col, os_first_region_col]].dropna()
        os_first_mapping['Branch'] = os_first_mapping[os_first_unit_col].apply(lambda x: map_branch(x, 'title'))
        os_first_mapping['Region'] = os_first_mapping[os_first_region_col].astype(str).str.strip()
        mappings_data.append(os_first_mapping[['Branch', 'Region']])
    if (os_second_region_col and os_second_region_col in os_second.columns and 
        os_second_unit_col in os_second.columns):
        os_second_mapping = os_second[[os_second_unit_col, os_second_region_col]].dropna()
        os_second_mapping['Branch'] = os_second_mapping[os_second_unit_col].apply(lambda x: map_branch(x, 'title'))
        os_second_mapping['Region'] = os_second_mapping[os_second_region_col].astype(str).str.strip()
        mappings_data.append(os_second_mapping[['Branch', 'Region']])
    if (sale_region_col and sale_region_col in total_sale.columns and 
        sale_branch_col in total_sale.columns):
        sale_mapping = total_sale[[sale_branch_col, sale_region_col]].dropna()
        sale_mapping['Branch'] = sale_mapping[sale_branch_col].apply(lambda x: map_branch(x, 'title'))
        sale_mapping['Region'] = sale_mapping[sale_region_col].astype(str).str.strip()
        mappings_data.append(sale_mapping[['Branch', 'Region']])
    if mappings_data:
        all_mappings = pd.concat(mappings_data, ignore_index=True).drop_duplicates()
        for region, group in all_mappings.groupby('Region'):
            if region and region.strip():  # Skip empty regions
                branches = group['Branch'].unique().tolist()
                branches = [b for b in branches if b and b != 'Unknown']  # Filter valid branches
                if branches:  # Only add if there are valid branches
                    region_branch_mapping[region.strip()] = sorted(branches)
    
    return region_branch_mapping
def create_dynamic_regional_summary(final_df, region_branch_mapping):
    """Create regional summary table using dynamic region-branch mapping."""
    if final_df.empty or not region_branch_mapping:
        return None
    
    # Create a copy of the input DataFrame
    df = final_df.copy()
    
    # Remove the TOTAL row if present
    if not df.empty and df.iloc[-1]['Branch'] == 'TOTAL':
        df = df[:-1].copy()
    
    # Create branch-to-region mapping
    branch_to_region = {}
    for region, branches in region_branch_mapping.items():
        for branch in branches:
            branch_to_region[branch] = region
    
    # Map branches to regions
    df['Region'] = df['Branch'].map(branch_to_region)
    df = df[df['Region'].notna()]
    
    if df.empty:
        return None
    
    # Define numeric columns
    numeric_cols = ["Due Target", "Collection Achieved", "For the month Overdue", "For the month Collection"]
    
    # Group by region and sum numeric columns (values are already in lakhs)
    regional_summary = df.groupby('Region')[numeric_cols].sum().reset_index()
    
    # Calculate percentage columns
    regional_summary["Overall % Achieved"] = np.where(
        regional_summary["Due Target"] > 0,
        (regional_summary["Collection Achieved"] / regional_summary["Due Target"]) * 100,
        0
    )
    regional_summary["% Achieved (Selected Month)"] = np.where(
        regional_summary["For the month Overdue"] > 0,
        (regional_summary["For the month Collection"] / regional_summary["For the month Overdue"]) * 100,
        0
    )
    
    # Round all columns to 2 decimal places consistently
    round_cols = numeric_cols + ["Overall % Achieved", "% Achieved (Selected Month)"]
    regional_summary[round_cols] = regional_summary[round_cols].round(2)
    
    # Create total row
    total_row = {'Region': 'TOTAL'}
    for col in numeric_cols:
        total_row[col] = round(regional_summary[col].sum(), 2)
    total_row["Overall % Achieved"] = round(
        regional_summary["Overall % Achieved"].replace([np.inf, -np.inf], 0).mean(), 2
    )
    total_row["% Achieved (Selected Month)"] = round(
        regional_summary["% Achieved (Selected Month)"].replace([np.inf, -np.inf], 0).mean(), 2
    )
    
    # Append total row
    regional_summary = pd.concat([regional_summary, pd.DataFrame([total_row])], ignore_index=True)
    
    # Reorder columns
    regional_summary = regional_summary[["Region", "Due Target", "Collection Achieved", "Overall % Achieved",
                                        "For the month Overdue", "For the month Collection", "% Achieved (Selected Month)"]]
    
    return regional_summary
def calculate_od_values_updated(os_first, os_second, total_sale, selected_month_str,
                              os_first_due_date_col, os_first_ref_date_col, os_first_unit_col, os_first_net_value_col, os_first_exec_col, os_first_region_col,
                              os_second_due_date_col, os_second_ref_date_col, os_second_unit_col, os_second_net_value_col, os_second_exec_col, os_second_region_col,
                              sale_bill_date_col, sale_due_date_col, sale_branch_col, sale_value_col, sale_exec_col, sale_region_col,
                              selected_executives, selected_branches, selected_regions):
   """Calculate OD Target vs Collection metrics with updated filtering and dynamic regional summary."""
   os_first = os_first.copy()
   os_second = os_second.copy()
   total_sale = total_sale.copy()
   
   # Validate numeric columns
   for df, col, file in [
       (os_first, os_first_net_value_col, "OS-First"),
       (os_second, os_second_net_value_col, "OS-Second"),
       (total_sale, sale_value_col, "Total Sale")
   ]:
       try:
           df[col] = pd.to_numeric(df[col], errors='coerce')
           if df[col].isna().all():
               st.error(f"Column '{col}' in {file} contains no valid numeric data.")
               return None, None, None
       except Exception as e:
           st.error(f"Error processing column '{col}' in {file}: {e}")
           return None, None, None

   # Filter out negative values
   os_first_initial_rows = os_first.shape[0]
   os_first = os_first[os_first[os_first_net_value_col] >= 0]
   os_first_filtered_rows = os_first.shape[0]
   logger.debug(f"OS-First: Filtered out {os_first_initial_rows - os_first_filtered_rows} rows with negative Net Value")

   os_second_initial_rows = os_second.shape[0]
   os_second = os_second[os_second[os_second_net_value_col] >= 0]
   os_second_filtered_rows = os_second.shape[0]
   logger.debug(f"OS-Second: Filtered out {os_second_initial_rows - os_second_filtered_rows} rows with negative Net Value")

   # Convert date columns and map branches
   os_first[os_first_due_date_col] = pd.to_datetime(os_first[os_first_due_date_col], errors='coerce')
   os_first[os_first_ref_date_col] = pd.to_datetime(os_first.get(os_first_ref_date_col), errors='coerce')
   os_first["Branch"] = os_first[os_first_unit_col].apply(map_branch, case='title')

   os_second[os_second_due_date_col] = pd.to_datetime(os_second[os_second_due_date_col], errors='coerce')
   os_second[os_second_ref_date_col] = pd.to_datetime(os_second.get(os_second_ref_date_col), errors='coerce')
   os_second["Branch"] = os_second[os_second_unit_col].apply(map_branch, case='title')

   total_sale[sale_bill_date_col] = pd.to_datetime(total_sale[sale_bill_date_col], errors='coerce')
   total_sale[sale_due_date_col] = pd.to_datetime(total_sale[sale_due_date_col], errors='coerce')
   total_sale["Branch"] = total_sale[sale_branch_col].apply(map_branch, case='title')

   # Create region-branch mapping BEFORE any further branch name changes
   region_branch_mapping = create_region_branch_mapping(
       os_first, os_second, total_sale,
       os_first_unit_col, os_first_region_col,
       os_second_unit_col, os_second_region_col,
       sale_branch_col, sale_region_col
   )

   # Apply executive filtering
   if selected_executives:
       os_first = os_first[os_first[os_first_exec_col].isin(selected_executives)]
       os_second = os_second[os_second[os_second_exec_col].isin(selected_executives)]
       total_sale = total_sale[total_sale[sale_exec_col].isin(selected_executives)]

   # Check for empty dataframes after executive filtering
   empty_dfs = []
   if os_first.empty:
       empty_dfs.append("OS-First")
   if os_second.empty:
       empty_dfs.append("OS-Second")
   if total_sale.empty:
       empty_dfs.append("Total Sale")
   
   if empty_dfs:
       st.error(f"No data remains after executive filtering for: {', '.join(empty_dfs)}")
       return None, None, None

   # Apply branch filtering
   if selected_branches:
       os_first = os_first[os_first["Branch"].isin(selected_branches)]
       os_second = os_second[os_second["Branch"].isin(selected_branches)]
       total_sale = total_sale[total_sale["Branch"].isin(selected_branches)]

   # Apply region filtering
   if selected_regions and region_branch_mapping:
       branches_in_selected_regions = []
       for region in selected_regions:
           if region in region_branch_mapping:
               branches_in_selected_regions.extend(region_branch_mapping[region])      
       if branches_in_selected_regions:
           os_first = os_first[os_first["Branch"].isin(branches_in_selected_regions)]
           os_second = os_second[os_second["Branch"].isin(branches_in_selected_regions)]
           total_sale = total_sale[total_sale["Branch"].isin(branches_in_selected_regions)]

   # Check for empty dataframes after all filtering
   empty_dfs = []
   if os_first.empty:
       empty_dfs.append("OS-First")
   if os_second.empty:
       empty_dfs.append("OS-Second")
   if total_sale.empty:
       empty_dfs.append("Total Sale")   
   if empty_dfs:
       st.error(f"No data remains after filtering for: {', '.join(empty_dfs)}")
       return None, None, None

   # Calculate metrics
   specified_date = pd.to_datetime("01-" + selected_month_str, format="%d-%b-%y")
   specified_month_end = specified_date + pd.offsets.MonthEnd(0)

   # Due target calculation
   due_target = os_first[os_first[os_first_due_date_col] <= specified_month_end]
   due_target_sum = due_target.groupby("Branch")[os_first_net_value_col].sum().reset_index()
   due_target_sum.columns = ["Branch", "Due Target"]

   # Collection calculation
   os_first_coll = os_first[os_first[os_first_due_date_col] <= specified_month_end]
   os_first_coll_sum = os_first_coll.groupby("Branch")[os_first_net_value_col].sum().reset_index()
   os_first_coll_sum.columns = ["Branch", "OS Jan Coll"]

   os_second_coll = os_second[(os_second[os_second_due_date_col] <= specified_month_end) & 
                             (os_second[os_second_ref_date_col] < specified_date)]
   os_second_coll_sum = os_second_coll.groupby("Branch")[os_second_net_value_col].sum().reset_index()
   os_second_coll_sum.columns = ["Branch", "OS Feb Coll"]

   collection = os_first_coll_sum.merge(os_second_coll_sum, on="Branch", how="outer").fillna(0)
   collection["Collection Achieved"] = collection["OS Jan Coll"] - collection["OS Feb Coll"]
   collection["Overall % Achieved"] = np.where(collection["OS Jan Coll"] > 0, 
                                              (collection["Collection Achieved"] / collection["OS Jan Coll"]) * 100, 
                                              0)
   collection = collection.merge(due_target_sum[["Branch", "Due Target"]], on="Branch", how="outer").fillna(0)

   # Overdue calculation
   overdue = total_sale[(total_sale[sale_bill_date_col].between(specified_date, specified_month_end)) & 
                       (total_sale[sale_due_date_col].between(specified_date, specified_month_end))]
   overdue_sum = overdue.groupby("Branch")[sale_value_col].sum().reset_index()
   overdue_sum.columns = ["Branch", "For the month Overdue"]

   # Sale value calculation
   sale_value = overdue.groupby("Branch")[sale_value_col].sum().reset_index()
   sale_value.columns = ["Branch", "Sale Value"]

   # Month collection calculation
   os_second_month = os_second[(os_second[os_second_ref_date_col].between(specified_date, specified_month_end)) & 
                              (os_second[os_second_due_date_col].between(specified_date, specified_month_end))]
   os_second_month_sum = os_second_month.groupby("Branch")[os_second_net_value_col].sum().reset_index()
   os_second_month_sum.columns = ["Branch", "OS Month Collection"]

   month_collection = sale_value.merge(os_second_month_sum, on="Branch", how="outer").fillna(0)
   month_collection["For the month Collection"] = month_collection["Sale Value"] - month_collection["OS Month Collection"]
   month_collection_final = month_collection[["Branch", "For the month Collection"]]

   # Final merge
   final = collection.drop(columns=["OS Jan Coll", "OS Feb Coll"]).merge(overdue_sum, on="Branch", how="outer")\
           .merge(month_collection_final, on="Branch", how="outer").fillna(0)
   
   final["% Achieved (Selected Month)"] = np.where(final["For the month Overdue"] > 0, 
                                                  (final["For the month Collection"] / final["For the month Overdue"]) * 100, 
                                                  0)

   # *** REMOVED THE PROBLEMATIC LINE ***
   # final["Branch"] = final["Branch"].replace({"Puducherry": "Pondicherry"})

   # Convert to lakhs and round
   val_cols = ["Due Target", "Collection Achieved", "For the month Overdue", "For the month Collection"]
   final[val_cols] = final[val_cols].div(100000)
   round_cols = val_cols + ["Overall % Achieved", "% Achieved (Selected Month)"]
   final[round_cols] = final[round_cols].round(2)

   # Reorder columns
   final = final[["Branch", "Due Target", "Collection Achieved", "Overall % Achieved", 
                 "For the month Overdue", "For the month Collection", "% Achieved (Selected Month)"]]
   
   final.sort_values("Branch", inplace=True)

   # Calculate regional summary BEFORE adding total row
   regional_summary = create_dynamic_regional_summary(final, region_branch_mapping)

   # Add total row
   total_row = {'Branch': 'TOTAL'}
   for col in final.columns[1:]:
       if col in ["Overall % Achieved", "% Achieved (Selected Month)"]:
           avg_val = final[col].replace([np.inf, -np.inf], 0).mean()
           total_row[col] = round(avg_val, 2)
       else:
           total_row[col] = round(final[col].sum(), 2)
   
   final = pd.concat([final, pd.DataFrame([total_row])], ignore_index=True)
   
   return final, regional_summary, region_branch_mapping
def create_od_ppt_updated(df, regional_df, title, logo_file=None):
   """Create a PPT for OD Target vs Collection Report with regional summary."""
   try:
       prs = Presentation()
       prs.slide_width = Inches(13.33)
       prs.slide_height = Inches(7.5)
       create_title_slide(prs, title, logo_file)
       add_table_slide(prs, df, f"Branch-wise Performance - {title}", percent_cols=[3, 6])
       if regional_df is not None and not regional_df.empty:
           add_table_slide(prs, regional_df, f"Regional Summary - {title}", percent_cols=[3, 6])
       ppt_buffer = BytesIO()
       prs.save(ppt_buffer)
       ppt_buffer.seek(0)
       return ppt_buffer
   except Exception as e:
       logger.error(f"Error creating OD PPT: {e}")
       st.error(f"Error creating OD PPT: {e}")
       return None
def tab_od_target():
    """OD Target vs Collection Report Tab with AUTO-MAPPING."""
    st.header("OD Target vs Collection Report")   
    dfs_info = []   
    os_jan_file = st.session_state.os_jan_file
    os_feb_file = st.session_state.os_feb_file
    sales_file = st.session_state.sales_file   
    if not os_jan_file or not os_feb_file or not sales_file:
        st.warning("Please upload all required files in the sidebar (OS Previous Month, OS Current Month and Sales files)")
        return dfs_info   
    try:
        os_jan_sheets = get_excel_sheets(os_jan_file)
        os_feb_sheets = get_excel_sheets(os_feb_file)
        sales_sheets = get_excel_sheets(sales_file)       
        if not os_jan_sheets or not os_feb_sheets or not sales_sheets:
            st.error("No sheets found in one or more files")
            return dfs_info
        st.subheader("Sheet Selection")
        col1, col2, col3 = st.columns(3)
        with col1:
            os_jan_sheet = st.selectbox("OS-Previous Month Sheet", os_jan_sheets, 
                                        index=os_jan_sheets.index('Sheet1') if 'Sheet1' in os_jan_sheets else 0, 
                                        key='od_os_jan_sheet')
            os_jan_header = st.number_input("OS-Previous Month Header Row (1-based)", min_value=1, max_value=11, value=1, key='od_os_jan_header') - 1
        with col2:
            os_feb_sheet = st.selectbox("OS-Current Month Sheet", os_feb_sheets, 
                                        index=os_feb_sheets.index('Sheet1') if 'Sheet1' in os_feb_sheets else 0, 
                                        key='od_os_feb_sheet')
            os_feb_header = st.number_input("OS-Current Month Header Row (1-based)", min_value=1, max_value=11, value=1, key='od_os_feb_header') - 1
        with col3:
            sales_sheet = st.selectbox("Sales Sheet", sales_sheets, 
                                        index=sales_sheets.index('Sheet1') if 'Sheet1' in sales_sheets else 0, 
                                        key='od_sales_sheet')
            sales_header = st.number_input("Sales Header Row (1-based)", min_value=1, max_value=11, value=1, key='od_sales_header') - 1
        os_jan = pd.read_excel(os_jan_file, sheet_name=os_jan_sheet, header=os_jan_header)
        os_feb = pd.read_excel(os_feb_file, sheet_name=os_feb_sheet, header=os_feb_header)
        total_sale = pd.read_excel(sales_file, sheet_name=sales_sheet, header=sales_header)
        os_jan_cols = os_jan.columns.tolist()
        os_feb_cols = os_feb.columns.tolist()
        sales_cols = total_sale.columns.tolist()
        os_first_mapping, os_second_mapping, sales_mapping = auto_map_od_columns(os_jan_cols, os_feb_cols, sales_cols)
        st.subheader("OS-Previous Month Column Mapping")
        col1, col2, col3 = st.columns(3)
        with col1:
            os_jan_due_date_col = st.selectbox("Due Date Column", os_jan_cols, 
                                                index=os_jan_cols.index(os_first_mapping['due_date']) if os_first_mapping['due_date'] else 0, 
                                                key='od_os_jan_due_date')
            os_jan_ref_date_col = st.selectbox("Reference Date Column", os_jan_cols, 
                                                index=os_jan_cols.index(os_first_mapping['ref_date']) if os_first_mapping['ref_date'] else 0, 
                                                key='od_os_jan_ref_date')
        with col2:
            os_jan_unit_col = st.selectbox("Branch Column", os_jan_cols, 
                                            index=os_jan_cols.index(os_first_mapping['branch']) if os_first_mapping['branch'] else 0, 
                                            key='od_os_jan_unit')
            os_jan_net_value_col = st.selectbox("Net Value Column", os_jan_cols, 
                                                index=os_jan_cols.index(os_first_mapping['net_value']) if os_first_mapping['net_value'] else 0, 
                                                key='od_os_jan_net_value')
        with col3:
            os_jan_exec_col = st.selectbox("Executive Column", os_jan_cols, 
                                            index=os_jan_cols.index(os_first_mapping['executive']) if os_first_mapping['executive'] else 0, 
                                            key='od_os_jan_exec')
            os_jan_region_col = st.selectbox("Region Column", ["None"] + os_jan_cols, 
                                            index=os_jan_cols.index(os_first_mapping['region']) + 1 if os_first_mapping['region'] else 0, 
                                            key='od_os_jan_region')
            os_jan_region_col = None if os_jan_region_col == "None" else os_jan_region_col      
        st.subheader("OS-Current Month Column Mapping")
        col1, col2, col3 = st.columns(3)
        with col1:
            os_feb_due_date_col = st.selectbox("Due Date Column", os_feb_cols, 
                                                index=os_feb_cols.index(os_second_mapping['due_date']) if os_second_mapping['due_date'] else 0, 
                                                key='od_os_feb_due_date')
            os_feb_ref_date_col = st.selectbox("Reference Date Column", os_feb_cols, 
                                                index=os_feb_cols.index(os_second_mapping['ref_date']) if os_second_mapping['ref_date'] else 0, 
                                                key='od_os_feb_ref_date')
        with col2:
            os_feb_unit_col = st.selectbox("Branch Column", os_feb_cols, 
                                            index=os_feb_cols.index(os_second_mapping['branch']) if os_second_mapping['branch'] else 0, 
                                            key='od_os_feb_unit')
            os_feb_net_value_col = st.selectbox("Net Value Column", os_feb_cols, 
                                                index=os_feb_cols.index(os_second_mapping['net_value']) if os_second_mapping['net_value'] else 0, 
                                                key='od_os_feb_net_value')
        with col3:
            os_feb_exec_col = st.selectbox("Executive Column", os_feb_cols, 
                                            index=os_feb_cols.index(os_second_mapping['executive']) if os_second_mapping['executive'] else 0, 
                                            key='od_os_feb_exec')
            os_feb_region_col = st.selectbox("Region Column", ["None"] + os_feb_cols, 
                                            index=os_feb_cols.index(os_second_mapping['region']) + 1 if os_second_mapping['region'] else 0, 
                                            key='od_os_feb_region')
            os_feb_region_col = None if os_feb_region_col == "None" else os_feb_region_col       
        st.subheader("Total Sale Column Mapping")
        col1, col2, col3 = st.columns(3)
        with col1:
            sale_bill_date_col = st.selectbox("Bill Date Column", sales_cols, 
                                                index=sales_cols.index(sales_mapping['bill_date']) if sales_mapping['bill_date'] else 0, 
                                                key='od_sale_bill_date')
            sale_due_date_col = st.selectbox("Due Date Column", sales_cols, 
                                            index=sales_cols.index(sales_mapping['due_date']) if sales_mapping['due_date'] else 0, 
                                            key='od_sale_due_date')
        with col2:
            sale_branch_col = st.selectbox("Branch Column", sales_cols, 
                                            index=sales_cols.index(sales_mapping['branch']) if sales_mapping['branch'] else 0, 
                                            key='od_sale_branch')
            sale_value_col = st.selectbox("Value Column", sales_cols, 
                                            index=sales_cols.index(sales_mapping['value']) if sales_mapping['value'] else 0, 
                                            key='od_sale_value')
        with col3:
            sale_exec_col = st.selectbox("Executive Column", sales_cols, 
                                        index=sales_cols.index(sales_mapping['executive']) if sales_mapping['executive'] else 0, 
                                        key='od_sale_exec')
            sale_region_col = st.selectbox("Region Column", ["None"] + sales_cols, 
                                            index=sales_cols.index(sales_mapping['region']) + 1 if sales_mapping['region'] else 0, 
                                            key='od_sale_region')
            sale_region_col = None if sale_region_col == "None" else sale_region_col
        available_months = get_available_months_from_sales(total_sale, sale_bill_date_col)
        
        if not available_months:
            st.error("No valid months found in the sales bill date column. Please check column selection.")
            return dfs_info          
        selected_month_str = st.selectbox("Select Sales Month", available_months, index=len(available_months)-1, key='od_month')
        cumulative_branches = get_cumulative_branches(
            os_jan, os_feb, total_sale, 
            os_jan_unit_col, os_feb_unit_col, sale_branch_col
        )
        cumulative_executives = set()
        cumulative_executives.update(os_jan[os_jan_exec_col].dropna().unique())
        cumulative_executives.update(os_feb[os_feb_exec_col].dropna().unique())
        cumulative_executives.update(total_sale[sale_exec_col].dropna().unique())
        cumulative_executives = sorted(list(cumulative_executives))
        cumulative_regions = get_cumulative_regions(
            os_jan, os_feb, total_sale,
            os_jan_region_col, os_feb_region_col, sale_region_col
        )
        if cumulative_regions:
            st.subheader("Region-Branch Mapping Preview")
            with st.expander("View Region-Branch Relationships", expanded=False):
                preview_mapping = create_region_branch_mapping(
                    os_jan, os_feb, total_sale,
                    os_jan_unit_col, os_jan_region_col,
                    os_feb_unit_col, os_feb_region_col,
                    sale_branch_col, sale_region_col
                )
                
                if preview_mapping:
                    for region, branches in preview_mapping.items():
                        st.write(f"**{region}:** {', '.join(branches)}")
                else:
                    st.write("No region-branch mapping could be created from the selected columns.")
        st.subheader("Filter Options")
        filter_tabs = st.tabs(["Cumulative Branches", "Cumulative Executives", "Dynamic Regions"])
        with filter_tabs[0]:
            branch_select_all = st.checkbox("Select All Branches", value=True, key='od_branch_all')
            if branch_select_all:
                selected_branches = cumulative_branches
            else:
                selected_branches = st.multiselect("Select Branches", cumulative_branches, key='od_branches')
        with filter_tabs[1]:
            exec_select_all = st.checkbox("Select All Executives", value=True, key='od_exec_all')
            if exec_select_all:
                selected_executives = cumulative_executives
            else:
                selected_executives = st.multiselect("Select Executives", cumulative_executives, key='od_executives')
        with filter_tabs[2]:
            if cumulative_regions:
                region_select_all = st.checkbox("Select All Regions", value=True, key='od_region_all')
                if region_select_all:
                    selected_regions = cumulative_regions
                else:
                    selected_regions = st.multiselect("Select Regions", cumulative_regions, key='od_regions')
            else:
                st.info("No region columns selected or no valid regions found in the data.")
                selected_regions = []
        if st.button("Generate OD Target vs Collection Report", key='od_generate'):
            final, regional_summary, region_mapping = calculate_od_values_updated(
                os_jan, os_feb, total_sale, selected_month_str,
                os_jan_due_date_col, os_jan_ref_date_col, os_jan_unit_col, os_jan_net_value_col, os_jan_exec_col, os_jan_region_col,
                os_feb_due_date_col, os_feb_ref_date_col, os_feb_unit_col, os_feb_net_value_col, os_feb_exec_col, os_feb_region_col,
                sale_bill_date_col, sale_due_date_col, sale_branch_col, sale_value_col, sale_exec_col, sale_region_col,
                selected_executives, selected_branches, selected_regions
            )          
            if final is not None and not final.empty:
                st.subheader("Results")
                if regional_summary is not None and not regional_summary.empty:
                    result_tabs = st.tabs(["Branch-wise Performance", "Regional Summary", "Region Mapping"])
                else:
                    result_tabs = st.tabs(["Branch-wise Performance", "Region Mapping"])
                
                with result_tabs[0]:
                    st.write("**Branch-wise OD Target vs Collection**")
                    st.dataframe(final)
                    img_buffer = create_table_image(final, f"OD TARGET VS COLLECTION (Branch-wise) - {selected_month_str}", percent_cols=[3, 6])
                    if img_buffer:
                        st.image(img_buffer, use_column_width=True)              
                tab_idx = 1
                if regional_summary is not None and not regional_summary.empty:
                    with result_tabs[1]:
                        st.write("**Regional Summary**")
                        st.dataframe(regional_summary)
                        img_buffer = create_table_image(regional_summary, f"OD TARGET VS COLLECTION (Regional Summary) - {selected_month_str}", percent_cols=[3, 6])
                        if img_buffer:
                            st.image(img_buffer, use_column_width=True)
                    tab_idx = 2
                with result_tabs[tab_idx]:
                    st.write("**Region-Branch Mapping Used in Analysis**")
                    if region_mapping:
                        mapping_df = pd.DataFrame([
                            {"Region": region, "Branches": ", ".join(branches)}
                            for region, branches in region_mapping.items()
                        ])
                        st.dataframe(mapping_df, use_container_width=True)
                    else:
                        st.write("No region mapping was created or used.")
                dfs_info = [
                    {'df': final, 'title': f"OD TARGET VS COLLECTION (Branch-wise) - {selected_month_str} (Value in Lakhs)", 'percent_cols': [3, 6]}
                ]
                
                if regional_summary is not None and not regional_summary.empty:
                    dfs_info.append({
                        'df': regional_summary, 
                        'title': f"OD TARGET VS COLLECTION (Regional Summary) - {selected_month_str} (Value in Lakhs)", 
                        'percent_cols': [3, 6]
                    })
                st.session_state.od_results = dfs_info
                ppt_buffer = create_od_ppt_updated(
                    final, 
                    regional_summary,
                    f"OD Target vs Collection - {selected_month_str}",
                    st.session_state.logo_file
                )              
                if ppt_buffer:
                    unique_id = str(uuid.uuid4())[:8]
                    st.download_button(
                        label="Download OD Target vs Collection PPT",
                        data=ppt_buffer,
                        file_name=f"OD_Target_vs_Collection_{selected_month_str}_{unique_id}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"od_download_{unique_id}"
                    )
            else:
                st.error("Failed to generate OD Target vs Collection report. Check your data and selections.")
    except Exception as e:
                st.error(f"Error in OD Target tab: {e}")
                logger.error(f"Error in OD Target tab: {e}")
        
    return dfs_info

#######################################################
# MODULE 4: PRODUCT GROWTH REPORT (UPDATED)
#######################################################

def auto_map_product_growth_columns(ly_columns, cy_columns, budget_columns):
    """Auto-map Product Growth columns."""
    ly_mapping = {
        'date': find_column_by_names(ly_columns, ['Date', 'Bill Date', 'Invoice Date']),
        'product_group': find_column_by_names(ly_columns, ['Type (Make)', 'Product Group', 'Type', 'Make']),
        'quantity': find_column_by_names(ly_columns, ['Actual Quantity', 'Quantity', 'Qty']),
        'company_group': find_column_by_names(ly_columns, ['Company Group', 'Company', 'Group']),
        'value': find_column_by_names(ly_columns, ['Value', 'Invoice Value', 'Amount']),
        'executive': find_column_by_names(ly_columns, ['Executive', 'Sales Executive', 'Executive Name'])
    }
    
    cy_mapping = {
        'date': find_column_by_names(cy_columns, ['Month Format', 'Date', 'Bill Date', 'Invoice Date']),
        'product_group': find_column_by_names(cy_columns, ['Type (Make)', 'Product Group', 'Type', 'Make']),
        'quantity': find_column_by_names(cy_columns, ['Actual Quantity', 'Quantity', 'Qty']),
        'company_group': find_column_by_names(cy_columns, ['Company Group', 'Company', 'Group']),
        'value': find_column_by_names(cy_columns, ['Value', 'Invoice Value', 'Amount']),
        'executive': find_column_by_names(cy_columns, ['Executive', 'Sales Executive', 'Executive Name'])
    }
    
    budget_mapping = {
        'quantity': find_column_by_names(budget_columns, ["Qty – Apr'25", 'Quantity', 'Qty', 'Budget Qty']),
        'company_group': find_column_by_names(budget_columns, ['Company Group', 'Company', 'Group']),
        'value': find_column_by_names(budget_columns, ["Value – Apr'25", 'Value', 'Budget Value', 'Amount']),
        'executive': find_column_by_names(budget_columns, ['Executive Name', 'Executive', 'Sales Executive']),
        'product_group': find_column_by_names(budget_columns, ['Product Group', 'Type(Make)', 'Product', 'Type'])
    }
    
    return ly_mapping, cy_mapping, budget_mapping

def standardize_name(name):
    """Standardize a name by removing extra spaces, special characters, and converting to title case."""
    if pd.isna(name) or not name:
        return ""
    name = str(name).strip().lower()
    name = ''.join(c for c in name if c.isalnum() or c.isspace())
    name = ' '.join(word.capitalize() for word in name.split())
    # Unify 'General' variants
    general_variants = ['general', 'gen', 'generals', 'general ', 'genral', 'generl']
    if any(variant in name.lower() for variant in general_variants):
        return 'General'
    return name

def log_non_numeric_values(df, col):
    """Log non-numeric values in a column before conversion."""
    if col in df.columns:
        non_numeric = df[df[col].apply(lambda x: not isinstance(x, (int, float)) and pd.notna(x))][col].unique()
        if len(non_numeric) > 0:
            logger.warning(f"Non-numeric values in {col}: {non_numeric}")

def calculate_product_growth(ly_df, cy_df, budget_df, ly_months, cy_months, ly_date_col, cy_date_col, 
                           ly_qty_col, cy_qty_col, ly_value_col, cy_value_col, 
                           budget_qty_col, budget_value_col, ly_product_col, cy_product_col, 
                           ly_company_group_col, cy_company_group_col, budget_company_group_col, budget_product_group_col,
                           ly_exec_col, cy_exec_col, budget_exec_col, selected_executives=None, 
                           selected_company_groups=None):
    """Calculate product growth metrics comparing last year, current year and budget data."""
    try:
        # Debug information to help troubleshoot
        st.write("Processing data...")
        
        # Create copies to avoid modifying original DataFrames
        ly_df = ly_df.copy()
        cy_df = cy_df.copy()
        budget_df = budget_df.copy()

        # Validate columns
        required_cols = [
            (ly_df, [ly_date_col, ly_qty_col, ly_value_col, ly_product_col, ly_company_group_col, ly_exec_col], "Last Year"),
            (cy_df, [cy_date_col, cy_qty_col, cy_value_col, cy_product_col, cy_company_group_col, cy_exec_col], "Current Year"),
            (budget_df, [budget_qty_col, budget_value_col, budget_product_group_col, budget_company_group_col, budget_exec_col], "Budget")
        ]
        
        for df, cols, df_name in required_cols:
            missing_cols = [col for col in cols if col not in df.columns]
            if missing_cols:
                st.error(f"Missing columns in {df_name} data: {missing_cols}")
                return None

        # Apply executive filter
        if selected_executives:
            st.write(f"Applying executive filter: {len(selected_executives)} executives selected")
            if ly_exec_col in ly_df.columns:
                ly_df = ly_df[ly_df[ly_exec_col].isin(selected_executives)]
            if cy_exec_col in cy_df.columns:
                cy_df = cy_df[cy_df[cy_exec_col].isin(selected_executives)]
            if budget_exec_col in budget_df.columns:
                budget_df = budget_df[budget_df[budget_exec_col].isin(selected_executives)]

        if ly_df.empty or cy_df.empty or budget_df.empty:
            empty_dfs = []
            if ly_df.empty: empty_dfs.append("Last Year")
            if cy_df.empty: empty_dfs.append("Current Year")
            if budget_df.empty: empty_dfs.append("Budget")
            
            st.warning(f"No data remains after executive filtering for: {', '.join(empty_dfs)}. Please check executive selections.")
            return None

        # Check and convert date columns
        st.write("Converting date columns...")
        try:
            ly_df[ly_date_col] = pd.to_datetime(ly_df[ly_date_col], dayfirst=True, errors='coerce')
            cy_df[cy_date_col] = pd.to_datetime(cy_df[cy_date_col], dayfirst=True, errors='coerce')
        except Exception as e:
            st.error(f"Error converting date columns: {e}")
            return None
        
        # Check for valid dates
        if ly_df[ly_date_col].isna().all() or cy_df[cy_date_col].isna().all():
            st.error("Date columns contain no valid dates. Check date formats.")
            return None
        
        available_ly_months = ly_df[ly_date_col].dt.strftime('%b %y').dropna().unique().tolist()
        available_cy_months = cy_df[cy_date_col].dt.strftime('%b %y').dropna().unique().tolist()

        st.write(f"Available LY months: {available_ly_months}")
        st.write(f"Available CY months: {available_cy_months}")
        
        if not available_ly_months or not available_cy_months:
            st.error("No valid months extracted from date columns. Please check the date column format.")
            return None

        # Apply month filter
        st.write(f"Filtering for LY months: {ly_months}")
        st.write(f"Filtering for CY months: {cy_months}")
        
        if not ly_months or not cy_months:
            st.error("No months selected for comparison. Please select months for both LY and CY.")
            return None
            
        try:
            ly_filtered_df = ly_df[ly_df[ly_date_col].dt.strftime('%b %y').isin(ly_months)]
            cy_filtered_df = cy_df[cy_df[cy_date_col].dt.strftime('%b %y').isin(cy_months)]
        except Exception as e:
            st.error(f"Error filtering by months: {e}")
            return None

        if ly_filtered_df.empty or cy_filtered_df.empty:
            st.warning(f"No data for selected months (LY: {', '.join(ly_months)}, CY: {', '.join(cy_months)})")
            return None

        # Standardize product and company group names
        st.write("Standardizing product and company group names...")
        
        # Handle possible missing values before applying standardize_name
        ly_filtered_df[ly_product_col] = ly_filtered_df[ly_product_col].fillna("")
        cy_filtered_df[cy_product_col] = cy_filtered_df[cy_product_col].fillna("")
        budget_df[budget_product_group_col] = budget_df[budget_product_group_col].fillna("")
        
        ly_filtered_df[ly_company_group_col] = ly_filtered_df[ly_company_group_col].fillna("")
        cy_filtered_df[cy_company_group_col] = cy_filtered_df[cy_company_group_col].fillna("")
        budget_df[budget_company_group_col] = budget_df[budget_company_group_col].fillna("")
        
        # Apply standardization
        ly_filtered_df[ly_product_col] = ly_filtered_df[ly_product_col].apply(standardize_name)
        cy_filtered_df[cy_product_col] = cy_filtered_df[cy_product_col].apply(standardize_name)
        budget_df[budget_product_group_col] = budget_df[budget_product_group_col].apply(standardize_name)

        ly_filtered_df[ly_company_group_col] = ly_filtered_df[ly_company_group_col].apply(standardize_name)
        cy_filtered_df[cy_company_group_col] = cy_filtered_df[cy_company_group_col].apply(standardize_name)
        budget_df[budget_company_group_col] = budget_df[budget_company_group_col].apply(standardize_name)

        # Apply company group filter
        if selected_company_groups:
            st.write(f"Applying company group filter: {selected_company_groups}")
            selected_company_groups = [standardize_name(g) for g in selected_company_groups]
            
            # Filter data but keep track of what's removed
            ly_before = len(ly_filtered_df)
            cy_before = len(cy_filtered_df)
            budget_before = len(budget_df)
            
            ly_filtered_df = ly_filtered_df[ly_filtered_df[ly_company_group_col].isin(selected_company_groups)]
            cy_filtered_df = cy_filtered_df[cy_filtered_df[cy_company_group_col].isin(selected_company_groups)]
            budget_df = budget_df[budget_df[budget_company_group_col].isin(selected_company_groups)]
            
            st.write(f"Filtered out {ly_before - len(ly_filtered_df)} rows from LY data")
            st.write(f"Filtered out {cy_before - len(cy_filtered_df)} rows from CY data")
            st.write(f"Filtered out {budget_before - len(budget_df)} rows from Budget data")

            if ly_filtered_df.empty or cy_filtered_df.empty or budget_df.empty:
                empty_dfs = []
                if ly_filtered_df.empty: empty_dfs.append("Last Year")
                if cy_filtered_df.empty: empty_dfs.append("Current Year")
                if budget_df.empty: empty_dfs.append("Budget")
                
                st.warning(f"No data remains after filtering for company groups: {selected_company_groups} in {', '.join(empty_dfs)}")
                return None

        # Convert quantity and value columns to numeric
        st.write("Converting quantity and value columns to numeric...")
        for df, qty_col, value_col, df_name in [
            (ly_filtered_df, ly_qty_col, ly_value_col, "Last Year"), 
            (cy_filtered_df, cy_qty_col, cy_value_col, "Current Year"), 
            (budget_df, budget_qty_col, budget_value_col, "Budget")
        ]:
            try:
                df[qty_col] = pd.to_numeric(df[qty_col], errors='coerce').fillna(0)
                df[value_col] = pd.to_numeric(df[value_col], errors='coerce').fillna(0)
                
                # Log some stats to help debug
                st.write(f"{df_name} qty range: {df[qty_col].min()} to {df[qty_col].max()}")
                st.write(f"{df_name} value range: {df[value_col].min()} to {df[value_col].max()}")
                
            except Exception as e:
                st.error(f"Error converting numeric columns in {df_name} data: {e}")
                return None

        # Get unique company groups from all datasets
        st.write("Determining company groups...")
        
        # Print what we have in each dataset
        st.write(f"LY company groups: {ly_filtered_df[ly_company_group_col].unique().tolist()}")
        st.write(f"CY company groups: {cy_filtered_df[cy_company_group_col].unique().tolist()}")
        st.write(f"Budget company groups: {budget_df[budget_company_group_col].unique().tolist()}")
        
        company_groups = selected_company_groups if selected_company_groups else pd.concat([
            ly_filtered_df[ly_company_group_col], 
            cy_filtered_df[cy_company_group_col], 
            budget_df[budget_company_group_col]
        ]).dropna().unique().tolist()
        
        st.write(f"Found {len(company_groups)} company groups: {company_groups}")

        if not company_groups:
            st.warning("No valid company groups found in the data. Please check company group columns.")
            return None

        result = {}
        st.write("Processing data by company group...")
        
        for group in company_groups:
            st.write(f"Processing group: {group}")
            # Filter data by company group
            ly_group_df = ly_filtered_df[ly_filtered_df[ly_company_group_col] == group]
            cy_group_df = cy_filtered_df[cy_filtered_df[cy_company_group_col] == group]
            budget_group_df = budget_df[budget_df[budget_company_group_col] == group]
            
            st.write(f"Group {group}: LY rows: {len(ly_group_df)}, CY rows: {len(cy_group_df)}, Budget rows: {len(budget_group_df)}")
            
            # Check if we have data for this group in all datasets
            if ly_group_df.empty or cy_group_df.empty or budget_group_df.empty:
                st.warning(f"Insufficient data for company group '{group}'. Skipping this group.")
                continue

            # Collect products specific to this company group
            ly_products = ly_group_df[ly_product_col].dropna().apply(standardize_name).unique().tolist()
            cy_products = cy_group_df[cy_product_col].dropna().apply(standardize_name).unique().tolist()
            budget_products = budget_group_df[budget_product_group_col].dropna().apply(standardize_name).unique().tolist()
            group_products = sorted(set(ly_products + cy_products + budget_products))
            
            st.write(f"Group {group}: Found {len(group_products)} products")
            
            # Restrict 'Gc' to 'General'
            if group != 'General':
                group_products = [p for p in group_products if p != 'Gc']

            if not group_products:
                st.warning(f"No products found for company group '{group}'. Skipping this group.")
                continue

            # Aggregate data
            st.write(f"Aggregating data for group {group}...")
            try:
                ly_qty = ly_group_df.groupby(ly_product_col)[ly_qty_col].sum().reset_index().rename(columns={ly_product_col: 'PRODUCT NAME', ly_qty_col: 'LY_QTY'})
                cy_qty = cy_group_df.groupby(cy_product_col)[cy_qty_col].sum().reset_index().rename(columns={cy_product_col: 'PRODUCT NAME', cy_qty_col: 'CY_QTY'})
                ly_value = ly_group_df.groupby(ly_product_col)[ly_value_col].sum().reset_index().rename(columns={ly_product_col: 'PRODUCT NAME', ly_value_col: 'LY_VALUE'})
                cy_value = cy_group_df.groupby(cy_product_col)[cy_value_col].sum().reset_index().rename(columns={cy_product_col: 'PRODUCT NAME', cy_value_col: 'CY_VALUE'})
                budget_qty = budget_group_df.groupby(budget_product_group_col)[budget_qty_col].sum().reset_index().rename(columns={budget_product_group_col: 'PRODUCT NAME', budget_qty_col: 'BUDGET_QTY'})
                budget_value = budget_group_df.groupby(budget_product_group_col)[budget_value_col].sum().reset_index().rename(columns={budget_product_group_col: 'PRODUCT NAME', budget_value_col: 'BUDGET_VALUE'})
            except Exception as e:
                st.error(f"Error aggregating data for group {group}: {e}")
                continue

            # Create a DataFrame with all product groups for this company group
            all_products_df = pd.DataFrame({'PRODUCT NAME': group_products})

            # Merge with left join to include all product groups
            qty_df = all_products_df.merge(ly_qty, on='PRODUCT NAME', how='left').merge(cy_qty, on='PRODUCT NAME', how='left').merge(budget_qty, on='PRODUCT NAME', how='left').fillna(0)
            value_df = all_products_df.merge(ly_value, on='PRODUCT NAME', how='left').merge(cy_value, on='PRODUCT NAME', how='left').merge(budget_value, on='PRODUCT NAME', how='left').fillna(0)

            # Define function for calculating achievement percentage
            def calc_achievement(row, cy_col, ly_col):
                if pd.isna(row[ly_col]) or row[ly_col] == 0:
                    return 0.00 if row[cy_col] == 0 else 100.00
                return ((row[cy_col] - row[ly_col]) / row[ly_col]) * 100

            # Calculate achievement percentages without rounding (we'll round all columns together later)
            qty_df['ACHIEVEMENT %'] = qty_df.apply(lambda row: calc_achievement(row, 'CY_QTY', 'LY_QTY'), axis=1)
            value_df['ACHIEVEMENT %'] = value_df.apply(lambda row: calc_achievement(row, 'CY_VALUE', 'LY_VALUE'), axis=1)

            # Round all numeric columns to 2 decimal places
            numeric_cols_qty = ['LY_QTY', 'BUDGET_QTY', 'CY_QTY', 'ACHIEVEMENT %']
            numeric_cols_value = ['LY_VALUE', 'BUDGET_VALUE', 'CY_VALUE', 'ACHIEVEMENT %']
            
            for col in numeric_cols_qty:
                qty_df[col] = qty_df[col].round(2)
            
            for col in numeric_cols_value:
                value_df[col] = value_df[col].round(2)
            
            # Rename columns to uppercase
            qty_df = qty_df[['PRODUCT NAME', 'LY_QTY', 'BUDGET_QTY', 'CY_QTY', 'ACHIEVEMENT %']]
            value_df = value_df[['PRODUCT NAME', 'LY_VALUE', 'BUDGET_VALUE', 'CY_VALUE', 'ACHIEVEMENT %']]

            # Add totals - make sure totals are also rounded
            qty_totals = pd.DataFrame({
                'PRODUCT NAME': ['TOTAL'],
                'LY_QTY': [qty_df['LY_QTY'].sum().round(2)],
                'BUDGET_QTY': [qty_df['BUDGET_QTY'].sum().round(2)],
                'CY_QTY': [qty_df['CY_QTY'].sum().round(2)],
                'ACHIEVEMENT %': [qty_df['ACHIEVEMENT %'].replace([np.inf, -np.inf], 0).mean().round(2) if len(qty_df) > 1 else 0]
            })
            qty_df = pd.concat([qty_df, qty_totals], ignore_index=True)

            value_totals = pd.DataFrame({
                'PRODUCT NAME': ['TOTAL'],
                'LY_VALUE': [value_df['LY_VALUE'].sum().round(2)],
                'BUDGET_VALUE': [value_df['BUDGET_VALUE'].sum().round(2)],
                'CY_VALUE': [value_df['CY_VALUE'].sum().round(2)],
                'ACHIEVEMENT %': [value_df['ACHIEVEMENT %'].replace([np.inf, -np.inf], 0).mean().round(2) if len(value_df) > 1 else 0]
            })
            value_df = pd.concat([value_df, value_totals], ignore_index=True)

            result[group] = {'qty_df': qty_df, 'value_df': value_df}

        if not result:
            st.error("No data available for report. Please check filters and data content.")
            return None

        # Remove debug statements in production
        st.write("Product growth calculation complete!")
        st.write(f"Generated reports for {len(result)} company groups")
        
        return result
    
    except Exception as e:
        st.error(f"Error in product growth calculation: {e}")
        import traceback
        st.code(traceback.format_exc())
        return None

def create_product_growth_ppt(group_results, month_title, logo_file=None):
    """Create a PPT for Product Growth analysis."""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide
        create_title_slide(prs, f"Product Growth by Company Group – {month_title}", logo_file)
        
        # Add data slides for each company group
        for group, data in group_results.items():
            qty_df = data['qty_df']
            value_df = data['value_df']
            
            add_table_slide(prs, qty_df, f"{group} - Quantity Growth (Qty in Mt)", percent_cols=[4])
            add_table_slide(prs, value_df, f"{group} - Value Growth (Value in Lakhs)", percent_cols=[4])
        
        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating Product Growth PPT: {e}")
        st.error(f"Error creating Product Growth PPT: {e}")
        return None

def tab_product_growth():
    """Product Growth Report Tab with AUTO-MAPPING."""
    st.header("Product Growth Dashboard")
    
    dfs_info = []
    
    sales_file = st.session_state.sales_file
    last_year_sales_file = st.session_state.last_year_sales_file
    budget_file = st.session_state.budget_file
    
    if not sales_file or not last_year_sales_file or not budget_file:
        st.warning("Please upload Sales, Last Year Sales, and Budget files in the sidebar")
        return dfs_info
    
    try:
        # Get sheet names
        ly_sheets = get_excel_sheets(last_year_sales_file)
        cy_sheets = get_excel_sheets(sales_file)
        budget_sheets = get_excel_sheets(budget_file)
        
        if not ly_sheets or not cy_sheets or not budget_sheets:
            st.error("No sheets found in files")
            return dfs_info
            
        # Sheet selection
        st.subheader("Configure Files")
        col1, col2 = st.columns(2)
        with col1:
            ly_sheet = st.selectbox("Last Year Sales Sheet", ly_sheets, 
                                   index=ly_sheets.index('Sheet1') if 'Sheet1' in ly_sheets else 0, 
                                   key='pg_ly_sheet')
            ly_header = st.number_input("Last Year Header Row (1-based)", min_value=1, max_value=11, value=1, key='pg_ly_header') - 1
            cy_sheet = st.selectbox("Current Year Sales Sheet", cy_sheets, 
                                   index=cy_sheets.index('Sheet1') if 'Sheet1' in cy_sheets else 0, 
                                   key='pg_cy_sheet')
            cy_header = st.number_input("Current Year Header Row (1-based)", min_value=1, max_value=11, value=1, key='pg_cy_header') - 1
        with col2:
            budget_sheet = st.selectbox("Budget Sheet", budget_sheets, 
                                      index=budget_sheets.index('Sheet1') if 'Sheet1' in budget_sheets else 0, 
                                      key='pg_budget_sheet')
            budget_header = st.number_input("Budget Header Row (1-based)", min_value=1, max_value=11, value=1, key='pg_budget_header') - 1
        
        # Load data
        ly_df = pd.read_excel(last_year_sales_file, sheet_name=ly_sheet, header=ly_header)
        cy_df = pd.read_excel(sales_file, sheet_name=cy_sheet, header=cy_header)
        budget_df = pd.read_excel(budget_file, sheet_name=budget_sheet, header=budget_header)
        
        # Get columns for auto-mapping
        ly_cols = ly_df.columns.tolist()
        cy_cols = cy_df.columns.tolist()
        budget_cols = budget_df.columns.tolist()
        
        # Auto-map columns
        ly_mapping, cy_mapping, budget_mapping = auto_map_product_growth_columns(ly_cols, cy_cols, budget_cols)
        
        # Column mappings with auto-mapping
        st.subheader("Last Year Column Mapping")
        col1, col2 = st.columns(2)
        with col1:
            ly_date_col = st.selectbox("Date Column", ly_cols, 
                                      index=ly_cols.index(ly_mapping['date']) if ly_mapping['date'] else 0, 
                                      key='pg_ly_date')
            ly_qty_col = st.selectbox("Quantity Column", ly_cols, 
                                     index=ly_cols.index(ly_mapping['quantity']) if ly_mapping['quantity'] else 0, 
                                     key='pg_ly_qty')
            ly_value_col = st.selectbox("Value Column", ly_cols, 
                                       index=ly_cols.index(ly_mapping['value']) if ly_mapping['value'] else 0, 
                                       key='pg_ly_value')
        with col2:
            ly_product_col = st.selectbox("Product Group Column", ly_cols, 
                                         index=ly_cols.index(ly_mapping['product_group']) if ly_mapping['product_group'] else 0, 
                                         key='pg_ly_product')
            ly_company_group_col = st.selectbox("Company Group Column", ly_cols, 
                                               index=ly_cols.index(ly_mapping['company_group']) if ly_mapping['company_group'] else 0, 
                                               key='pg_ly_company_group')
            ly_exec_col = st.selectbox("Executive Column", ly_cols, 
                                      index=ly_cols.index(ly_mapping['executive']) if ly_mapping['executive'] else 0, 
                                      key='pg_ly_exec')
        
        st.subheader("Current Year Column Mapping")
        col1, col2 = st.columns(2)
        with col1:
            cy_date_col = st.selectbox("Date Column", cy_cols, 
                                      index=cy_cols.index(cy_mapping['date']) if cy_mapping['date'] else 0, 
                                      key='pg_cy_date')
            cy_qty_col = st.selectbox("Quantity Column", cy_cols, 
                                     index=cy_cols.index(cy_mapping['quantity']) if cy_mapping['quantity'] else 0, 
                                     key='pg_cy_qty')
            cy_value_col = st.selectbox("Value Column", cy_cols, 
                                       index=cy_cols.index(cy_mapping['value']) if cy_mapping['value'] else 0, 
                                       key='pg_cy_value')
        with col2:
            cy_product_col = st.selectbox("Product Group Column", cy_cols, 
                                         index=cy_cols.index(cy_mapping['product_group']) if cy_mapping['product_group'] else 0, 
                                         key='pg_cy_product')
            cy_company_group_col = st.selectbox("Company Group Column", cy_cols, 
                                               index=cy_cols.index(cy_mapping['company_group']) if cy_mapping['company_group'] else 0, 
                                               key='pg_cy_company_group')
            cy_exec_col = st.selectbox("Executive Column", cy_cols, 
                          index=cy_cols.index(cy_mapping['executive']) if cy_mapping['executive'] else 0, 
                          key='pg_cy_exec')
        
        st.subheader("Budget Column Mapping")
        col1, col2 = st.columns(2)
        with col1:
            budget_qty_col = st.selectbox("Quantity Column", budget_cols, 
                                         index=budget_cols.index(budget_mapping['quantity']) if budget_mapping['quantity'] else 0, 
                                         key='pg_budget_qty')
            budget_value_col = st.selectbox("Value Column", budget_cols, 
                                           index=budget_cols.index(budget_mapping['value']) if budget_mapping['value'] else 0, 
                                           key='pg_budget_value')
            budget_product_group_col = st.selectbox("Product Group Column", budget_cols, 
                                                   index=budget_cols.index(budget_mapping['product_group']) if budget_mapping['product_group'] else 0, 
                                                   key='pg_budget_product')
        with col2:
            budget_company_group_col = st.selectbox("Company Group Column", budget_cols, 
                                                   index=budget_cols.index(budget_mapping['company_group']) if budget_mapping['company_group'] else 0, 
                                                   key='pg_budget_company_group')
            budget_exec_col = st.selectbox("Executive Column", budget_cols, 
                                          index=budget_cols.index(budget_mapping['executive']) if budget_mapping['executive'] else 0, 
                                          key='pg_budget_exec')
            
        # Remove empty company groups and standardize names
        for df, col in [(ly_df, ly_company_group_col), (cy_df, cy_company_group_col), (budget_df, budget_company_group_col)]:
            df[col] = df[col].replace("", np.nan)
            df.dropna(subset=[col], inplace=True)
            df[col] = df[col].apply(standardize_name)
        
        # Get unique standardized company groups
        ly_groups = sorted(ly_df[ly_company_group_col].dropna().unique().tolist())
        cy_groups = sorted(cy_df[cy_company_group_col].dropna().unique().tolist())
        budget_groups = sorted(budget_df[budget_company_group_col].dropna().unique().tolist())
        
        all_company_groups = sorted(set(ly_groups + cy_groups + budget_groups))
        
        # Month selection
        st.subheader("Select Month Range")
        
        # Get available months
        ly_months = pd.to_datetime(ly_df[ly_date_col], dayfirst=True, errors='coerce', format='mixed').dt.strftime('%b %y').dropna().unique().tolist()
        cy_months = pd.to_datetime(cy_df[cy_date_col], dayfirst=True, errors='coerce', format='mixed').dt.strftime('%b %y').dropna().unique().tolist()
        
        if not ly_months or not cy_months:
            st.error("No valid months found in LY or CY data. Please check date columns.")
            return dfs_info
        
        col1, col2 = st.columns(2)
        with col1:
            selected_ly_months = st.multiselect("Last Year Months", sorted(ly_months), key='pg_ly_months')
        with col2:
            selected_cy_months = st.multiselect("Current Year Months", sorted(cy_months), key='pg_cy_months')
        
        # Filters for executives and company groups
        st.subheader("Filter Options")
        filter_tabs = st.tabs(["Executives", "Company Groups"])
        
        # Executive selection
        with filter_tabs[0]:
            all_executives = sorted(set(pd.concat([ly_df[ly_exec_col], cy_df[cy_exec_col], budget_df[budget_exec_col]]).dropna().unique().tolist()))
            
            exec_select_all = st.checkbox("Select All Executives", value=True, key='pg_exec_all')
            if exec_select_all:
                selected_executives = all_executives
            else:
                selected_executives = st.multiselect("Select Executives", all_executives, key='pg_executives')
        
        # Company group selection
        with filter_tabs[1]:
            group_select_all = st.checkbox("Select All Company Groups", value=True, key='pg_group_all')
            if group_select_all:
                selected_company_groups = all_company_groups
            else:
                selected_company_groups = st.multiselect("Select Company Groups", all_company_groups, key='pg_groups')
        
        # Generate report button
        if st.button("Generate Product Growth Report", key='pg_generate'):
            if selected_ly_months and selected_cy_months:
                month_title = f"LY: {', '.join(selected_ly_months)} vs CY: {', '.join(selected_cy_months)}"
                
                group_results = calculate_product_growth(
                    ly_df, cy_df, budget_df, selected_ly_months, selected_cy_months,
                    ly_date_col, cy_date_col, ly_qty_col, cy_qty_col, ly_value_col, cy_value_col,
                    budget_qty_col, budget_value_col, ly_product_col, cy_product_col,
                    ly_company_group_col, cy_company_group_col, budget_company_group_col, budget_product_group_col,
                    ly_exec_col, cy_exec_col, budget_exec_col,
                    selected_executives, selected_company_groups
                )
                
                if group_results:
                    st.subheader("Results")
                    
                    # Display results and collect data for consolidated PPT
                    dfs_info = []
                    
                    for group, data in group_results.items():
                        st.write(f"**{group} - Quantity Growth (Qty in Mt)**")
                        st.dataframe(data['qty_df'])
                        
                        st.write(f"**{group} - Value Growth (Value in Lakhs)**")
                        st.dataframe(data['value_df'])
                        
                        # Add to dfs_info for consolidated PPT
                        dfs_info.append({'df': data['qty_df'], 'title': f"{group} - Quantity Growth (Qty in Mt)", 'percent_cols': [4]})
                        dfs_info.append({'df': data['value_df'], 'title': f"{group} - Value Growth (Value in Lakhs)", 'percent_cols': [4]})
                    
                    # Store results in session state
                    st.session_state.product_results = dfs_info
                    
                    # Create individual PPT
                    ppt_buffer = create_product_growth_ppt(group_results, month_title, st.session_state.logo_file)
                    
                    if ppt_buffer:
                        unique_id = str(uuid.uuid4())[:8]
                        st.download_button(
                            label="Download Product Growth PPT",
                            data=ppt_buffer,
                            file_name=f"Product_Growth_{month_title.replace(', ', '_')}_{unique_id}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key=f"pg_download_{unique_id}"
                        )
                else:
                    st.error("Failed to generate Product Growth report. Check your data and selections.")
            else:
                st.warning("Select at least one month for LY and CY.")
    except Exception as e:
        st.error(f"Error in Product Growth tab: {e}")
        logger.error(f"Error in Product Growth tab: {e}")
    
    return dfs_info

#######################################################
# SIDEBAR AND MAIN APPLICATION
#######################################################

def sidebar_ui():
    """Sidebar UI for file uploads and app info."""
    with st.sidebar:
        st.title("ACCLLP Dashboard")
        st.subheader("File Uploads")
        
        # Sales file upload
        sales_file = st.file_uploader("Upload Current Year Sales Excel File", type=["xlsx"], key="upload_sales")
        if sales_file:
            st.session_state.sales_file = sales_file
            st.success("✅ Current Year Sales file uploaded")
        
        # Last Year Sales file upload
        last_year_sales_file = st.file_uploader("Upload Last Year Sales Excel File", type=["xlsx"], key="upload_last_year_sales")
        if last_year_sales_file:
            st.session_state.last_year_sales_file = last_year_sales_file
            st.success("✅ Last Year Sales file uploaded")
        
        # Budget file upload
        budget_file = st.file_uploader("Upload Budget Excel File", type=["xlsx"], key="upload_budget")
        if budget_file:
            st.session_state.budget_file = budget_file
            st.success("✅ Budget file uploaded")
        
        # OS Jan file upload
        os_jan_file = st.file_uploader("Upload OS-Previous Month Excel File", type=["xlsx"], key="upload_os_jan")
        if os_jan_file:
            st.session_state.os_jan_file = os_jan_file
            st.success("✅ OS-Previous Month file uploaded")
        
        # OS Feb file upload
        os_feb_file = st.file_uploader("Upload OS-Current Month Excel File", type=["xlsx"], key="upload_os_feb")
        if os_feb_file:
            st.session_state.os_feb_file = os_feb_file
            st.success("✅ OS-Current Month file uploaded")
        
        # Logo file upload
        logo_file = st.file_uploader("Upload Logo", type=["png", "jpg", "jpeg"], key="upload_logo")
        if logo_file:
            st.session_state.logo_file = logo_file
            st.image(logo_file, width=100, caption="Logo Preview")
            st.success("✅ Logo uploaded")
        
        st.divider()

def main():
    """Main application function."""
    sidebar_ui()
    
    st.title("ACCLLP Integrated Dashboard")
    
    # Check if files are uploaded and show appropriate messages
    if (not st.session_state.sales_file or 
        not st.session_state.last_year_sales_file or
        not st.session_state.budget_file or 
        not st.session_state.os_jan_file or 
        not st.session_state.os_feb_file):
        st.warning("Please upload all required files in the sidebar to access full functionality")
        
        # Show file status
        col1, col2 = st.columns(2)
        with col1:
            st.markdown("#### Required Files:")
            st.markdown(f"- Current Year Sales File: {'✅ Uploaded' if st.session_state.sales_file else '❌ Missing'}")
            st.markdown(f"- Last Year Sales File: {'✅ Uploaded' if st.session_state.last_year_sales_file else '❌ Missing'}")
            st.markdown(f"- Budget File: {'✅ Uploaded' if st.session_state.budget_file else '❌ Missing'}")
        with col2:
            st.markdown("####  ")
            st.markdown(f"- OS-Previous Month File: {'✅ Uploaded' if st.session_state.os_jan_file else '❌ Missing'}")
            st.markdown(f"- OS-Current Month File: {'✅ Uploaded' if st.session_state.os_feb_file else '❌ Missing'}")
    
    # Create tabs for each module
    tabs = st.tabs([
        "Budget vs Billed",
        "OD Target vs Collection",
        "Product Growth",
        "Number of Billed Customers"
    ])
    
    # Tab content
    with tabs[0]:
        budget_dfs = tab_budget_vs_billed()
    
    with tabs[1]:
        od_dfs = tab_od_target()
    
    with tabs[2]:
        product_dfs = tab_product_growth()  
    with tabs[3]:
        customers_dfs = tab_billed_customers()
        
    
    # Add a consolidated PPT download option if any reports have been generated
    st.divider()
    st.subheader("Consolidated Report")

    # Collect all available reports from session state
    all_dfs_info = []
    collected_sections = []

    # Budget vs Billed
    if hasattr(st.session_state, 'budget_results') and st.session_state.budget_results:
        all_dfs_info.extend(st.session_state.budget_results)
        collected_sections.append(f"Budget: {len(st.session_state.budget_results)} reports")
    # OD Target vs Collection
    if hasattr(st.session_state, 'od_results') and st.session_state.od_results:
        # Check if it's already a list or a dictionary
        if isinstance(st.session_state.od_results, list):
            all_dfs_info.extend(st.session_state.od_results)
            collected_sections.append(f"OD: {len(st.session_state.od_results)} reports")
        elif isinstance(st.session_state.od_results, dict) and 'df' in st.session_state.od_results:
            all_dfs_info.append({
                'df': st.session_state.od_results['df'], 
                'title': st.session_state.od_results.get('title', 'OD Target')
            })
            collected_sections.append("OD: 1 report")
    # Product Growth
    if hasattr(st.session_state, 'product_results') and st.session_state.product_results:
        all_dfs_info.extend(st.session_state.product_results)
        collected_sections.append(f"Product: {len(st.session_state.product_results)} reports")
    # Number of Billed Customers
    if hasattr(st.session_state, 'customers_results') and st.session_state.customers_results:
        all_dfs_info.extend(st.session_state.customers_results)
        collected_sections.append(f"Customers: {len(st.session_state.customers_results)} reports")

    # Display consolidated report options
    if all_dfs_info:
        st.info(f"Reports collected: {', '.join(collected_sections)}")
        title = st.text_input("Enter Consolidated Report Title", "ACCLLP Consolidated Report")
        
        if st.button("Generate Consolidated PPT"):
            with st.spinner("Creating consolidated PowerPoint..."):
                consolidated_ppt = create_consolidated_ppt(
                    all_dfs_info,
                    st.session_state.logo_file,
                    title
                )
                
                if consolidated_ppt:
                    unique_id = str(uuid.uuid4())[:8]
                    st.success(f"PowerPoint created with {len(all_dfs_info)} slides!")
                    st.download_button(
                        label="Download Consolidated PPT",
                        data=consolidated_ppt,
                        file_name=f"ACCLLP_Consolidated_Report_{unique_id}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"consolidated_download_{unique_id}"
                    )
                else:
                    st.error("Failed to create consolidated PowerPoint. Check the reports data.")
    else:
        st.info("Generate at least one report from any tab to enable the consolidated report option")

if __name__ == "__main__":
    main()
